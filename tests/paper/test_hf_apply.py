"""Contract tests: real fields and footer machinery.

`Presentation.apply_footers` / `Slide.apply_footers` must persist exactly what
PowerPoint's Insert > Header & Footer dialog does (mechanism pinned by the mechanism
findings and the frozen `footers_applied` fixture): materialized minimal placeholders
bound to layout furniture by idx, real `a:fld` elements for slide number and automatic
date (cached text is a consumer-refreshed hint), literal runs for footer text and fixed
dates, placeholder removal for unchecked elements, and untouched `p:hf` flags.
"""

from __future__ import annotations

import io
from datetime import datetime

import pytest

from pptx import Presentation
from pptx.errors import PaperRefusal, UnsupportedStructureError

from . import corpus
from .contract import (
    assert_changed_parts,
    assert_refusal_atomic,
    save_reopen,
    save_to_bytes,
)
from .lo import lo_load_smoke

_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
_P = "http://schemas.openxmlformats.org/presentationml/2006/main"

GAUNTLET = "self_generated/gauntlet.pptx"
MINIMAL = "self_generated/minimal_clean.pptx"
HF_FLAGS = "self_generated/hf_flags.pptx"
FOOTERS_FIXTURE = "self_generated/footers_applied.pptx"

NOW = datetime(2026, 7, 8, 14, 30, 45)


def _open(relpath):
    return Presentation(str(corpus.fixture_path(relpath)))


def _hf_placeholders(slide):
    """Map ph-type token -> (idx, fld types, text) for hf placeholders on `slide`."""
    found = {}
    for sp in slide._element.iter("{%s}sp" % _P):
        ph = sp.find(".//{%s}nvSpPr/{%s}nvPr/{%s}ph" % (_P, _P, _P))
        if ph is None or ph.get("type") not in ("dt", "ftr", "sldNum"):
            continue
        flds = [f.get("type") for f in sp.findall(".//{%s}fld" % _A)]
        text = "".join(t.text or "" for t in sp.iter("{%s}t" % _A))
        found[ph.get("type")] = (int(ph.get("idx")), flds, text)
    return found


# ------------------------------------------------------------------------- apply to all


def test_apply_to_all_materializes_dialog_structure():
    prs = _open(GAUNTLET)
    prs.apply_footers(
        footer="Paper Confidential", slide_number=True, date_format="datetime1", now=NOW
    )
    reopened = save_reopen(prs)
    for ordinal, slide in enumerate(reopened.slides, start=1):
        found = _hf_placeholders(slide)
        assert set(found) == {"dt", "ftr", "sldNum"}
        assert found["sldNum"][1] == ["slidenum"]
        assert found["sldNum"][2] == str(ordinal)  # -- cached text seeds current position
        assert found["dt"][1] == ["datetime1"]
        assert found["dt"][2] == "07/08/2026"
        assert found["ftr"][1] == []  # -- footer is a literal, never a field
        assert found["ftr"][2] == "Paper Confidential"


def test_apply_matches_the_frozen_dialog_fixture_structure():
    """The API's output must be structurally identical to the `footers_applied`
    fixture (the LibreOffice-validated reproduction of the dialog's persistence):
    same placeholder set, same layout-bound idx values, same fld types."""
    prs = Presentation()
    layout = prs.slide_layouts[1]
    for n in range(1, 6):
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text_frame.paragraphs[0].add_run().text = "Footered slide %d" % n
    prs.apply_footers(footer="Paper Fixture Footer", slide_number=True, now=NOW)
    reopened = save_reopen(prs)

    fixture = _open(FOOTERS_FIXTURE)
    for ours, theirs in zip(reopened.slides, fixture.slides):
        ours_found = _hf_placeholders(ours)
        theirs_found = _hf_placeholders(theirs)
        for kind in ("ftr", "sldNum"):
            if kind in theirs_found:
                assert ours_found[kind][0] == theirs_found[kind][0]  # -- same bound idx
                assert ours_found[kind][1] == theirs_found[kind][1]  # -- same fld types
        # -- and identical footer text where the fixture has a footer
        if "ftr" in theirs_found:
            assert ours_found["ftr"][2] == theirs_found["ftr"][2]


def test_apply_has_exact_slide_only_budget():
    prs = _open(GAUNTLET)
    before = save_to_bytes(prs)
    prs.apply_footers(footer="Budget", slide_number=True, now=NOW)
    assert_changed_parts(
        before,
        save_to_bytes(prs),
        expect_changed=["ppt/slides/slide%d.xml" % n for n in (1, 2, 3, 4)],
    )


def test_reapplying_same_state_is_a_complete_noop():
    """Field ids persist per ISO 29500 (same token until the field is removed), so an
    identical re-application must not rewrite anything."""
    prs = _open(GAUNTLET)
    prs.apply_footers(footer="Same", slide_number=True, date_format="datetime3", now=NOW)
    before = save_to_bytes(prs)
    prs2 = Presentation(io.BytesIO(before))
    prs2.apply_footers(footer="Same", slide_number=True, date_format="datetime3", now=NOW)
    assert_changed_parts(before, save_to_bytes(prs2))  # -- empty budget


def test_unchecked_elements_are_removed():
    prs = _open(GAUNTLET)
    prs.apply_footers(footer="Everything", slide_number=True, date_format="datetime1", now=NOW)
    prs.apply_footers(slide_number=True)  # -- footer and date now unchecked
    reopened = save_reopen(prs)
    for slide in reopened.slides:
        found = _hf_placeholders(slide)
        assert set(found) == {"sldNum"}


def test_fixed_date_is_a_literal_never_a_field():
    prs = _open(MINIMAL)
    prs.apply_footers(fixed_date="8 July 2026")
    reopened = save_reopen(prs)
    found = _hf_placeholders(reopened.slides[0])
    assert found["dt"][1] == []
    assert found["dt"][2] == "8 July 2026"


def test_skip_title_slides_gives_title_slides_the_removed_state():
    prs = _open(MINIMAL)  # -- its one slide is on the type="title" layout
    prs.apply_footers(footer="X", slide_number=True, now=NOW)
    prs.apply_footers(footer="X", slide_number=True, skip_title_slides=True, now=NOW)
    reopened = save_reopen(prs)
    assert _hf_placeholders(reopened.slides[0]) == {}


def test_first_slide_number_offset_is_honored():
    prs = _open(GAUNTLET)
    prs._element.set("firstSlideNum", "5")
    prs.apply_footers(slide_number=True)
    reopened = save_reopen(prs)
    cached = [_hf_placeholders(s)["sldNum"][2] for s in reopened.slides]
    assert cached == ["5", "6", "7", "8"]


def test_hf_flags_are_never_written():
    prs = _open(GAUNTLET)
    prs.apply_footers(footer="No flags", slide_number=True, date_format="datetime1", now=NOW)
    saved = save_to_bytes(prs)
    import zipfile

    with zipfile.ZipFile(io.BytesIO(saved)) as zipf:
        for name in zipf.namelist():
            if name.startswith(
                ("ppt/slides/", "ppt/slideLayouts/", "ppt/slideMasters/")
            ) and name.endswith(".xml"):
                assert b"<p:hf" not in zipf.read(name), name


# --------------------------------------------------------------------- per-slide override


def test_per_slide_override_removes_only_that_slides_footer():
    prs = _open(GAUNTLET)
    prs.apply_footers(footer="Deck footer", slide_number=True, now=NOW)
    prs.slides[2].apply_footers(slide_number=True)  # -- footer unchecked on slide 3 only
    reopened = save_reopen(prs)
    for ordinal, slide in enumerate(reopened.slides, start=1):
        found = _hf_placeholders(slide)
        if ordinal == 3:
            assert "ftr" not in found
        else:
            assert found["ftr"][2] == "Deck footer"
        assert found["sldNum"][2] == str(ordinal)


def test_per_slide_apply_budgets_to_exactly_that_slide():
    prs = _open(GAUNTLET)
    before = save_to_bytes(prs)
    prs.slides[1].apply_footers(footer="Only two", now=NOW)
    assert_changed_parts(
        before, save_to_bytes(prs), expect_changed=["ppt/slides/slide2.xml"]
    )


def test_existing_placeholder_formatting_survives_reapplication():
    """Replacing content preserves the first paragraph's pPr and first run's rPr."""
    prs = _open(MINIMAL)
    prs.apply_footers(footer="Styled", now=NOW)
    slide = prs.slides[0]
    ftr_shape = next(
        s
        for s in slide.shapes
        if s.is_placeholder and s.element.ph_type.name == "FOOTER"
    )
    ftr_shape.text_frame.paragraphs[0].runs[0].font.bold = True
    prs.slides[0].apply_footers(footer="Restyled", now=NOW)

    reopened = save_reopen(prs)
    found = _hf_placeholders(reopened.slides[0])
    assert found["ftr"][2] == "Restyled"
    ftr = next(
        s
        for s in reopened.slides[0].shapes
        if s.is_placeholder and s.element.ph_type.name == "FOOTER"
    )
    assert ftr.text_frame.paragraphs[0].runs[0].font.bold is True


# -------------------------------------------------------------------------------- refusals


def test_layout_without_furniture_refuses_atomically_deck_wide():
    """Validation must complete across the whole deck before the first write: slide 1's
    layout is fine, slide 2's layout lacks a footer placeholder - nothing may change."""
    prs = _open(GAUNTLET)
    layout = prs.slides[2].slide_layout  # -- gauntlet slide 3 is on "Blank"
    for sp in list(layout._element.iter("{%s}sp" % _P)):
        ph = sp.find(".//{%s}nvSpPr/{%s}nvPr/{%s}ph" % (_P, _P, _P))
        if ph is not None and ph.get("type") == "ftr":
            sp.getparent().remove(sp)

    raised = assert_refusal_atomic(
        prs,
        lambda p: p.apply_footers(footer="Will not land", now=NOW),
        UnsupportedStructureError,
    )
    assert "footer placeholder" in str(raised)
    assert isinstance(raised, PaperRefusal)


def test_disabling_hf_flags_refuse_atomically():
    """hf_flags fixture: the master's p:hf disables ftr (and sldNum); the layout is
    silent on ftr, so the master's disable is the nearest declaration - refuse."""
    prs = _open(HF_FLAGS)
    raised = assert_refusal_atomic(
        prs, lambda p: p.apply_footers(footer="Nope"), UnsupportedStructureError
    )
    assert "footer" in str(raised)
    assert "master" in str(raised)
    # -- dt is explicitly enabled on the master (dt="1"), so a date applies fine
    prs.apply_footers(date_format="datetime1", now=NOW)
    reopened = save_reopen(prs)
    assert "dt" in _hf_placeholders(reopened.slides[0])


def test_layout_hf_flag_overrides_masters():
    """Nearest declaration wins: the hf_flags fixture's layout 1 re-enables sldNum
    (sldNum="1") over the master's sldNum="0", so slide-number application succeeds."""
    prs = _open(HF_FLAGS)
    prs.apply_footers(slide_number=True)
    reopened = save_reopen(prs)
    assert "sldNum" in _hf_placeholders(reopened.slides[0])


@pytest.mark.parametrize(
    "kwargs",
    [
        {"footer": ""},
        {"footer": b"bytes"},
        {"footer": "ctrl\x07char"},
        {"slide_number": "yes"},
        {"slide_number": 1},
        {"date_format": "datetime99"},
        {"date_format": "datetime1", "fixed_date": "8 July"},
        {"fixed_date": ""},
        {"now": "2026-07-08"},
        {"skip_title_slides": "no"},
    ],
)
def test_bad_arguments_raise_valueerror_before_any_change(kwargs):
    prs = _open(MINIMAL)
    before = save_to_bytes(prs)
    with pytest.raises(ValueError):
        prs.apply_footers(**kwargs)
    assert_changed_parts(before, save_to_bytes(prs))  # -- empty budget


# ------------------------------------------------------------- regressions


def test_duplicate_furniture_placeholders_converge_to_dialog_state():
    """Regression: duplicated furniture placeholders (schema-legal) must
    converge to the dialog's one-per-kind state on apply, and ALL must go on uncheck."""
    import copy

    prs = _open(MINIMAL)
    prs.apply_footers(footer="Original", now=NOW)
    slide = prs.slides[0]
    ftr_sp = next(
        s._element for s in slide.shapes
        if s.is_placeholder and s.element.ph.get("type") == "ftr"
    )
    duplicate = copy.deepcopy(ftr_sp)
    duplicate.nvSpPr.cNvPr.set("id", str(int(duplicate.nvSpPr.cNvPr.get("id")) + 50))
    ftr_sp.addnext(duplicate)

    prs.slides[0].apply_footers(footer="Converged", now=NOW)
    reopened = save_reopen(prs)
    footers = [
        s for s in reopened.slides[0].shapes
        if s.is_placeholder and s.element.ph.get("type") == "ftr"
    ]
    assert len(footers) == 1
    assert footers[0].text_frame.text == "Converged"

    prs.slides[0].apply_footers()  # -- uncheck removes every footer placeholder
    reopened2 = save_reopen(prs)
    assert not any(
        s.is_placeholder and s.element.ph.get("type") == "ftr"
        for s in reopened2.slides[0].shapes
    )


def test_field_formatting_survives_reapplication():
    """Regression: rPr preservation must hold on the FIELD path too, not
    just the literal-footer path."""
    prs = _open(MINIMAL)
    prs.apply_footers(slide_number=True, footer="F1", now=NOW)
    slide = prs.slides[0]
    sldnum_sp = next(
        s._element for s in slide.shapes
        if s.is_placeholder and s.element.ph.get("type") == "sldNum"
    )
    fld = sldnum_sp.findall(".//{%s}fld" % _A)[0]
    rPr = fld.find("{%s}rPr" % _A)
    if rPr is None:  # -- a freshly-applied field carries no rPr yet
        rPr = fld.makeelement("{%s}rPr" % _A, {})
        fld.insert(0, rPr)
    rPr.set("b", "1")  # -- user styles the page number bold

    prs.slides[0].apply_footers(slide_number=True, footer="F2", now=NOW)
    reopened = save_reopen(prs)
    reopened_fld = next(
        s._element for s in reopened.slides[0].shapes
        if s.is_placeholder and s.element.ph.get("type") == "sldNum"
    ).findall(".//{%s}fld" % _A)[0]
    assert reopened_fld.find("{%s}rPr" % _A).get("b") == "1"
    assert reopened_fld.get("type") == "slidenum"


# ---------------------------------------------------------------- fields stay fields


def test_inspect_text_reports_applied_fields_as_fields():
    from pptx.inspect import inspect_text

    prs = _open(MINIMAL)
    prs.apply_footers(footer="Literal footer", slide_number=True, date_format="datetime1", now=NOW)
    reopened = save_reopen(prs)
    blocks = inspect_text(reopened.slides[0]).blocks
    field_tokens = sorted(token for b in blocks for token in b.fields)
    assert field_tokens == ["datetime1", "slidenum"]
    # -- volatile field display text stays out of block text (and therefore anchors)
    assert not any("07/08/2026" in b.text for b in blocks)
    assert any(b.text == "Literal footer" for b in blocks)


def test_slide_number_field_renumbers_by_consumers_after_move():
    """The package never rewrites cached text on reorder - consumers refresh it (the
    LibreOffice probe proved they do). What must hold structurally: the fld
    element travels with its slide, and its type stays 'slidenum'."""
    prs = _open(GAUNTLET)
    prs.apply_footers(slide_number=True)
    prs.slides.move(3, 0)  # -- last slide to the front
    reopened = save_reopen(prs)
    found = _hf_placeholders(reopened.slides[0])
    assert found["sldNum"][1] == ["slidenum"]
    assert found["sldNum"][2] == "4"  # -- stale cache by design; consumers renumber


# --------------------------------------------------------------------------------- lo_smoke


@pytest.mark.lo_smoke
def test_applied_footers_load_in_libreoffice(tmp_path):
    prs = _open(GAUNTLET)
    prs.apply_footers(
        footer="LO smoke footer", slide_number=True, date_format="datetime2", now=NOW
    )
    out = tmp_path / "footers.pptx"
    prs.save(str(out))
    lo_load_smoke(out, tmp_path)
