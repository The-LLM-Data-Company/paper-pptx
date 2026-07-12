"""Contract tests: chart addressing + safe data replacement (route, don't build).

The organ is `SlideShapes.chart_by_name` plus `Chart.replace_data_safe`: full validation,
typed refusals for structures the mechanism can't honor, then routing to upstream's public
`replace_data`. The quarantined chart XML-writer module gains no code.
"""

from __future__ import annotations

import copy
import io

import pytest

from pptx import Presentation
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.errors import (
    AmbiguousTargetError,
    TargetNotFoundError,
    UnsupportedStructureError,
)
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml.ns import qn
from pptx.util import Inches

from . import corpus
from .contract import (
    assert_changed_parts,
    assert_refusal_atomic,
    save_to_bytes,
    zip_member_map,
)
from .lo import lo_load_smoke

CHART_NOTES = "self_generated/chart_notes.pptx"
LO_CHART_NOTES = "libreoffice_export/lo_chart_notes.pptx"
MINIMAL = "self_generated/minimal_clean.pptx"
CHART_NAME = "clone_fixture_chart"


def _open(relpath):
    return Presentation(str(corpus.fixture_path(relpath)))


# ----------------------------------------------------------------------------- addressing


def test_chart_by_name_finds_the_frozen_fixture_chart():
    chart = _open(CHART_NOTES).slides[0].shapes.chart_by_name(CHART_NAME)
    assert chart.chart_type == XL_CHART_TYPE.COLUMN_CLUSTERED


def test_chart_by_name_reports_missing_and_wrong_kind_distinctly():
    shapes = _open(CHART_NOTES).slides[0].shapes
    with pytest.raises(TargetNotFoundError, match="no shape named"):
        shapes.chart_by_name("does_not_exist")
    with pytest.raises(TargetNotFoundError, match="holds no chart"):
        shapes.chart_by_name("Title 1")


def test_chart_by_name_refuses_to_pick_between_duplicates():
    prs = _open(CHART_NOTES)
    chart_data = CategoryChartData()
    chart_data.categories = ["a"]
    chart_data.add_series("s", (1,))
    frame = prs.slides[0].shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1), Inches(3), Inches(2), chart_data
    )
    frame.name = CHART_NAME
    with pytest.raises(AmbiguousTargetError):
        prs.slides[0].shapes.chart_by_name(CHART_NAME)


# ------------------------------------------------------------------------------ replacing


def test_replace_data_safe_round_trips_with_exact_budget():
    prs = _open(CHART_NOTES)
    chart = prs.slides[0].shapes.chart_by_name(CHART_NAME)
    before = save_to_bytes(prs)
    chart.replace_data_safe(
        ["North", "South"],
        [("FY25", (10.5, None)), ("FY26", (12, 13))],
        number_format="0.0",
    )
    after = save_to_bytes(prs)
    assert_changed_parts(
        before,
        after,
        expect_changed=[
            "ppt/charts/chart1.xml",
            "ppt/embeddings/Microsoft_Excel_Sheet1.xlsx",
        ],
    )

    reopened_chart = (
        Presentation(io.BytesIO(after)).slides[0].shapes.chart_by_name(CHART_NAME)
    )
    assert [(s.name, tuple(s.values)) for s in reopened_chart.series] == [
        ("FY25", (10.5, None)),
        ("FY26", (12.0, 13.0)),
    ]
    assert list(reopened_chart.plots[0].categories) == ["North", "South"]


def test_replace_data_safe_refuses_a_workbook_shared_by_another_reachable_chart():
    prs = _open(CHART_NOTES)
    slide = prs.slides[0]
    chart = slide.shapes.chart_by_name(CHART_NAME)
    chart_data = CategoryChartData()
    chart_data.categories = ["old"]
    chart_data.add_series("other", (2,))
    other = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(5),
        Inches(1),
        Inches(2),
        Inches(2),
        chart_data,
    ).chart
    shared_workbook = chart._workbook.xlsx_part
    old_rId = other._chartSpace.xlsx_part_rId
    shared_rId = other.part.relate_to(shared_workbook, RT.PACKAGE)
    other._chartSpace.externalData.set(qn("r:id"), shared_rId)
    if old_rId != shared_rId:
        other.part.drop_rel(old_rId)
    before = save_to_bytes(prs)

    with pytest.raises(UnsupportedStructureError, match="workbook is shared"):
        chart.replace_data_safe(["new"], [("series", (9,))])

    assert zip_member_map(save_to_bytes(prs)) == zip_member_map(before)
    assert other._workbook.xlsx_part is shared_workbook


def test_replace_data_safe_refuses_a_chart_part_shared_by_two_shapes():
    prs = _open(CHART_NOTES)
    source_slide = prs.slides[0]
    chart_shape = source_slide.shapes.shape_by_name(CHART_NAME)
    chart = chart_shape.chart
    other_slide = prs.slides.add_slide(prs.slide_layouts[6])
    copied_frame = copy.deepcopy(chart_shape._element)
    copied_frame.nvGraphicFramePr.cNvPr.set("id", "999")
    copied_frame.nvGraphicFramePr.cNvPr.set("name", "Shared chart view")
    rId = other_slide.part.relate_to(chart.part, RT.CHART)
    copied_frame.find(".//%s" % qn("c:chart")).set(qn("r:id"), rId)
    other_slide._element.cSld.spTree.append(copied_frame)

    with pytest.raises(UnsupportedStructureError, match="shared by multiple"):
        chart.replace_data_safe(["new"], [("series", (9,))])


def test_replace_data_safe_supports_a_chart_on_a_reachable_layout():
    prs = _open(CHART_NOTES)
    slide = prs.slides[0]
    chart_shape = slide.shapes.shape_by_name(CHART_NAME)
    chart_part = chart_shape.chart.part
    chart_ref = chart_shape._element.find(".//%s" % qn("c:chart"))
    old_rId = chart_ref.get(qn("r:id"))
    layout = slide.slide_layout
    layout_rId = layout.part.relate_to(chart_part, RT.CHART)
    chart_ref.set(qn("r:id"), layout_rId)
    layout._element.cSld.spTree.append(chart_shape._element)
    slide.part.drop_rel(old_rId)
    layout_chart = next(shape.chart for shape in layout.shapes if shape.has_chart)

    layout_chart.replace_data_safe(["new"], [("series", (9,))])

    assert tuple(layout_chart.series[0].values) == (9.0,)


def test_replace_data_safe_refuses_a_stale_chart_root():
    prs = _open(CHART_NOTES)
    chart = prs.slides[0].shapes.chart_by_name(CHART_NAME)
    chart.part._element = copy.deepcopy(chart.part._element)

    with pytest.raises(TargetNotFoundError, match="chart is stale"):
        chart.replace_data_safe(["x"], [("series", (1,))])


def test_replace_data_safe_refuses_a_chart_removed_from_its_slide():
    prs = _open(CHART_NOTES)
    slide = prs.slides[0]
    chart_shape = slide.shapes.shape_by_name(CHART_NAME)
    chart = chart_shape.chart
    slide.shapes.delete(chart_shape)

    with pytest.raises(TargetNotFoundError, match="chart is stale"):
        chart.replace_data_safe(["x"], [("series", (1,))])


@pytest.mark.parametrize(
    "bad_call",
    [
        lambda c: c.replace_data_safe([], [("a", ())]),
        lambda c: c.replace_data_safe(["x"], []),
        lambda c: c.replace_data_safe(["x"], [("a", (1, 2))]),
        lambda c: c.replace_data_safe(["x"], [("a", ("y",))]),
        lambda c: c.replace_data_safe(["x"], [("a", (True,))]),
        lambda c: c.replace_data_safe(["x"], [("a", (1,)), ("a", (2,))]),
        lambda c: c.replace_data_safe([1], [("a", (1,))]),
        lambda c: c.replace_data_safe(["x"], [("", (1,))]),
        lambda c: c.replace_data_safe(["x"], ["not-a-pair"]),
        lambda c: c.replace_data_safe(["x"], [("a", (1,))], number_format=7),
    ],
)
def test_data_shape_problems_are_valueerrors_and_leave_the_chart_untouched(bad_call):
    prs = _open(CHART_NOTES)
    chart = prs.slides[0].shapes.chart_by_name(CHART_NAME)
    before = save_to_bytes(prs)
    with pytest.raises(ValueError):
        bad_call(chart)
    assert_changed_parts(before, save_to_bytes(prs))  # -- empty budget


def test_updates_workbookless_libreoffice_chart_xml_only():
    """Previously a refusal: LibreOffice charts carry no c:externalData; the
    series rewriter runs against chart XML alone and no workbook is invented."""
    prs = _open(LO_CHART_NOTES)
    chart_shape_name = next(s for s in prs.slides[0].shapes if s.has_chart).name
    before = save_to_bytes(prs)

    chart = prs.slides[0].shapes.chart_by_name(chart_shape_name)
    chart.replace_data_safe(["North", "South"], [("FY25", (3.25, 4.5))])
    after = save_to_bytes(prs)
    assert_changed_parts(before, after, expect_changed=["ppt/charts/chart1.xml"])
    assert not any(
        n.startswith("ppt/embeddings/") for n in zip_member_map(after)
    )  # -- still workbook-less: nothing invented

    reopened_chart = (
        Presentation(io.BytesIO(after)).slides[0].shapes.chart_by_name(chart_shape_name)
    )
    assert [(s.name, tuple(s.values)) for s in reopened_chart.series] == [
        ("FY25", (3.25, 4.5))
    ]
    plot = reopened_chart.plots[0]
    assert list(plot.categories) == ["North", "South"]


def test_refuses_unsupported_chart_types_atomically():
    prs = _open(MINIMAL)
    xy_data = XyChartData()
    series = xy_data.add_series("s1")
    series.add_data_point(1, 2)
    series.add_data_point(3, 4)
    frame = prs.slides[0].shapes.add_chart(
        XL_CHART_TYPE.XY_SCATTER, Inches(1), Inches(1), Inches(4), Inches(3), xy_data
    )
    frame.name = "xy_chart"

    def operation(p):
        chart = p.slides[0].shapes.chart_by_name("xy_chart")
        chart.replace_data_safe(["x"], [("a", (1,))])

    raised = assert_refusal_atomic(prs, operation, UnsupportedStructureError)
    assert "not supported" in str(raised)


def test_refuses_multi_plot_combo_charts_atomically():
    """A combo chart (two plots in one plotArea) refuses rather than desyncing one plot."""
    import copy

    from pptx.oxml.ns import qn

    prs = _open(CHART_NOTES)
    chart = prs.slides[0].shapes.chart_by_name(CHART_NAME)
    plotArea = chart._chartSpace.chart.plotArea
    barChart = plotArea.find(qn("c:barChart"))
    barChart.addnext(copy.deepcopy(barChart))  # -- now a two-plot chart
    assert len(chart.plots) == 2

    def operation(p):
        p.slides[0].shapes.chart_by_name(CHART_NAME).replace_data_safe(["x"], [("a", (1,))])

    raised = assert_refusal_atomic(prs, operation, UnsupportedStructureError)
    assert "multi-plot" in str(raised)


def test_nonrepresentable_numerics_are_rejected_before_any_mutation():
    """Regression (review): 10**400 passed isinstance(int) then raised OverflowError AFTER
    the chart XML was rewritten — the exact XML/workbook desync the API exists to prevent.
    inf/nan would serialize as schema-invalid lexical values."""
    prs = _open(CHART_NOTES)
    chart = prs.slides[0].shapes.chart_by_name(CHART_NAME)
    before = save_to_bytes(prs)
    for bad_value in (10**400, float("inf"), float("-inf"), float("nan")):
        with pytest.raises(ValueError):
            chart.replace_data_safe(["x"], [("s", (bad_value,))])
    assert_changed_parts(before, save_to_bytes(prs))  # -- empty budget


def test_lone_surrogate_strings_are_rejected_before_any_mutation():
    """Regression: a str containing a lone surrogate passed isinstance validation, then
    exploded during serialization AFTER the chart XML had already been rewritten."""
    prs = _open(CHART_NOTES)
    chart = prs.slides[0].shapes.chart_by_name(CHART_NAME)
    before = save_to_bytes(prs)
    with pytest.raises(ValueError, match="not encodable"):
        chart.replace_data_safe(["ok"], [("bad\udc80name", (1,))])
    with pytest.raises(ValueError, match="not encodable"):
        chart.replace_data_safe(["bad\udc80cat"], [("name", (1,))])
    assert_changed_parts(before, save_to_bytes(prs))  # -- empty budget


@pytest.mark.parametrize(
    ("categories", "series", "number_format"),
    [
        (["bad\x01category"], [("name", (1,))], None),
        (["category"], [("bad\x01name", (1,))], None),
        (["category"], [("name", (1,))], "bad\x01format"),
    ],
)
def test_xml_1_0_control_characters_are_rejected_before_mutation(
    categories, series, number_format
):
    prs = _open(CHART_NOTES)
    chart = prs.slides[0].shapes.chart_by_name(CHART_NAME)
    before = save_to_bytes(prs)
    with pytest.raises(ValueError, match="XML 1.0"):
        chart.replace_data_safe(categories, series, number_format=number_format)
    assert_changed_parts(before, save_to_bytes(prs))


def test_excel_series_capacity_is_rejected_before_mutation():
    prs = _open(CHART_NOTES)
    chart = prs.slides[0].shapes.chart_by_name(CHART_NAME)
    before = save_to_bytes(prs)
    too_many = [("s%d" % idx, (idx,)) for idx in range(16_384)]
    with pytest.raises(ValueError, match="data-column limit"):
        chart.replace_data_safe(["category"], too_many)
    assert_changed_parts(before, save_to_bytes(prs))


def test_workbook_commit_failure_rolls_chart_and_workbook_back(monkeypatch):
    prs = _open(CHART_NOTES)
    chart = prs.slides[0].shapes.chart_by_name(CHART_NAME)
    before = save_to_bytes(prs)

    def fail_update(self, blob):
        raise RuntimeError("forced workbook failure")

    monkeypatch.setattr(type(chart._workbook), "update_from_xlsx_blob", fail_update)
    with pytest.raises(RuntimeError, match="forced workbook failure"):
        chart.replace_data_safe(["x"], [("series", (1,))])
    assert_changed_parts(before, save_to_bytes(prs))


def test_default_number_format_path_round_trips():
    prs = _open(CHART_NOTES)
    chart = prs.slides[0].shapes.chart_by_name(CHART_NAME)
    chart.replace_data_safe(["One", "Two"], [("Only", (4.5, 6.5))])
    reopened_chart = (
        Presentation(io.BytesIO(save_to_bytes(prs))).slides[0].shapes.chart_by_name(CHART_NAME)
    )
    assert [(s.name, tuple(s.values)) for s in reopened_chart.series] == [("Only", (4.5, 6.5))]


# --------------------------------------------------------------------------------- lo_smoke


@pytest.mark.lo_smoke
def test_replaced_chart_output_loads_in_libreoffice(tmp_path):
    prs = _open(CHART_NOTES)
    chart = prs.slides[0].shapes.chart_by_name(CHART_NAME)
    chart.replace_data_safe(["Alpha", "Beta", "Gamma"], [("S1", (1.5, 2.5, 3.5))])
    out = tmp_path / "chart_replaced.pptx"
    prs.save(str(out))
    lo_load_smoke(out, tmp_path)
