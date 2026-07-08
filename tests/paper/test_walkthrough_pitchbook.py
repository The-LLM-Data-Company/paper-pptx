"""The pitch-book assembly eval (v0.11): composition as a permanent, job-shaped test.

Nobody builds a pitch book from scratch: bank-overview pages come from the master deck,
tombstones from the credentials library, sector pages from the sector team. This module
freezes that job with shipped API only: assemble from a library deck plus a second source
deck (one slide per import mode), rebind a library page (the rebind job), renumber via
Phase 2, scrub via Phase 3 — and end with the release keystone: the operations' own
reports and `diff_decks(input, output)` agree.
"""

from __future__ import annotations

import io

import pytest

from pptx import Presentation
from pptx.diff import diff_decks

from . import corpus
from .contract import zip_member_map
from .idlists import dangling_section_slide_ids, duplicate_section_slide_ids
from .lo import lo_load_smoke
from .relint import dangling_relationship_targets, missing_relationship_references

ALPHA = "self_generated/template_alpha.pptx"  # -- the house library deck
BETA = "self_generated/template_beta.pptx"  # -- the incoming source deck


def _build_pitchbook(tmp_path):
    """Run the assembly job; return (library_path, out_path, evidence dict)."""
    library_path = tmp_path / "library.pptx"
    library_path.write_bytes(corpus.fixture_path(ALPHA).read_bytes())
    prs = Presentation(str(library_path))
    source = Presentation(str(corpus.fixture_path(BETA)))

    # -- one slide per import mode: a conscious choice each time -----------------------
    adopt = prs.import_slide(source, 0, mode="adopt_theme")  # -- takes the house look
    keep = prs.import_slide(source, 2, mode="keep_appearance")  # -- chart keeps its own
    bake = prs.import_slide(source, 1, mode="bake")  # -- frozen look, no new master

    # -- the rebind job: promote the library's content page onto Two Content ------------
    rebind = prs.slides[1].rebind_layout(
        next(layout for layout in prs.slide_layouts if layout.name == "Two Content")
    )

    # -- renumber (Phase 2): real fields across the assembled deck ----------------------
    prs.apply_footers(footer="Paper Pitch Book", slide_number=True)

    # -- exit gate (Phase 3): metadata and unused furniture go, content stays -----------
    scrub = prs.scrub(metadata=True, unused_layouts=True, unreachable_media=True)

    out_path = tmp_path / "pitchbook.pptx"
    prs.save(str(out_path))
    return library_path, out_path, {
        "adopt": adopt,
        "keep": keep,
        "bake": bake,
        "rebind": rebind,
        "scrub": scrub,
    }


def _assert_package_integrity(pptx_bytes):
    zip_map = zip_member_map(pptx_bytes)
    assert dangling_relationship_targets(zip_map) == []
    assert missing_relationship_references(zip_map) == []
    assert dangling_section_slide_ids(zip_map) == []
    assert duplicate_section_slide_ids(zip_map) == []


def test_pitchbook_end_to_end(tmp_path):
    _, out_path, evidence = _build_pitchbook(tmp_path)
    saved = out_path.read_bytes()
    _assert_package_integrity(saved)

    reopened = Presentation(io.BytesIO(saved))
    assert len(reopened.slides) == 6  # -- 3 library + 3 imported

    # -- adopt_theme: the imported title page speaks the house language now
    adopted = reopened.slides[3]
    title_font = adopted.shapes.title.text_frame.paragraphs[0].runs[0].effective_font()
    assert title_font.name.value == "Georgia"  # -- alpha's major font
    assert evidence["adopt"].run_shifts  # -- and the shift was reported, never silent

    # -- keep_appearance: the chart page kept beta's look via its transplanted chain
    kept = reopened.slides[4]
    kept_font = kept.shapes.title.text_frame.paragraphs[0].runs[0].effective_font()
    assert kept_font.name.value == "Courier New"  # -- beta's major font travelled
    assert evidence["keep"].run_shifts == ()  # -- the zero-shift invariant
    assert len(reopened.slide_masters) == 2
    chart = kept.shapes.chart_by_name("beta_chart")
    assert [series.name for plot in chart.plots for series in plot.series] == ["FY26"]

    # -- bake: frozen look, no third master; all CONTENT shapes are free now (the later
    # -- footer pass legitimately re-added dt/ftr/sldNum furniture placeholders)
    baked = reopened.slides[5]
    assert all(
        s.element.ph.get("type") in ("dt", "ftr", "sldNum")
        for s in baked.shapes
        if s.is_placeholder
    )
    baked_run = next(
        s for s in baked.shapes if s.name == "Title 1"
    ).text_frame.paragraphs[0].runs[0]
    assert baked_run.font.name == "Courier New"  # -- explicit, not inherited

    # -- renumbering: real slidenum fields on all six slides, cached 1..6
    from pptx.inspect import inspect_text

    for ordinal, slide in enumerate(reopened.slides, start=1):
        blocks = inspect_text(slide).blocks
        assert any("slidenum" in b.fields for b in blocks), "slide %d" % ordinal
        assert any(b.text == "Paper Pitch Book" for b in blocks)

    # -- scrub held: metadata gone, the assembled content untouched
    assert reopened.core_properties.author == ""
    assert evidence["scrub"].unused_layouts_removed  # -- alpha's unused furniture went


def test_pitchbook_self_consistency_against_deck_diff(tmp_path):
    """The release keystone: every operation report agrees with the deck diff."""
    library_path, out_path, evidence = _build_pitchbook(tmp_path)
    report = diff_decks(str(library_path), str(out_path), detail="text")

    # -- import reports and the diff agree on exactly which slides appeared
    added_ids = {ref.slide_id for ref in report.slides_added}
    assert added_ids == {
        evidence["adopt"].dest_slide_id,
        evidence["keep"].dest_slide_id,
        evidence["bake"].dest_slide_id,
    }
    assert report.slides_removed == ()
    assert report.slides_moved == ()

    # -- the rebound library page: diff sees the same slide the rebind report names
    rebound_id = Presentation(str(library_path)).slides[1].slide_id
    changes = {change.slide_id: change for change in report.slide_changes}
    assert rebound_id in changes  # -- footer blocks + any rebind-visible change

    # -- the footer pass shows on every ORIGINAL slide as added text blocks
    for slide_id, change in changes.items():
        if change.text_changes:
            assert all(c["before"] is None for c in change.text_changes), (
                "unexpected edits to library content on slide %d" % slide_id
            )
            assert any(
                c["after"] == "Paper Pitch Book" for c in change.text_changes
            ), "footer text missing on slide %d" % slide_id


def test_pitchbook_rebind_report_agrees_with_full_diff(tmp_path):
    """The rebind job's shift report and a full-detail diff reach the same answer."""
    library_path = tmp_path / "library.pptx"
    library_path.write_bytes(corpus.fixture_path(ALPHA).read_bytes())
    prs = Presentation(str(library_path))
    rebind_report = prs.slides[1].rebind_layout(
        next(layout for layout in prs.slide_layouts if layout.name == "Two Content")
    )
    out_path = tmp_path / "rebound.pptx"
    prs.save(str(out_path))

    diff = diff_decks(str(library_path), str(out_path), detail="full")
    rebound_id = Presentation(str(library_path)).slides[1].slide_id
    changes = {change.slide_id: change for change in diff.slide_changes}
    diff_shifts = {
        (s.text, s.before["size"]["value"], s.after["size"]["value"])
        for s in (changes[rebound_id].effective_shifts if rebound_id in changes else ())
    }
    report_shifts = {
        (s.text, s.before["size"]["value"], s.after["size"]["value"])
        for s in rebind_report.run_shifts
    }
    assert diff_shifts == report_shifts


@pytest.mark.lo_smoke
def test_pitchbook_loads_in_libreoffice(tmp_path):
    _, out_path, _ = _build_pitchbook(tmp_path)
    lo_load_smoke(out_path, tmp_path)
