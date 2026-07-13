"""Contract tests: diff_decks, the verification mirror.

Required by the plan: diff(A, A) empty across the ENTIRE corpus; the reorder-only
fixture reads as moves (never delete-plus-add); determinism goldens; the declared
id-based matching contract. The release-level self-consistency checks (operation report
vs diff of input/output) live with the standing eval jobs in test_walkthrough_*.
"""

from __future__ import annotations

import io
import json
import zipfile

import pytest
from lxml import etree

from pptx import Presentation
from pptx.diff import diff_decks
from pptx.errors import UnsupportedStructureError
from pptx.oxml.ns import qn

from . import corpus

V1 = "self_generated/lineage_v1.pptx"
V2 = "self_generated/lineage_v2.pptx"
REORDER = "self_generated/lineage_reorder.pptx"

NONCORRUPT_RELPATHS = [
    r for r in corpus.iter_fixture_relpaths() if not corpus.is_corrupt_fixture(r)
]


def _path(relpath):
    return str(corpus.fixture_path(relpath))


# ------------------------------------------------------------------- the core invariants


@pytest.mark.parametrize("relpath", NONCORRUPT_RELPATHS)
def test_self_diff_is_empty_across_entire_corpus(relpath):
    report = diff_decks(_path(relpath), _path(relpath), detail="text")
    assert report.is_empty, report.to_dict()


def test_reorder_only_fixture_reads_as_moves_never_delete_plus_add():
    report = diff_decks(_path(V1), _path(REORDER))
    assert report.slides_added == ()
    assert report.slides_removed == ()
    assert [(m.slide_id, m.from_position, m.to_position) for m in report.slides_moved] == [
        (260, 4, 0)
    ]
    assert report.slide_changes == ()


def test_lineage_pair_reports_the_exact_edit_list():
    """Every edit the v2 sidecar documents, attributed to the right slide id."""
    report = diff_decks(_path(V1), _path(V2), detail="text")

    assert [(r.slide_id, r.title) for r in report.slides_added] == [
        (261, "Lineage slide six, new")
    ]
    assert [(r.slide_id, r.title) for r in report.slides_removed] == [
        (260, "Lineage slide five")
    ]
    assert len(report.slides_moved) == 1  # -- one displacement (256<->257 tie, declared)

    changes = {change.slide_id: change for change in report.slide_changes}
    assert changes[256].text_changes == (
        {
            "shape_id": 2,
            "shape_name": "Title 1",
            "block_ordinal": 0,
            "before": "Lineage slide one",
            "after": "Lineage slide one, retitled",
        },
    )
    assert changes[257].notes_change == {
        "before": "Original notes for slide two.",
        "after": "Updated notes for slide two.",
    }
    assert changes[258].chart_data_changes == (
        {
            "chart": "lineage_chart",
            "series": "FY",
            "category": "South",
            "before": 20.0,
            "after": 25.0,
        },
    )
    assert changes[259].images_replaced == ("lineage_pic",)
    assert changes[259].geometry_changes == (
        {"shape": "lineage_box", "facet": "left", "before": 3657600, "after": 4572000},
    )


def test_diff_matches_frozen_golden():
    report = diff_decks(_path(V1), _path(V2), detail="text")
    actual = (
        json.dumps(report.to_dict(), indent=2, ensure_ascii=False) + "\n"
    ).encode("utf-8")
    golden_path = corpus.FIXTURES_DIR.parent / "goldens" / "lineage_v1_v2.diff.json"
    assert actual == golden_path.read_bytes()


def test_diff_is_deterministic_across_runs():
    first = diff_decks(_path(V1), _path(V2), detail="full").to_dict()
    second = diff_decks(_path(V1), _path(V2), detail="full").to_dict()
    assert first == second


def test_package_changes_make_metadata_only_edit_nonempty():
    changed = Presentation(_path("self_generated/minimal_clean.pptx"))
    changed.core_properties.author = "Different author"
    report = diff_decks(_path("self_generated/minimal_clean.pptx"), changed, detail="full")
    assert not report.is_empty
    assert [delta.partname for delta in report.package_changes] == [
        "/docProps/core.xml"
    ]


def test_package_changes_make_z_order_only_edit_nonempty():
    changed = Presentation(_path("self_generated/gauntlet.pptx"))
    slide = changed.slides[2]
    slide.shapes.move(slide.shapes[0], len(slide.shapes) - 1)
    report = diff_decks(_path("self_generated/gauntlet.pptx"), changed, detail="full")
    assert not report.is_empty
    assert "/ppt/slides/slide3.xml" in {
        delta.partname for delta in report.package_changes
    }


@pytest.mark.parametrize("source_kind", ["path", "stream"])
def test_package_changes_use_original_package_members(tmp_path, source_kind):
    source = corpus.fixture_path("self_generated/minimal_clean.pptx")
    modified = io.BytesIO()
    with zipfile.ZipFile(source) as before, zipfile.ZipFile(modified, "w") as after:
        for info in before.infolist():
            data = before.read(info.filename)
            if info.filename == "[Content_Types].xml":
                root = etree.fromstring(data)
                namespace = root.tag.partition("}")[0] + "}"
                etree.SubElement(
                    root,
                    namespace + "Default",
                    Extension="paper-unused",
                    ContentType="application/x-paper-unused",
                )
                data = etree.tostring(root, xml_declaration=True, encoding="UTF-8")
            after.writestr(info, data)
    modified.seek(0)

    if source_kind == "path":
        changed_source = tmp_path / "unused-content-type.pptx"
        changed_source.write_bytes(modified.getvalue())
    else:
        changed_source = modified

    report = diff_decks(source, changed_source)
    assert [delta.partname for delta in report.package_changes] == ["/[Content_Types].xml"]


def test_seekable_stream_positions_survive_success_and_failure():
    data = corpus.fixture_path("self_generated/minimal_clean.pptx").read_bytes()
    before = io.BytesIO(data)
    after = io.BytesIO(data)
    before.seek(11)
    after.seek(23)

    assert diff_decks(before, after).is_empty
    assert before.tell() == 11
    assert after.tell() == 23

    unreadable = io.BytesIO(b"not a presentation")
    before.seek(17)
    unreadable.seek(3)
    with pytest.raises(UnsupportedStructureError, match="not a readable presentation"):
        diff_decks(before, unreadable)
    assert before.tell() == 17
    assert unreadable.tell() == 3


def test_exact_package_stream_read_failure_is_typed_and_restores_position():
    class FailingExactRead(io.BytesIO):
        def read(self, size=-1):
            if self.tell() == 0 and size == 1024 * 1024:
                raise OSError("forced exact-map read failure")
            return super().read(size)

    data = corpus.fixture_path("self_generated/minimal_clean.pptx").read_bytes()
    source = FailingExactRead(data)
    source.seek(17)

    with pytest.raises(UnsupportedStructureError, match="cannot read before input stream"):
        diff_decks(source, _path("self_generated/minimal_clean.pptx"))

    assert source.tell() == 17


# ------------------------------------------------------------------------------ detail levels


def test_structure_detail_excludes_text_chart_and_notes():
    report = diff_decks(_path(V1), _path(V2), detail="structure")
    changed_ids = {change.slide_id for change in report.slide_changes}
    assert 259 in changed_ids  # -- geometry + image replacement are structural
    assert 256 not in changed_ids  # -- pure text edit invisible at structure level
    assert 257 not in changed_ids  # -- notes edit invisible at structure level
    assert 258 not in changed_ids  # -- chart data invisible at structure level


def test_full_detail_reports_effective_shifts():
    """Rebind a slide in a saved copy; diff(original, rebound) at full detail must
    carry the same run-level shifts the rebind report declared."""
    import tempfile
    from pathlib import Path

    source = Presentation(_path("self_generated/gauntlet.pptx"))
    rebind_report = source.slides[0].rebind_layout(
        source.slide_layouts[3]
    )  # -- Two Content: three genuine size shifts
    with tempfile.TemporaryDirectory() as tmp:
        out = Path(tmp) / "rebound.pptx"
        source.save(str(out))
        diff = diff_decks(_path("self_generated/gauntlet.pptx"), str(out), detail="full")

    changed = {change.slide_id: change for change in diff.slide_changes}
    slide_id = Presentation(_path("self_generated/gauntlet.pptx")).slides[0].slide_id
    shifts = changed[slide_id].effective_shifts
    assert {(s.text, s.before["size"]["value"], s.after["size"]["value"]) for s in shifts} == {
        (s.text, s.before["size"]["value"], s.after["size"]["value"])
        for s in rebind_report.run_shifts
    }


def test_diff_reports_geometry_changed_by_layout_rebind():
    before = Presentation(_path("self_generated/minimal_clean.pptx"))
    after = Presentation(_path("self_generated/minimal_clean.pptx"))
    after.slides[0].rebind_layout(after.slide_layouts[1])

    report = diff_decks(before, after, detail="structure")

    geometry = report.slide_changes[0].geometry_changes
    assert geometry
    assert {change["facet"] for change in geometry} == {
        "left", "top", "width", "height"
    }


def test_bad_detail_raises_valueerror():
    with pytest.raises(ValueError, match="detail"):
        diff_decks(_path(V1), _path(V2), detail="everything")


# ---------------------------------------------------------------------- matching contract


def test_id_matching_contract_on_rebuilt_and_re_idd_decks():
    """The declared contract, asserted for real (the previous
    version was vacuous): colliding ids on a rebuilt deck MATCH and surface the content
    difference honestly; distinct ids read as full replacement."""
    rebuilt = Presentation()
    slide = rebuilt.slides.add_slide(rebuilt.slide_layouts[0])
    slide.shapes.title.text_frame.paragraphs[0].add_run().text = "Rebuilt title"

    # -- id collision (both decks start at 256): matched slide, text delta surfaced
    report = diff_decks(_path("self_generated/minimal_clean.pptx"), rebuilt)
    report_text = diff_decks(
        _path("self_generated/minimal_clean.pptx"), rebuilt, detail="text"
    )
    assert report.slides_added == ()
    assert report.slides_removed == ()
    changed_texts = [
        (c["before"], c["after"])
        for change in report_text.slide_changes
        for c in change.text_changes
    ]
    assert ("Minimal clean fixture", "Rebuilt title") in changed_texts

    # -- distinct ids: full replacement, never a spurious match
    rebuilt._element.sldIdLst.sldId_lst[0].set("id", "300")
    report2 = diff_decks(_path("self_generated/minimal_clean.pptx"), rebuilt)
    assert [ref.slide_id for ref in report2.slides_added] == [300]
    assert [ref.slide_id for ref in report2.slides_removed] == [256]


def test_unnamed_shape_fallback_keys_are_deterministic():
    """Shapes with empty or duplicated names key as `<kind>#<ordinal>` (declared)."""
    prs_a = Presentation(_path("self_generated/minimal_clean.pptx"))
    slide = prs_a.slides[0]
    box_one = slide.shapes.add_textbox(0, 0, 914400, 914400)
    box_two = slide.shapes.add_textbox(914400, 0, 914400, 914400)
    box_one.name = "Duplicate"
    box_two.name = "Duplicate"
    box_two.text_frame.paragraphs[0].add_run().text = "changed"
    import tempfile
    from pathlib import Path

    with tempfile.TemporaryDirectory() as tmp:
        out = Path(tmp) / "dupes.pptx"
        prs_a.save(str(out))
        report = diff_decks(_path("self_generated/minimal_clean.pptx"), str(out))
    change = report.slide_changes[0]
    assert set(change.shapes_added) == {"sp#2", "sp#3"}  # -- synthetic keys, stable


def test_user_shape_name_cannot_collide_with_a_fallback_key():
    base = Presentation(_path("self_generated/minimal_clean.pptx"))
    slide = base.slides[0]
    named = slide.shapes.add_textbox(0, 0, 914400, 914400)
    named.name = "sp#1"
    for left in (914400, 1828800):
        duplicate = slide.shapes.add_textbox(left, 0, 914400, 914400)
        duplicate.name = "Duplicate"
    stream = io.BytesIO()
    base.save(stream)
    changed = Presentation(io.BytesIO(stream.getvalue()))
    changed.slides[0].shapes.shape_by_name("sp#1").left += 100

    report = diff_decks(base, changed)
    geometry = [
        delta
        for change in report.slide_changes
        for delta in change.geometry_changes
    ]
    assert any(delta["shape"] == "sp#1" for delta in geometry)


# ------------------------------------------------------------- regressions


def test_trailing_whitespace_edit_is_a_reported_text_change():
    """Regression: whitespace is content. The frozen
    trailing-space pair must diff as a text change, never as 'identical'."""
    report = diff_decks(
        _path("self_generated/whitespace_trailing_a.pptx"),
        _path("self_generated/whitespace_trailing_b.pptx"),
        detail="text",
    )
    deltas = [
        (c["before"], c["after"])
        for change in report.slide_changes
        for c in change.text_changes
    ]
    assert ("Trailing space ", "Trailing space") in deltas


def test_field_type_change_is_reported_when_visible_text_is_unchanged():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paragraph = slide.shapes.add_textbox(
        0, 0, 914400, 914400
    ).text_frame.paragraphs[0]
    paragraph.add_slide_number_field()
    before = io.BytesIO()
    prs.save(before)

    changed = Presentation(io.BytesIO(before.getvalue()))
    field = changed.slides[0].shapes[0].text_frame.paragraphs[0]._p.find(qn("a:fld"))
    field.set("type", "datetime1")

    report = diff_decks(io.BytesIO(before.getvalue()), changed, detail="text")
    change = report.slide_changes[0].text_changes[0]
    assert change["before"] == change["after"] == ""
    assert change["field_types_before"] == ["slidenum"]
    assert change["field_types_after"] == ["datetime1"]


def test_field_movement_is_reported_without_false_formatting_shifts():
    before = Presentation()
    slide = before.slides.add_slide(before.slide_layouts[6])
    paragraph = slide.shapes.add_textbox(0, 0, 914400, 914400).text_frame.paragraphs[0]
    first = paragraph.add_run()
    first.text = "Page "
    first.font.bold = True
    paragraph.add_slide_number_field()
    paragraph.add_run().text = " of total"
    stream = io.BytesIO()
    before.save(stream)
    after = Presentation(io.BytesIO(stream.getvalue()))
    changed_paragraph = after.slides[0].shapes[0].text_frame.paragraphs[0]._p
    field = changed_paragraph.find(qn("a:fld"))
    changed_paragraph.remove(field)
    pPr = changed_paragraph.find(qn("a:pPr"))
    changed_paragraph.insert(1 if pPr is not None else 0, field)

    report = diff_decks(before, after, detail="full")
    change = report.slide_changes[0].text_changes[0]

    assert change["fields_before"] == [{"offset": 5, "type": "slidenum"}]
    assert change["fields_after"] == [{"offset": 0, "type": "slidenum"}]
    assert report.slide_changes[0].effective_shifts == ()


def test_inserted_literal_run_does_not_create_false_formatting_shifts():
    before = Presentation()
    slide = before.slides.add_slide(before.slide_layouts[6])
    paragraph = slide.shapes.add_textbox(0, 0, 914400, 914400).text_frame.paragraphs[0]
    first = paragraph.add_run()
    first.text = "A"
    first.font.bold = True
    paragraph.add_run().text = "B"
    stream = io.BytesIO()
    before.save(stream)
    after = Presentation(io.BytesIO(stream.getvalue()))
    changed = after.slides[0].shapes[0].text_frame.paragraphs[0]
    inserted = changed.add_run()
    inserted.text = "X"
    inserted.font.italic = True
    changed._p.remove(inserted._r)
    first_run = changed._p.find(qn("a:r"))
    first_run.addprevious(inserted._r)

    report = diff_decks(before, after, detail="full")

    assert report.slide_changes[0].effective_shifts == ()


def test_full_detail_sees_emphasis_shifts():
    """Regression: bold/italic/underline participate in resolution-state
    comparison - a mutant dropping them must fail here."""
    prs_b = Presentation(_path("self_generated/minimal_clean.pptx"))
    run = prs_b.slides[0].shapes.title.text_frame.paragraphs[0].runs[0]
    run.font.bold = True
    run.font.italic = True
    diff = diff_decks(_path("self_generated/minimal_clean.pptx"), prs_b, detail="full")
    shifts = [s for change in diff.slide_changes for s in change.effective_shifts]
    assert len(shifts) == 1
    shift = shifts[0]
    assert shift.before["bold"]["value"] is False
    assert shift.after["bold"]["value"] is True
    assert shift.before["italic"]["value"] is False
    assert shift.after["italic"]["value"] is True


def test_shape_removal_does_not_misattribute_later_blocks():
    """Regression: text/effective comparison keys are shape-scoped, so
    removing an early shape must read as that shape's blocks disappearing - never as
    edits to the shapes below it."""
    prs_b = Presentation(_path("self_generated/gauntlet.pptx"))
    slide = prs_b.slides[0]
    title = slide.shapes.title
    slide.shapes.delete(title)
    diff = diff_decks(_path("self_generated/gauntlet.pptx"), prs_b, detail="full")
    change = next(c for c in diff.slide_changes)
    assert "Title 1" in change.shapes_removed
    # -- the title's block reads as removed (after=None); the body blocks are untouched
    for delta in change.text_changes:
        assert delta["after"] is None, delta
    assert change.effective_shifts == ()  # -- no phantom shifts from index slippage


def test_diff_accepts_presentation_objects():
    """Regression: the natural first attempt - passing open decks."""
    prs_a = Presentation(_path(V1))
    prs_b = Presentation(_path(V2))
    report = diff_decks(prs_a, prs_b, detail="structure")
    assert [ref.slide_id for ref in report.slides_added] == [261]


def test_unreadable_package_refuses_typed(tmp_path):
    from pptx.errors import UnsupportedStructureError

    garbage = tmp_path / "not-a-deck.pptx"
    garbage.write_bytes(b"this is not a zip archive")
    with pytest.raises(UnsupportedStructureError, match="not a readable presentation"):
        diff_decks(str(garbage), _path(V1))


def test_text_diff_refuses_notes_without_a_body_placeholder():
    from pptx.enum.shapes import PP_PLACEHOLDER

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    notes_slide = slide.notes_slide
    body = next(
        shape
        for shape in notes_slide.placeholders
        if shape.placeholder_format.type == PP_PLACEHOLDER.BODY
    )
    body._element.getparent().remove(body._element)

    with pytest.raises(UnsupportedStructureError, match="no body placeholder"):
        diff_decks(prs, prs, detail="text")


# ------------------------------------------------------------------ typed refusals on bad input


def test_corrupt_input_speaks_as_typed_refusal():
    """Bad input produces typed, specific refusals from the
    organs - never raw tracebacks. (Upstream loader behavior is unchanged, additively.)"""
    from pptx.errors import UnsupportedStructureError

    corrupt = _path("self_generated/corrupt_dangling_sldid.pptx")
    with pytest.raises(UnsupportedStructureError, match="relationship graph is broken"):
        diff_decks(corrupt, corrupt)
    with pytest.raises(UnsupportedStructureError, match="relationship graph is broken"):
        Presentation(corrupt).scrub(notes=True)
    with pytest.raises(UnsupportedStructureError, match="relationship graph is broken"):
        Presentation(corrupt).apply_footers(slide_number=True)
    dest = Presentation(_path("self_generated/template_alpha.pptx"))
    with pytest.raises(UnsupportedStructureError, match="relationship graph is broken"):
        dest.import_slide(Presentation(corrupt), 0, mode="bake")
