"""v0.11 Phase 1 contract tests: table structure operations.

insert_row / delete_row / insert_column / delete_column with grid-width bookkeeping and
cell-wise merged-region guards. The grid invariant (every `a:tr` holds exactly one `a:tc`
per `a:gridCol`) is asserted directly after every mutation, per the plan.
"""

from __future__ import annotations

import io

import pytest

from pptx import Presentation
from pptx.errors import PaperRefusal, UnsupportedStructureError
from pptx.util import Emu, Inches

from . import corpus
from .contract import (
    assert_changed_parts,
    assert_refusal_atomic,
    save_reopen,
    save_to_bytes,
)
from .fragval import assert_tbl_fragment_valid
from .lo import lo_load_smoke

MERGED = "self_generated/merged_tables.pptx"
LO_MERGED = "libreoffice_export/lo_merged_tables.pptx"
GAUNTLET = "self_generated/gauntlet.pptx"


def _open(relpath):
    return Presentation(str(corpus.fixture_path(relpath)))


def _merged_table(prs):
    return prs.slides[0].shapes.shape_by_name("merged_table").table


def _gauntlet_table(prs):
    return prs.slides[2].shapes.table_by_name("gauntlet_table")


def assert_grid_consistent(table):
    """The Phase 1 invariant: one a:tc per a:gridCol in every row, continuations included."""
    col_count = len(table._tbl.tblGrid.gridCol_lst)
    for tr in table._tbl.tr_lst:
        assert len(tr.tc_lst) == col_count
    assert_tbl_fragment_valid(table)


def _cell_texts(table, row_idx):
    return [table.cell(row_idx, c).text_frame.text for c in range(len(table.columns))]


# ------------------------------------------------------------------------------ insert_row


def test_insert_row_after_last_survives_save_reopen():
    prs = _open(GAUNTLET)
    table = _gauntlet_table(prs)
    new_row = table.insert_row(2)
    new_row.cells[0].text_frame.paragraphs[0].add_run().text = "appended"
    assert_grid_consistent(table)

    reopened_table = _gauntlet_table(save_reopen(prs))
    assert len(reopened_table.rows) == 4
    assert reopened_table.cell(3, 0).text_frame.text == "appended"
    assert_grid_consistent(reopened_table)


def test_insert_row_at_top_with_after_minus_one():
    prs = _open(GAUNTLET)
    table = _gauntlet_table(prs)
    table.insert_row(-1)
    assert_grid_consistent(table)
    reopened_table = _gauntlet_table(save_reopen(prs))
    assert reopened_table.cell(0, 0).text_frame.text == ""
    assert reopened_table.cell(1, 0).text_frame.text == "r0c0"  # -- old first row shifted down


def test_insert_row_has_exact_part_budget():
    prs = _open(GAUNTLET)
    before = save_to_bytes(prs)
    _gauntlet_table(prs).insert_row(2)
    assert_changed_parts(before, save_to_bytes(prs), expect_changed=["ppt/slides/slide3.xml"])


def test_insert_row_default_height_copies_neighbor():
    """The neighbor's height must be DISTINCT from every other row's, or a wrong-source
    mutant passes (final-review mutation finding); asserted through save->reopen."""
    prs = _open(GAUNTLET)
    table = _gauntlet_table(prs)
    table.rows[1].height = Emu(777240)  # -- make the neighbor unmistakable
    table.insert_row(1)
    reopened_table = _gauntlet_table(save_reopen(prs))
    assert reopened_table.rows[2].height == Emu(777240)  # -- the new row, at index 2
    assert reopened_table.rows[1].height == Emu(777240)  # -- the neighbor it copied
    assert reopened_table.rows[0].height != Emu(777240)


def test_insert_column_default_width_copies_the_distinct_neighbor():
    prs = _open(GAUNTLET)
    table = _gauntlet_table(prs)
    table.columns[1].width = Emu(555120)  # -- distinct from columns 0 and 2
    table.insert_column(1)
    reopened_table = _gauntlet_table(save_reopen(prs))
    assert reopened_table.columns[2].width == Emu(555120)  # -- the new column
    assert reopened_table.columns[0].width != Emu(555120)


def test_insert_row_copy_format_from_copies_tcPr_but_never_merges_or_text():
    prs = _open(MERGED)
    table = _merged_table(prs)
    from pptx.dml.color import RGBColor

    for col_idx in range(4):
        cell = table.cell(1, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x11, 0x22, 0x33)
    template_height = table.rows[1].height

    new_row = table.insert_row(4, copy_format_from=1)
    assert new_row.height == template_height
    assert_grid_consistent(table)

    reopened_table = _merged_table(save_reopen(prs))
    for col_idx in range(4):
        tc = reopened_table.cell(5, col_idx)._tc
        assert tc.tcPr is not None
        fill = tc.tcPr.find(
            "{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill"
        )
        assert fill is not None, "cell (5, %d) did not inherit the template fill" % col_idx
        assert tc.gridSpan == 1
        assert tc.rowSpan == 1
        assert not tc.hMerge
        assert not tc.vMerge
        assert reopened_table.cell(5, col_idx).text_frame.text == ""


def test_insert_row_boundary_inside_vertical_merge_refuses_atomically():
    prs = _open(MERGED)
    raised = assert_refusal_atomic(
        prs, lambda p: _merged_table(p).insert_row(2), UnsupportedStructureError
    )
    assert "(2, 0)" in str(raised)
    assert isinstance(raised, PaperRefusal)


def test_insert_row_below_merged_header_is_allowed():
    """The cell-wise rule: a horizontally-merged header must not poison row operations
    whose boundary only touches the region's edge."""
    prs = _open(MERGED)
    table = _merged_table(prs)
    table.insert_row(0)
    assert_grid_consistent(table)
    reopened_table = _merged_table(save_reopen(prs))
    assert reopened_table.cell(0, 0)._tc.gridSpan == 4  # -- header intact above the new row
    assert reopened_table.cell(1, 0).text_frame.text == ""


# ------------------------------------------------------------------------------ delete_row


def test_delete_row_survives_save_reopen_with_exact_budget():
    prs = _open(GAUNTLET)
    before = save_to_bytes(prs)
    table = _gauntlet_table(prs)
    deleted_texts = _cell_texts(table, 1)
    table.delete_row(1)
    assert_grid_consistent(table)
    after = save_to_bytes(prs)
    assert_changed_parts(before, after, expect_changed=["ppt/slides/slide3.xml"])

    reopened_table = _gauntlet_table(Presentation(io.BytesIO(after)))
    assert len(reopened_table.rows) == 2
    assert deleted_texts[0] not in [c.text_frame.text for c in reopened_table.iter_cells()]


def test_delete_row_containing_only_horizontal_merge_is_allowed():
    """Deleting the merged header row removes the whole merge with it - allowed."""
    prs = _open(MERGED)
    table = _merged_table(prs)
    table.delete_row(0)
    assert_grid_consistent(table)
    reopened_table = _merged_table(save_reopen(prs))
    assert len(reopened_table.rows) == 4
    assert reopened_table.cell(0, 0).text_frame.text == "r1c0"
    # -- the vertical merge (previously rows 2..3) moved up intact
    assert reopened_table.cell(1, 0)._tc.rowSpan == 2


@pytest.mark.parametrize("row_idx", [2, 3])
def test_delete_row_intersecting_vertical_merge_refuses_atomically(row_idx):
    prs = _open(MERGED)
    raised = assert_refusal_atomic(
        prs, lambda p: _merged_table(p).delete_row(row_idx), UnsupportedStructureError
    )
    assert "spanning rows 2..3" in str(raised)


def test_delete_last_remaining_row_raises_valueerror_atomically():
    prs = _open(GAUNTLET)
    table = _gauntlet_table(prs)
    table.delete_row(2)
    table.delete_row(1)
    before = save_to_bytes(prs)
    with pytest.raises(ValueError, match="last remaining row"):
        _gauntlet_table(prs).delete_row(0)
    assert_changed_parts(before, save_to_bytes(prs))  # -- empty budget


# --------------------------------------------------------------------------- insert_column


def test_insert_column_survives_save_reopen_with_exact_budget():
    prs = _open(GAUNTLET)
    before = save_to_bytes(prs)
    table = _gauntlet_table(prs)
    new_column = table.insert_column(0, width=Inches(1))
    assert new_column.width == Inches(1)
    assert_grid_consistent(table)
    after = save_to_bytes(prs)
    assert_changed_parts(before, after, expect_changed=["ppt/slides/slide3.xml"])

    reopened_table = _gauntlet_table(Presentation(io.BytesIO(after)))
    assert len(reopened_table.columns) == 4
    assert reopened_table.cell(0, 1).text_frame.text == ""
    assert reopened_table.cell(0, 0).text_frame.text == "r0c0"
    assert reopened_table.cell(0, 2).text_frame.text == "r0c1"


def test_insert_column_default_width_copies_neighbor_and_updates_frame():
    prs = _open(GAUNTLET)
    table = _gauntlet_table(prs)
    neighbor_width = table.columns[1].width
    table.insert_column(1)
    assert table.columns[2].width == neighbor_width
    shape = save_reopen(prs).slides[2].shapes.shape_by_name("gauntlet_table")
    assert shape.width == Emu(sum(c.width for c in shape.table.columns))


@pytest.mark.parametrize("after", [0, 1, 2])
def test_insert_column_boundary_inside_horizontal_merge_refuses_atomically(after):
    """Every interior boundary of the header merge (cols 0..3) must refuse - including
    the left edge (after=0), where an off-by-one in the intersection test would
    silently split the merge (final-review mutation finding)."""
    prs = _open(MERGED)
    raised = assert_refusal_atomic(
        prs, lambda p: _merged_table(p).insert_column(after), UnsupportedStructureError
    )
    assert "columns 0..3" in str(raised)


def test_insert_column_at_right_edge_of_merge_is_allowed():
    prs = _open(MERGED)
    table = _merged_table(prs)
    table.insert_column(3)
    assert_grid_consistent(table)
    reopened_table = _merged_table(save_reopen(prs))
    assert len(reopened_table.columns) == 5
    assert reopened_table.cell(0, 0)._tc.gridSpan == 4  # -- header merge untouched


# --------------------------------------------------------------------------- delete_column


def test_delete_column_survives_save_reopen_with_exact_budget():
    prs = _open(GAUNTLET)
    before = save_to_bytes(prs)
    table = _gauntlet_table(prs)
    table.delete_column(1)
    assert_grid_consistent(table)
    after = save_to_bytes(prs)
    assert_changed_parts(before, after, expect_changed=["ppt/slides/slide3.xml"])

    reopened_table = _gauntlet_table(Presentation(io.BytesIO(after)))
    assert len(reopened_table.columns) == 2
    assert [c.text_frame.text for c in reopened_table.rows[0].cells] == ["r0c0", "r0c2"]
    shape = Presentation(io.BytesIO(after)).slides[2].shapes.shape_by_name("gauntlet_table")
    assert shape.width == Emu(sum(c.width for c in shape.table.columns))


def test_delete_column_containing_whole_vertical_merge_is_allowed():
    """A vertical merge wholly inside the deleted column goes with it - cell-wise rule."""
    prs = _open(GAUNTLET)
    table = _gauntlet_table(prs)
    table.cell(0, 1).merge(table.cell(1, 1))  # -- rowSpan=2 merge contained in column 1
    table.delete_column(1)
    assert_grid_consistent(table)
    reopened_table = _gauntlet_table(save_reopen(prs))
    assert len(reopened_table.columns) == 2
    for cell in reopened_table.iter_cells():
        assert not cell._tc.is_spanned
        assert not cell._tc.is_merge_origin


@pytest.mark.parametrize("col_idx", [0, 1, 3])
def test_delete_column_intersecting_horizontal_merge_refuses_atomically(col_idx):
    prs = _open(MERGED)
    raised = assert_refusal_atomic(
        prs, lambda p: _merged_table(p).delete_column(col_idx), UnsupportedStructureError
    )
    assert "columns 0..3" in str(raised)


def test_delete_last_remaining_column_raises_valueerror_atomically():
    prs = _open(GAUNTLET)
    table = _gauntlet_table(prs)
    table.delete_column(2)
    table.delete_column(1)
    before = save_to_bytes(prs)
    with pytest.raises(ValueError, match="last remaining column"):
        _gauntlet_table(prs).delete_column(0)
    assert_changed_parts(before, save_to_bytes(prs))  # -- empty budget


# ----------------------------------------------------------------- arguments and invariants


@pytest.mark.parametrize("bad", [True, False, "1", 1.0, None, 99, -2])
def test_bad_indices_raise_valueerror_before_any_change(bad):
    prs = _open(MERGED)
    before = save_to_bytes(prs)
    table = _merged_table(prs)
    operations = [
        lambda: table.insert_row(bad),
        lambda: table.delete_row(bad),
        lambda: table.insert_column(bad),
        lambda: table.delete_column(bad),
    ]
    if bad is not None:  # -- None is copy_format_from's legitimate default
        operations.append(lambda: table.insert_row(0, copy_format_from=bad))
    for operation in operations:
        with pytest.raises(ValueError):
            operation()
    assert_changed_parts(before, save_to_bytes(prs))  # -- empty budget


def test_insert_column_rejects_bad_width():
    prs = _open(MERGED)
    table = _merged_table(prs)
    for bad in (True, 0, -914400, 1.5, "wide"):
        with pytest.raises(ValueError):
            table.insert_column(3, width=bad)


def test_insert_then_delete_row_is_a_complete_noop():
    prs = _open(MERGED)
    before = save_to_bytes(prs)
    table = _merged_table(prs)
    table.insert_row(4)
    table.delete_row(5)
    assert_changed_parts(before, save_to_bytes(prs))  # -- empty budget


def test_insert_then_delete_column_is_a_complete_noop():
    prs = _open(MERGED)
    before = save_to_bytes(prs)
    table = _merged_table(prs)
    table.insert_column(3)
    table.delete_column(4)
    assert_changed_parts(before, save_to_bytes(prs))  # -- empty budget


def test_operations_on_libreoffice_authored_merged_table():
    """Producer diversity: the same guards and surgery on LibreOffice-written bytes."""
    prs = Presentation(str(corpus.fixture_path(LO_MERGED)))
    table = prs.slides[0].shapes.shape_by_name("merged_table").table
    raised = assert_refusal_atomic(
        prs,
        lambda p: p.slides[0].shapes.shape_by_name("merged_table").table.delete_row(2),
        UnsupportedStructureError,
    )
    assert "spanning rows 2..3" in str(raised)

    table.insert_row(4)
    table.insert_column(3)
    assert_grid_consistent(table)
    reopened = save_reopen(prs)
    reopened_table = reopened.slides[0].shapes.shape_by_name("merged_table").table
    assert len(reopened_table.rows) == 6
    assert len(reopened_table.columns) == 5
    assert reopened_table.cell(0, 0)._tc.gridSpan == 4


# ------------------------------------------------------- anchored writes reach table cells


def test_anchored_replace_reaches_table_cells_including_merge_origin():
    """PLAN v0.11 Phase 1: cell text edits route through the v0.1 anchored-write path -
    prove that path reaches table cells, merged origins included."""
    from pptx.edit import replace_text, replace_text_at
    from pptx.inspect import inspect_text

    prs = _open(MERGED)
    result = replace_text(prs, "Merged header", "Merged HEADING")
    assert result.replacements == 1

    reopened = save_reopen(prs)
    reopened_table = _merged_table(reopened)
    assert reopened_table.cell(0, 0).text_frame.text == "Merged HEADING"

    # -- and the single-block anchored path, via an inspect_text table-cell anchor
    blocks = inspect_text(reopened.slides[0])
    cell_block = next(
        b for b in blocks.blocks if b.container == "table-cell" and b.text == "r1c2"
    )
    result = replace_text_at(reopened, cell_block.anchor, "r1c2", "R1C2 EDITED")
    assert result.replacements == 1
    final = save_reopen(reopened)
    assert _merged_table(final).cell(1, 2).text_frame.text == "R1C2 EDITED"


# --------------------------------------------------------------------------------- lo_smoke


@pytest.mark.lo_smoke
def test_surgered_table_loads_in_libreoffice(tmp_path):
    prs = _open(MERGED)
    table = _merged_table(prs)
    table.insert_row(4, copy_format_from=1)
    table.insert_column(3)
    table.delete_row(1)
    assert_grid_consistent(table)
    out = tmp_path / "surgered.pptx"
    prs.save(str(out))
    lo_load_smoke(out, tmp_path)
