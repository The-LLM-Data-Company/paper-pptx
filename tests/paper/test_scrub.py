"""v0.11 Phase 3 contract tests: scrub, the exit gate.

Every removal is relationship-graph surgery, so reachable parts structurally cannot be
removed; the |ScrubReport| carries the exact zip-member budget and every test here holds
the actual diff to it member for member. relint + the section-integrity scan run on every
scrub output, per the plan.
"""

from __future__ import annotations

import io
import json

import pytest

from pptx import Presentation

from . import corpus
from .contract import assert_changed_parts, save_to_bytes, zip_member_map
from .idlists import dangling_section_slide_ids, duplicate_section_slide_ids
from .lo import lo_load_smoke
from .relint import dangling_relationship_targets, missing_relationship_references

SCRUB_GAUNTLET = "self_generated/scrub_gauntlet.pptx"
GAUNTLET = "self_generated/gauntlet.pptx"

EVERYTHING = {
    "notes": True,
    "comments": True,
    "metadata": True,
    "hidden_slides": True,
    "unused_layouts": True,
    "unused_masters": True,
    "unreachable_media": True,
    "embedded_fonts": True,
}


def _open(relpath):
    return Presentation(str(corpus.fixture_path(relpath)))


def _assert_clean(after_bytes):
    zip_map = zip_member_map(after_bytes)
    assert dangling_relationship_targets(zip_map) == []
    assert missing_relationship_references(zip_map) == []
    assert dangling_section_slide_ids(zip_map) == []
    assert duplicate_section_slide_ids(zip_map) == []


def _scrub_with_budget(relpath, **toggles):
    """Scrub, assert the changed-part diff matches the report exactly, return both."""
    prs = _open(relpath)
    before = save_to_bytes(prs)
    report = prs.scrub(**toggles)
    after = save_to_bytes(prs)
    assert_changed_parts(
        before,
        after,
        expect_changed=report.parts_modified,
        expect_removed=report.parts_removed,
    )
    _assert_clean(after)
    return report, before, after


# ----------------------------------------------------------------------- the acceptance job


def test_scrub_everything_acceptance():
    """The plan's job-shaped acceptance: gauntlet deck -> scrub(everything) -> reopens
    clean, visibly identical retained slides, zero notes/comments/metadata, smaller file."""
    from pptx.inspect import inspect_text

    original = _open(SCRUB_GAUNTLET)
    visible_before = [
        [b.text for b in inspect_text(s).blocks]
        for s in list(original.slides)[:3]  # -- slide 4 is hidden and will be removed
    ]

    report, before, after = _scrub_with_budget(SCRUB_GAUNTLET, **EVERYTHING)

    assert len(after) < len(before)
    reopened = Presentation(io.BytesIO(after))
    assert len(reopened.slides) == 3
    visible_after = [[b.text for b in inspect_text(s).blocks] for s in reopened.slides]
    assert visible_after == visible_before  # -- visible content untouched

    members = zip_member_map(after)
    assert not any(m.startswith("ppt/notesSlides/") for m in members)
    assert not any(m.startswith("ppt/comments/") for m in members)
    assert not any(m.startswith("ppt/fonts/") for m in members)
    assert "docProps/custom.xml" not in members
    assert "docProps/app.xml" not in members
    assert b"embeddedFontLst" not in members["ppt/presentation.xml"]

    core = reopened.core_properties
    assert core.author == ""
    assert core.last_modified_by == ""
    assert core.comments == ""

    # -- the report is the operation's evidence object
    assert report.hidden_slides_removed == ("/ppt/slides/slide4.xml",)
    assert report.embedded_font_parts_removed == ("/ppt/fonts/font1.fntdata",)
    assert report.notes_master_retained is True
    assert "ppt/media/image_unused_layout.png" in report.parts_removed


def test_scrub_report_matches_frozen_golden():
    """Byte-identical to the reviewed golden; update ONLY via update_goldens.py + PR."""
    prs = _open(SCRUB_GAUNTLET)
    report = prs.scrub(**EVERYTHING)
    actual = (
        json.dumps(report.to_dict(), indent=2, ensure_ascii=False) + "\n"
    ).encode("utf-8")
    golden_path = corpus.FIXTURES_DIR.parent / "goldens" / "scrub_gauntlet.scrub.json"
    assert actual == golden_path.read_bytes()


def test_noop_scrub_changes_nothing():
    prs = _open(SCRUB_GAUNTLET)
    before = save_to_bytes(prs)
    report = prs.scrub()
    assert_changed_parts(before, save_to_bytes(prs))  # -- empty budget
    assert report.parts_removed == ()
    assert report.parts_modified == ()


# ------------------------------------------------------------------- individual toggles


def test_notes_only_removes_notes_and_keeps_everything_else():
    report, _, after = _scrub_with_budget(SCRUB_GAUNTLET, notes=True)
    assert report.notes_slides_removed == (
        "/ppt/notesSlides/notesSlide1.xml",
        "/ppt/notesSlides/notesSlide2.xml",
    )
    assert report.notes_master_retained is True
    members = zip_member_map(after)
    assert not any(m.startswith("ppt/notesSlides/") for m in members)
    assert "ppt/notesMasters/notesMaster1.xml" in members  # -- declared retention
    assert "ppt/comments/comment1.xml" in members
    reopened = Presentation(io.BytesIO(after))
    assert not reopened.slides[0].has_notes_slide


def test_comments_only_removes_comment_parts_and_author_registry():
    report, _, after = _scrub_with_budget(SCRUB_GAUNTLET, comments=True)
    assert report.comment_parts_removed == ("/ppt/comments/comment1.xml",)
    assert report.comment_author_parts_removed == ("/ppt/commentAuthors.xml",)
    members = zip_member_map(after)
    assert not any(m.startswith("ppt/comments") for m in members)
    assert "ppt/commentAuthors.xml" not in members
    assert any(m.startswith("ppt/notesSlides/") for m in members)  # -- notes untouched


def test_metadata_only_clears_core_fields_and_removes_props_parts():
    report, _, after = _scrub_with_budget(SCRUB_GAUNTLET, metadata=True)
    assert report.metadata_fields_cleared == ("author", "comments", "last_modified_by")
    assert "/docProps/app.xml" in report.metadata_parts_removed
    assert "/docProps/custom.xml" in report.metadata_parts_removed
    reopened = Presentation(io.BytesIO(after))
    core = reopened.core_properties
    assert core.author == ""
    assert core.comments == ""
    assert core.last_modified_by == ""
    # -- declared: created/modified survive (pipeline-relevant, not personal)
    assert core.created is not None


def test_hidden_slides_only():
    report, _, after = _scrub_with_budget(SCRUB_GAUNTLET, hidden_slides=True)
    assert report.hidden_slides_removed == ("/ppt/slides/slide4.xml",)
    reopened = Presentation(io.BytesIO(after))
    assert len(reopened.slides) == 3
    assert all(s._element.get("show") is None for s in reopened.slides)


def test_unused_layouts_only_removes_them_with_their_exclusive_media():
    report, _, after = _scrub_with_budget(SCRUB_GAUNTLET, unused_layouts=True)
    # -- hidden slide retained, so its layout (slideLayout7) stays used
    assert "/ppt/slideLayouts/slideLayout7.xml" not in report.unused_layouts_removed
    assert "/ppt/slideLayouts/slideLayout4.xml" in report.unused_layouts_removed
    assert "ppt/media/image_unused_layout.png" in report.parts_removed
    members = zip_member_map(after)
    assert "ppt/slideLayouts/slideLayout2.xml" in members  # -- used layouts survive
    reopened = Presentation(io.BytesIO(after))
    assert len(reopened.slides) == 4  # -- slides untouched


def test_unused_masters_negative_master_in_use_survives():
    report, _, after = _scrub_with_budget(SCRUB_GAUNTLET, unused_masters=True)
    assert report.unused_masters_removed == ()
    assert "ppt/slideMasters/slideMaster1.xml" in zip_member_map(after)


def test_unreachable_media_drops_only_unreferenced_relationships():
    """A media rel whose rId no XML reference uses is dropped; referenced media never."""
    prs = _open(SCRUB_GAUNTLET)
    slide = prs.slides[1]  # -- slide 2 has no picture yet
    picture = slide.shapes.add_picture(
        io.BytesIO(_distinct_png()), 0, 0, 914400
    )
    # -- strip the shape but leave the relationship: the classic leak
    picture._element.getparent().remove(picture._element)

    before = save_to_bytes(prs)
    report = prs.scrub(unreachable_media=True)
    after = save_to_bytes(prs)
    assert_changed_parts(
        before,
        after,
        expect_changed=report.parts_modified,
        expect_removed=report.parts_removed,
    )
    assert len(report.unreachable_media_rels_dropped) == 1
    # -- slide 1's referenced picture is untouched
    reopened = Presentation(io.BytesIO(after))
    live_pic = reopened.slides[0].shapes.picture_by_name("live_media_pic")
    assert live_pic.image.blob  # -- loads fine
    _assert_clean(after)


def test_referenced_media_is_never_dropped():
    report, _, after = _scrub_with_budget(SCRUB_GAUNTLET, unreachable_media=True)
    assert report.unreachable_media_rels_dropped == ()
    assert report.parts_removed == ()


def test_toggle_validation_rejects_non_bool_atomically():
    prs = _open(SCRUB_GAUNTLET)
    before = save_to_bytes(prs)
    with pytest.raises(ValueError):
        prs.scrub(notes="yes")
    with pytest.raises(ValueError):
        prs.scrub(unused_layouts=1)
    assert_changed_parts(before, save_to_bytes(prs))  # -- empty budget


def test_scrub_report_is_deterministic():
    first = _open(SCRUB_GAUNTLET).scrub(**EVERYTHING).to_dict()
    second = _open(SCRUB_GAUNTLET).scrub(**EVERYTHING).to_dict()
    assert first == second


def test_scrub_everything_on_clean_deck_is_mostly_noop():
    """The plain gauntlet has notes + unused layouts but no comments/fonts/hidden slides:
    scrub removes exactly what exists and reports nothing for absent categories."""
    report, _, after = _scrub_with_budget(GAUNTLET, **EVERYTHING)
    assert report.comment_parts_removed == ()
    assert report.hidden_slides_removed == ()
    assert report.embedded_font_parts_removed == ()
    assert report.notes_slides_removed == ("/ppt/notesSlides/notesSlide1.xml",)
    reopened = Presentation(io.BytesIO(after))
    assert len(reopened.slides) == 4


def _distinct_png():
    import io as _io

    from PIL import Image

    img = Image.new("RGB", (8, 8), (7, 77, 177))
    buf = _io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


# ------------------------------------------------------------- final-review regressions


def test_metadata_clears_every_documented_field():
    """Regression (final review): all 11 documented core fields clear, not just the 3
    the fixture happens to populate."""
    from pptx.scrub import _CLEARED_CORE_FIELDS

    prs = _open(SCRUB_GAUNTLET)
    core = prs.core_properties
    for field_name in _CLEARED_CORE_FIELDS:
        setattr(core, field_name, "populated-%s" % field_name)
    before = save_to_bytes(prs)
    prs2 = Presentation(io.BytesIO(before))
    report = prs2.scrub(metadata=True)
    assert report.metadata_fields_cleared == tuple(sorted(_CLEARED_CORE_FIELDS))
    after = save_to_bytes(prs2)
    assert_changed_parts(
        before, after, expect_changed=report.parts_modified,
        expect_removed=report.parts_removed,
    )
    reopened = Presentation(io.BytesIO(after))
    for field_name in _CLEARED_CORE_FIELDS:
        assert getattr(reopened.core_properties, field_name) == "", field_name


def test_scrub_with_broken_layout_relationship():
    """Regression (final review): a broken slide->layout relationship must (a) not stop
    the all-False scrub from returning its promised empty report, and (b) refuse TYPED
    and ATOMICALLY before the layout-usage passes - never mutate-then-KeyError."""
    from pptx.errors import UnsupportedStructureError
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT

    def _broken_deck():
        prs = Presentation()
        slide_one = prs.slides.add_slide(prs.slide_layouts[6])
        slide_two = prs.slides.add_slide(prs.slide_layouts[6])
        slide_two.notes_slide.notes_text_frame.text = "must survive the refusal"
        for rId, rel in list(slide_one.part.rels.items()):
            if rel.reltype == RT.SLIDE_LAYOUT:
                slide_one.part.rels.pop(rId)
        return prs

    prs = _broken_deck()
    report = prs.scrub()  # -- all toggles False: empty report, no crash
    assert report.parts_removed == ()
    assert report.parts_modified == ()

    prs2 = _broken_deck()
    with pytest.raises(UnsupportedStructureError, match="layout relationship is broken"):
        prs2.scrub(notes=True, unused_layouts=True)
    # -- atomic: the notes pass must NOT have run before the refusal
    assert prs2.slides[1].has_notes_slide


def test_scrub_metadata_without_core_properties_part_is_a_noop():
    """Regression (final review): a package with no core-properties part must not have
    one CREATED by the metadata pass (scrub may never create parts)."""
    prs = Presentation("tests/test_files/no-core-props.pptx")
    before = save_to_bytes(prs)
    report = prs.scrub(metadata=True)
    assert report.metadata_fields_cleared == ()
    assert_changed_parts(
        before,
        save_to_bytes(prs),
        expect_changed=report.parts_modified,
        expect_removed=report.parts_removed,
    )
    assert "docProps/core.xml" not in zip_member_map(save_to_bytes(prs))


def test_content_types_budget_covers_override_typed_media():
    """Regression (final review): a removed media part carrying its own Override (e.g.
    image/svg+xml) changes [Content_Types].xml even when other parts share its
    extension - the budget must say so (and the exact-budget helper proves it)."""
    prs = _open(SCRUB_GAUNTLET)
    package = prs.part.package
    unused_media = next(
        part
        for part in package.iter_parts()
        if str(part.partname) == "/ppt/media/image_unused_layout.png"
    )
    unused_media._content_type = "image/svg+xml"  # -- forces an Override for this part
    before = save_to_bytes(prs)
    report = prs.scrub(unused_layouts=True)
    assert "ppt/media/image_unused_layout.png" in report.parts_removed
    assert "[Content_Types].xml" in report.parts_modified
    assert_changed_parts(
        before,
        save_to_bytes(prs),
        expect_changed=report.parts_modified,
        expect_removed=report.parts_removed,
    )


def test_notes_master_retained_reports_false_when_none_exists():
    prs = _open("self_generated/minimal_clean.pptx")
    report = prs.scrub(notes=True)
    assert report.notes_slides_removed == ()
    assert report.notes_master_retained is False


# --------------------------------------------------------------------------------- lo_smoke


@pytest.mark.lo_smoke
def test_scrubbed_deck_loads_in_libreoffice(tmp_path):
    prs = _open(SCRUB_GAUNTLET)
    prs.scrub(**EVERYTHING)
    out = tmp_path / "scrubbed.pptx"
    prs.save(str(out))
    lo_load_smoke(out, tmp_path)
