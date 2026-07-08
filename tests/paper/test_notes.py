"""Phase 6 contract tests: speaker notes on existing notes parts only.

The upstream `notes_slide` accessor auto-creates a whole notes part graph on access; these
APIs never do — a slide without notes refuses, atomically, and the refusal provably does not
trigger creation. Replacement touches only the notes body placeholder, preserving formatting.
"""

from __future__ import annotations

import pytest

from pptx import Presentation
from pptx.errors import PaperRefusal, UnsupportedStructureError
from pptx.util import Pt

from . import corpus
from .contract import (
    assert_changed_parts,
    assert_refusal_atomic,
    save_reopen,
    save_to_bytes,
)
from .lo import lo_load_smoke

CHART_NOTES = "self_generated/chart_notes.pptx"
LO_CHART_NOTES = "libreoffice_export/lo_chart_notes.pptx"
MINIMAL = "self_generated/minimal_clean.pptx"


def _open(relpath):
    return Presentation(str(corpus.fixture_path(relpath)))


# ------------------------------------------------------------------------------- reading


def test_reads_notes_text_from_frozen_fixture():
    expected = corpus.load_sidecar(CHART_NOTES)["ground_truth"]["notes_text"]
    assert _open(CHART_NOTES).slides[0].read_notes_text() == expected


def test_reads_notes_from_a_libreoffice_authored_deck():
    """LO rewrote the notes slide (added a <number> field placeholder); body still reads."""
    assert (
        _open(LO_CHART_NOTES).slides[0].read_notes_text()
        == "Speaker notes for the clone fixture."
    )


def test_reading_empty_notes_body_returns_empty_string():
    prs = _open(CHART_NOTES)
    prs.slides[0].replace_notes_text("")
    assert save_reopen(prs).slides[0].read_notes_text() == ""


# ------------------------------------------------------------------------------ replacing


def test_replace_round_trips_with_exact_part_budget():
    prs = _open(CHART_NOTES)
    before = save_to_bytes(prs)
    prs.slides[0].replace_notes_text("Replaced notes.")
    assert_changed_parts(
        before, save_to_bytes(prs), expect_changed=["ppt/notesSlides/notesSlide1.xml"]
    )
    assert save_reopen(prs).slides[0].read_notes_text() == "Replaced notes."


def test_replace_preserves_first_run_formatting_and_splits_paragraphs():
    prs = _open(CHART_NOTES)
    notes_tf = prs.slides[0].notes_slide.notes_text_frame
    notes_tf.paragraphs[0].runs[0].font.size = Pt(14)
    notes_tf.paragraphs[0].runs[0].font.bold = True

    prs.slides[0].replace_notes_text("First line\nSecond line\n\nFourth line")

    reopened_tf = save_reopen(prs).slides[0].notes_slide.notes_text_frame
    assert len(reopened_tf.paragraphs) == 4
    assert reopened_tf.paragraphs[0].runs[0].font.size.pt == 14.0
    assert reopened_tf.paragraphs[0].runs[0].font.bold is True
    assert reopened_tf.paragraphs[1].runs[0].font.size.pt == 14.0  # -- template carried over
    assert reopened_tf.paragraphs[2].runs == ()  # -- empty line = empty paragraph


def test_replace_leaves_other_notes_placeholders_untouched():
    prs = _open(CHART_NOTES)
    other_xml_before = [
        ph._element.xml
        for ph in prs.slides[0].notes_slide.placeholders
        if ph.placeholder_format.type.name != "BODY"
    ]
    prs.slides[0].replace_notes_text("Only the body changes.")
    reopened = save_reopen(prs)
    other_xml_after = [
        ph._element.xml
        for ph in reopened.slides[0].notes_slide.placeholders
        if ph.placeholder_format.type.name != "BODY"
    ]
    assert other_xml_after == other_xml_before


# -------------------------------------------------------------------- refusals, atomically


def test_read_refuses_without_notes_and_never_creates_one():
    prs = _open(MINIMAL)
    raised = assert_refusal_atomic(
        prs, lambda p: p.slides[0].read_notes_text(), UnsupportedStructureError
    )
    assert isinstance(raised, PaperRefusal)
    assert prs.slides[0].has_notes_slide is False


def test_replace_refuses_without_notes_and_never_creates_one():
    prs = _open(MINIMAL)
    assert_refusal_atomic(
        prs, lambda p: p.slides[0].replace_notes_text("nope"), UnsupportedStructureError
    )
    assert prs.slides[0].has_notes_slide is False


def test_replace_rejects_non_string_before_touching_anything():
    prs = _open(CHART_NOTES)
    before = save_to_bytes(prs)
    with pytest.raises(ValueError):
        prs.slides[0].replace_notes_text(42)
    assert_changed_parts(before, save_to_bytes(prs))  # -- empty budget


# -------------------------------------------------------------------------------- lo_smoke


@pytest.mark.lo_smoke
def test_replaced_notes_output_loads_in_libreoffice(tmp_path):
    prs = _open(CHART_NOTES)
    prs.slides[0].replace_notes_text("LibreOffice smoke notes.")
    out = tmp_path / "notes.pptx"
    prs.save(str(out))
    lo_load_smoke(out, tmp_path)
