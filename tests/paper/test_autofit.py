"""Phase 3 contract tests: autofit reading and `TextFrame.normalize_autofit`.

The organ extends upstream's `auto_size` property (never a parallel API): `font_scale` /
`line_space_reduction` expose what `a:normAutofit` records, and `normalize_autofit` freezes
rendered metrics into explicit values before setting `a:noAutofit`. Where freezing would
require information that is not locally resolvable, it refuses — proven atomic here.
"""

from __future__ import annotations

import pytest

from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.errors import PaperRefusal, UnsupportedStructureError
from pptx.util import Pt

from . import corpus
from .contract import (
    assert_changed_parts,
    assert_refusal_atomic,
    save_reopen,
    save_to_bytes,
)
from .fragval import assert_bodyPr_fragment_valid
from .lo import lo_load_smoke

NORMAL = "self_generated/autofit_normal.pptx"
NONE = "self_generated/autofit_none.pptx"
SHAPE = "self_generated/autofit_shape.pptx"
LO_NORMAL = "libreoffice_export/lo_autofit_normal.pptx"


def _autofit_frame(prs):
    return next(s for s in prs.slides[0].shapes if s.name == "autofit_box").text_frame


def _open(relpath):
    return Presentation(str(corpus.fixture_path(relpath)))


# ------------------------------------------------------------------------- reading details


def test_reads_font_scale_and_reduction_from_frozen_fixture():
    tf = _autofit_frame(_open(NORMAL))
    assert tf.auto_size == MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    assert tf.font_scale == 62.5
    assert tf.line_space_reduction == 20.0


def test_reads_defaults_when_normautofit_carries_no_attributes():
    """The LibreOffice-authored fixture has bare <a:normAutofit/> (attrs stripped by LO)."""
    tf = _autofit_frame(_open(LO_NORMAL))
    assert tf.font_scale == 100.0
    assert tf.line_space_reduction == 0.0


@pytest.mark.parametrize(
    ("relpath", "expected_mode"),
    [(NONE, MSO_AUTO_SIZE.NONE), (SHAPE, MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)],
)
def test_non_normautofit_modes_report_none_details(relpath, expected_mode):
    tf = _autofit_frame(_open(relpath))
    assert tf.auto_size == expected_mode
    assert tf.font_scale is None
    assert tf.line_space_reduction is None


# ------------------------------------------------------------------- refusals, atomically


def test_refuses_to_guess_inherited_font_sizes():
    """The frozen normAutofit fixture's run carries no explicit size: freezing must refuse."""
    prs = _open(NORMAL)

    def operation(prs):
        _autofit_frame(prs).normalize_autofit()

    raised = assert_refusal_atomic(prs, operation, UnsupportedStructureError)
    assert "font size" in str(raised)
    assert isinstance(raised, PaperRefusal)


def test_refuses_to_guess_inherited_line_spacing():
    prs = _open(NORMAL)
    tf = _autofit_frame(prs)
    tf.paragraphs[0].runs[0].font.size = Pt(24)  # -- sizes resolvable, spacing still not

    def operation(prs):
        _autofit_frame(prs).normalize_autofit()

    raised = assert_refusal_atomic(prs, operation, UnsupportedStructureError)
    assert "line spacing" in str(raised)


def test_bad_min_font_size_is_a_valueerror():
    tf = _autofit_frame(_open(NORMAL))
    with pytest.raises(ValueError):
        tf.normalize_autofit(min_font_size=-3)
    with pytest.raises(ValueError):
        tf.normalize_autofit(min_font_size=11.0)  # -- floats are ambiguous; require Length


# ------------------------------------------------------------------------------- freezing


def test_normalize_freezes_scale_and_spacing_and_round_trips():
    prs = _open(NORMAL)
    tf = _autofit_frame(prs)
    tf.paragraphs[0].runs[0].font.size = Pt(24)
    tf.paragraphs[0].line_spacing = 1.0
    before = save_to_bytes(prs)

    tf.normalize_autofit()
    assert_bodyPr_fragment_valid(tf)
    assert_changed_parts(before, save_to_bytes(prs), expect_changed=["ppt/slides/slide1.xml"])

    reopened_tf = _autofit_frame(save_reopen(prs))
    assert reopened_tf.auto_size == MSO_AUTO_SIZE.NONE
    assert reopened_tf.font_scale is None
    assert reopened_tf.paragraphs[0].runs[0].font.size.pt == 15.0  # -- 24pt × 62.5%
    assert reopened_tf.paragraphs[0].line_spacing == 0.8  # -- 1.0 × (100−20)%


def test_normalize_scales_paragraph_default_sizes_without_touching_inherited_runs():
    """A paragraph-level default size makes runs resolvable; the run itself stays local-free."""
    prs = _open(NORMAL)
    tf = _autofit_frame(prs)
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].line_spacing = 1.5
    tf.normalize_autofit()

    reopened_tf = _autofit_frame(save_reopen(prs))
    assert reopened_tf.paragraphs[0].font.size.pt == 12.5  # -- 20pt × 62.5%
    assert reopened_tf.paragraphs[0].runs[0].font.size is None  # -- still inherits from pPr


def test_min_font_size_floor_applies_after_scaling():
    prs = _open(NORMAL)
    tf = _autofit_frame(prs)
    tf.paragraphs[0].runs[0].font.size = Pt(12)  # -- 12 × 62.5% = 7.5pt, below the floor
    tf.paragraphs[0].line_spacing = 1.0
    tf.normalize_autofit(min_font_size=Pt(11))
    assert _autofit_frame(save_reopen(prs)).paragraphs[0].runs[0].font.size.pt == 11.0


def test_normalize_on_spautofit_changes_only_the_autofit_element():
    prs = _open(SHAPE)
    tf = _autofit_frame(prs)
    before = save_to_bytes(prs)
    tf.normalize_autofit()
    assert_changed_parts(before, save_to_bytes(prs), expect_changed=["ppt/slides/slide1.xml"])
    assert _autofit_frame(save_reopen(prs)).auto_size == MSO_AUTO_SIZE.NONE


def test_normalize_on_noautofit_is_a_complete_noop():
    prs = _open(NONE)
    before = save_to_bytes(prs)
    _autofit_frame(prs).normalize_autofit()
    assert_changed_parts(before, save_to_bytes(prs))  # -- empty budget: nothing changed


def test_normalize_with_bare_normautofit_drops_element_without_touching_text():
    """LibreOffice's attribute-less normAutofit: scale 100/reduction 0 → only the mode flips."""
    prs = Presentation(str(corpus.fixture_path(LO_NORMAL)))
    tf = _autofit_frame(prs)
    sizes_before = [run.font.size for run in tf.paragraphs[0].runs]
    before = save_to_bytes(prs)
    tf.normalize_autofit()
    assert_changed_parts(before, save_to_bytes(prs), expect_changed=["ppt/slides/slide1.xml"])
    reopened_tf = _autofit_frame(save_reopen(prs))
    assert reopened_tf.auto_size == MSO_AUTO_SIZE.NONE
    assert [run.font.size for run in reopened_tf.paragraphs[0].runs] == sizes_before


def test_normalize_autofit_resolve_freezes_inherited_placeholder_sizes():
    """PLAN-v0.1 0.4: the template-placeholder case. Without resolution the refusal stands;
    with resolve=True the sizes come from the master through the effective walk and are
    frozen at size x fontScale (2600->1625, 2200->1375 per the fixture sidecar)."""
    INHERITED = "self_generated/autofit_inherited.pptx"

    def body_tf(prs):
        return prs.slides[0].placeholders[1].text_frame

    # -- default behavior unchanged: refuses, atomically
    prs = _open(INHERITED)
    raised = assert_refusal_atomic(
        prs, lambda p: body_tf(p).normalize_autofit(), UnsupportedStructureError
    )
    assert "resolve" in str(raised)

    # -- resolve=True freezes what the reader sees
    prs = _open(INHERITED)
    before = save_to_bytes(prs)
    body_tf(prs).normalize_autofit(resolve=True)
    after = save_to_bytes(prs)
    assert_changed_parts(before, after, expect_changed=["ppt/slides/slide1.xml"])

    import io

    reopened_tf = body_tf(Presentation(io.BytesIO(after)))
    assert reopened_tf.auto_size == MSO_AUTO_SIZE.NONE
    assert reopened_tf.font_scale is None
    sizes = [p.runs[0].font.size.centipoints for p in reopened_tf.paragraphs]
    assert sizes == [1625, 1375]


def test_normalize_autofit_resolve_still_refuses_unresolvable_spacing():
    """resolve covers sizes only: line-spacing reduction with no explicit spacing refuses."""
    prs = _open(NORMAL)  # -- fixture carries lnSpcReduction=20000
    tf = _autofit_frame(prs)
    raised = assert_refusal_atomic(
        prs, lambda p: _autofit_frame(p).normalize_autofit(resolve=True),
        UnsupportedStructureError,
    )
    assert "line spacing" in str(raised) or "line-spacing" in str(raised)


def test_normalize_autofit_resolve_rejects_non_bool():
    prs = _open(NORMAL)
    with pytest.raises(ValueError):
        _autofit_frame(prs).normalize_autofit(resolve="yes")


def test_fields_are_scaled_and_validated_like_runs():
    """Regression: a:fld (slide number/date fields) render text exactly like runs; their
    explicit sizes must be frozen too, and an unsized field must refuse, not silently render
    at full size after normalization."""
    from lxml import etree

    _A = "http://schemas.openxmlformats.org/drawingml/2006/main"

    def add_field(tf, with_size):
        p = tf.paragraphs[0]._p
        fld = etree.SubElement(p, "{%s}fld" % _A)
        fld.set("id", "{11111111-2222-3333-4444-555555555555}")
        fld.set("type", "slidenum")
        if with_size:
            rPr = etree.SubElement(fld, "{%s}rPr" % _A)
            rPr.set("sz", "2000")
        t = etree.SubElement(fld, "{%s}t" % _A)
        t.text = "1"

    # -- unsized field refuses atomically --
    prs = _open(NORMAL)
    tf = _autofit_frame(prs)
    tf.paragraphs[0].runs[0].font.size = Pt(24)
    tf.paragraphs[0].line_spacing = 1.0
    add_field(tf, with_size=False)

    raised = assert_refusal_atomic(
        prs, lambda p: _autofit_frame(p).normalize_autofit(), UnsupportedStructureError
    )
    assert "field" in str(raised)

    # -- sized field is frozen along with the runs --
    prs = _open(NORMAL)
    tf = _autofit_frame(prs)
    tf.paragraphs[0].runs[0].font.size = Pt(24)
    tf.paragraphs[0].line_spacing = 1.0
    add_field(tf, with_size=True)
    tf.normalize_autofit()

    reopened_tf = _autofit_frame(save_reopen(prs))
    fld_rPr = reopened_tf.paragraphs[0]._p.find(
        "{%s}fld/{%s}rPr" % (_A, _A)
    )
    assert int(fld_rPr.get("sz")) == 1250  # -- 20pt × 62.5%


def test_upstream_auto_size_setter_behavior_is_unchanged():
    """§1.1 spot check: the existing property still round-trips all three modes + None."""
    prs = _open(NONE)
    tf = _autofit_frame(prs)
    for mode in (
        MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT,
        MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE,
        MSO_AUTO_SIZE.NONE,
        None,
    ):
        tf.auto_size = mode
        assert tf.auto_size == mode


@pytest.mark.lo_smoke
def test_normalized_output_loads_in_libreoffice(tmp_path):
    prs = _open(NORMAL)
    tf = _autofit_frame(prs)
    tf.paragraphs[0].runs[0].font.size = Pt(24)
    tf.paragraphs[0].line_spacing = 1.0
    tf.normalize_autofit(min_font_size=Pt(11))
    out = tmp_path / "normalized.pptx"
    prs.save(str(out))
    lo_load_smoke(out, tmp_path)
