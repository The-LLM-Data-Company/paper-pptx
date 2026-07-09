"""Contract tests: `paragraph.bullet` (real bullets and numbering).

Every mutation asserts through save -> reopen, budgets its changed parts,
and schema-validates the emitted `a:pPr` fragment. Bullets define no PaperRefusal
conditions — bad arguments are programmer errors — but every ValueError path is still proven
to leave the XML tree untouched (§1.3 spirit).
"""

from __future__ import annotations

import pytest
from lxml import etree

from pptx import Presentation
from pptx.enum.text import PP_BULLET_TYPE
from pptx.text.bullet import BulletFormat
from pptx.util import Emu

from . import corpus
from .contract import assert_changed_parts, save_reopen, save_to_bytes, snapshot_parts
from .fragval import _validation_errors, assert_pPr_fragment_valid
from .lo import lo_load_smoke

MINIMAL = "self_generated/minimal_clean.pptx"
_A = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _open_minimal():
    return Presentation(str(corpus.fixture_path(MINIMAL)))


def _add_box_with_text(prs, text="bullet me"):
    box = prs.slides[0].shapes.add_textbox(0, 0, Emu(914400 * 4), Emu(914400))
    box.text_frame.paragraphs[0].add_run().text = text
    return box


def _bullet_choice_children(paragraph):
    pPr = paragraph._p.pPr
    if pPr is None:
        return []
    tags = {"{%s}buNone" % _A, "{%s}buChar" % _A, "{%s}buAutoNum" % _A}
    return [child for child in pPr if child.tag in tags]


# ----------------------------------------------------------------------- reading local state


def test_fresh_paragraph_reports_fully_inherited_bullet_state():
    paragraph = _add_box_with_text(_open_minimal()).text_frame.paragraphs[0]
    bullet = paragraph.bullet
    assert isinstance(bullet, BulletFormat)
    assert bullet.type is None
    assert bullet.char is None
    assert bullet.number_scheme is None
    assert bullet.start_at is None
    assert bullet.font_name is None
    assert bullet.size_percent is None


def test_reads_real_bullet_from_frozen_gauntlet_fixture():
    """The gauntlet's real_bullet_box carries buFont Arial + buChar • (frozen real bytes)."""
    prs = Presentation(str(corpus.fixture_path("self_generated/gauntlet.pptx")))
    slide3 = prs.slides[2]
    box = next(shape for shape in slide3.shapes if shape.name == "real_bullet_box")
    bullet = box.text_frame.paragraphs[0].bullet
    assert bullet.type == PP_BULLET_TYPE.CHARACTER
    assert bullet.char == "•"
    assert bullet.font_name == "Arial"


def test_reads_legacy_thousandths_buszpct_form():
    """Real ECMA-376:2006-era files carry val="75000"; the reader must understand it."""
    paragraph = _add_box_with_text(_open_minimal()).text_frame.paragraphs[0]
    pPr = paragraph._p.get_or_add_pPr()
    buSzPct = etree.SubElement(pPr, "{%s}buSzPct" % _A)
    buSzPct.set("val", "75000")
    assert paragraph.bullet.size_percent == 0.75


def test_recognizes_picture_bullet_read_only():
    paragraph = _add_box_with_text(_open_minimal()).text_frame.paragraphs[0]
    pPr = paragraph._p.get_or_add_pPr()
    etree.SubElement(pPr, "{%s}buBlip" % _A)
    assert paragraph.bullet.type == PP_BULLET_TYPE.PICTURE


# --------------------------------------------------------------------------------- mutation


def test_set_character_round_trips_with_exact_part_budget():
    prs = _open_minimal()
    before = save_to_bytes(prs)
    box = _add_box_with_text(prs)
    paragraph = box.text_frame.paragraphs[0]
    paragraph.bullet.set_character(font_name="Arial", size_percent=0.75)
    assert_pPr_fragment_valid(paragraph)
    after = save_to_bytes(prs)
    assert_changed_parts(before, after, expect_changed=["ppt/slides/slide1.xml"])

    reopened = save_reopen(prs)
    shape = next(s for s in reopened.slides[0].shapes if s.shape_id == box.shape_id)
    bullet = shape.text_frame.paragraphs[0].bullet
    assert bullet.type == PP_BULLET_TYPE.CHARACTER
    assert bullet.char == "•"
    assert bullet.font_name == "Arial"
    assert bullet.size_percent == 0.75
    pPr = shape.text_frame.paragraphs[0]._p.pPr
    assert pPr.marL == 342900
    assert pPr.indent == -171450
    # -- the bullet is paragraph *property* state: no fake glyph enters the text itself
    assert shape.text_frame.text == "bullet me"


def test_set_numbered_round_trips():
    prs = _open_minimal()
    box = _add_box_with_text(prs)
    paragraph = box.text_frame.paragraphs[0]
    paragraph.bullet.set_numbered("romanUcParenR", start_at=3)
    assert_pPr_fragment_valid(paragraph)

    reopened = save_reopen(prs)
    shape = next(s for s in reopened.slides[0].shapes if s.shape_id == box.shape_id)
    bullet = shape.text_frame.paragraphs[0].bullet
    assert bullet.type == PP_BULLET_TYPE.NUMBERED
    assert bullet.number_scheme == "romanUcParenR"
    assert bullet.start_at == 3


def test_set_none_overrides_inherited_bullet_on_body_placeholder():
    """The branded template's body placeholder inherits master bullets; buNone beats them."""
    prs = Presentation(str(corpus.fixture_path("self_generated/branded_template.pptx")))
    paragraph = prs.slides[0].placeholders[1].text_frame.paragraphs[0]
    paragraph.bullet.set_none()
    assert_pPr_fragment_valid(paragraph)

    reopened = save_reopen(prs)
    reopened_paragraph = reopened.slides[0].placeholders[1].text_frame.paragraphs[0]
    assert reopened_paragraph.bullet.type == PP_BULLET_TYPE.NONE


def test_transitions_keep_exactly_one_bullet_choice_element():
    prs = _open_minimal()
    box = _add_box_with_text(prs)
    paragraph = box.text_frame.paragraphs[0]
    paragraph.bullet.set_character()
    paragraph.bullet.set_numbered()
    paragraph.bullet.set_none()
    paragraph.bullet.set_character("→")
    assert_pPr_fragment_valid(paragraph)

    reopened = save_reopen(prs)
    shape = next(s for s in reopened.slides[0].shapes if s.shape_id == box.shape_id)
    reopened_paragraph = shape.text_frame.paragraphs[0]
    assert len(_bullet_choice_children(reopened_paragraph)) == 1
    assert reopened_paragraph.bullet.type == PP_BULLET_TYPE.CHARACTER
    assert reopened_paragraph.bullet.char == "→"


def test_setters_replace_an_existing_picture_bullet():
    """Regression: buBlip is the fourth member of the schema's bullet choice; setters used to
    install a second bullet element beside it, producing schema-invalid XML."""
    prs = _open_minimal()
    box = _add_box_with_text(prs)
    paragraph = box.text_frame.paragraphs[0]
    pPr = paragraph._p.get_or_add_pPr()
    etree.SubElement(pPr, "{%s}buBlip" % _A)
    assert paragraph.bullet.type == PP_BULLET_TYPE.PICTURE

    paragraph.bullet.set_character()
    assert_pPr_fragment_valid(paragraph)

    reopened = save_reopen(prs)
    shape = next(s for s in reopened.slides[0].shapes if s.shape_id == box.shape_id)
    bullet = shape.text_frame.paragraphs[0].bullet
    assert bullet.type == PP_BULLET_TYPE.CHARACTER
    assert len(_bullet_choice_children(shape.text_frame.paragraphs[0])) == 1


def test_none_margins_leave_existing_attributes_untouched():
    prs = _open_minimal()
    box = _add_box_with_text(prs)
    paragraph = box.text_frame.paragraphs[0]
    paragraph.bullet.set_character(left_margin=Emu(500000), hanging_indent=Emu(250000))
    paragraph.bullet.set_numbered(left_margin=None, hanging_indent=None)

    reopened = save_reopen(prs)
    shape = next(s for s in reopened.slides[0].shapes if s.shape_id == box.shape_id)
    pPr = shape.text_frame.paragraphs[0]._p.pPr
    assert pPr.marL == 500000
    assert pPr.indent == -250000


def test_set_none_touches_no_margins():
    prs = _open_minimal()
    box = _add_box_with_text(prs)
    box.text_frame.paragraphs[0].bullet.set_none()

    reopened = save_reopen(prs)
    shape = next(s for s in reopened.slides[0].shapes if s.shape_id == box.shape_id)
    pPr = shape.text_frame.paragraphs[0]._p.pPr
    assert pPr.marL is None
    assert pPr.indent is None


# ------------------------------------------------------------------- validation is atomic


@pytest.mark.parametrize(
    "bad_call",
    [
        lambda b: b.set_character(""),
        lambda b: b.set_character(42),
        lambda b: b.set_character(size_percent=9.0),
        lambda b: b.set_character(size_percent=0.1),
        lambda b: b.set_character(left_margin="wide"),
        lambda b: b.set_numbered("notAScheme"),
        lambda b: b.set_numbered(start_at=0),
        lambda b: b.set_numbered(start_at=True),
        lambda b: b.set_character(left_margin=Emu(99999999999)),
        lambda b: b.set_numbered(font_name=42),
        lambda b: b.set_character(size_percent=0.755),  # -- schema admits whole percents only
        lambda b: b.set_character("bad\udc80char"),  # -- lone surrogate: unserializable
        lambda b: b.set_numbered(font_name="bad\udc80font"),
    ],
)
def test_bad_arguments_raise_valueerror_and_leave_the_tree_untouched(bad_call):
    prs = _open_minimal()
    paragraph = _add_box_with_text(prs).text_frame.paragraphs[0]
    before = snapshot_parts(prs)
    with pytest.raises(ValueError):
        bad_call(paragraph.bullet)
    assert snapshot_parts(prs) == before


# ------------------------------------------------------------------------- fragment oracle


def test_fragment_oracle_rejects_malformed_bullet_xml():
    """Failing direction of the oracle: required-attribute and child-order violations."""
    no_char = etree.fromstring('<a:pPr xmlns:a="%s"><a:buChar/></a:pPr>' % _A)
    assert _validation_errors(no_char, "pPr", "CT_TextParagraphProperties")
    bad_order = etree.fromstring(
        '<a:pPr xmlns:a="%s"><a:buChar char="x"/><a:buFont typeface="Arial"/></a:pPr>' % _A
    )
    assert _validation_errors(bad_order, "pPr", "CT_TextParagraphProperties")


def test_upstream_written_paragraphs_also_pass_the_fragment_oracle():
    """Sanity: paragraphs upstream code writes (no bullets) validate too."""
    prs = _open_minimal()
    for shape in prs.slides[0].shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                assert_pPr_fragment_valid(paragraph)


# ------------------------------------------------------------------------------ lo_smoke


@pytest.mark.lo_smoke
def test_bulleted_output_loads_in_libreoffice(tmp_path):
    prs = _open_minimal()
    paragraph = _add_box_with_text(prs).text_frame.paragraphs[0]
    paragraph.bullet.set_character(font_name="Arial")
    out = tmp_path / "bullets_lo.pptx"
    prs.save(str(out))
    lo_load_smoke(out, tmp_path)
