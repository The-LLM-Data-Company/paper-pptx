"""Self-tests for the contract harness: prove the harness catches what it exists to catch.

A harness that silently passes on violations is worse than no harness - every organ test to
come leans on these assertions. Assertions 1 (save->reopen), 3 (changed-part budget),
4 (independent-loader smoke), and 5 (refusal atomicity, both halves) are exercised here in
both their passing and failing directions; assertion 2 (intended effect present) is a
caller-side assert on the reopened document and has no harness utility to self-test.
"""

from __future__ import annotations

import io
import zipfile

import pytest

from pptx import Presentation

from . import corpus
from .clock import PAPER_TEST_INSTANT, FrozenClock
from .contract import (
    assert_changed_parts,
    assert_file_bytes_unchanged,
    assert_refusal_atomic,
    diff_zip_members,
    save_reopen,
    save_to_bytes,
    zip_member_map,
)
from .idlists import dangling_section_slide_ids, duplicate_section_slide_ids
from .lo import lo_load_smoke, soffice_path
from .relint import dangling_relationship_targets, missing_relationship_references

MINIMAL = "self_generated/minimal_clean.pptx"


def _open_minimal():
    return Presentation(str(corpus.fixture_path(MINIMAL)))


def _zip_bytes(members):
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w") as zipf:
        for name, data in members.items():
            zipf.writestr(name, data)
    return buf.getvalue()


# ----------------------------------------------------------------- save_reopen (assertion 1)


def test_save_reopen_returns_a_fresh_presentation_with_the_same_content():
    prs = _open_minimal()
    reopened = save_reopen(prs)
    assert reopened is not prs
    assert reopened.slides[0].shapes.title.text == "Minimal clean fixture"


def test_save_reopen_sees_only_what_reached_the_saved_bytes():
    prs = _open_minimal()
    prs.slides[0].shapes.title.text_frame.paragraphs[0].runs[0].text = "Edited title"
    assert save_reopen(prs).slides[0].shapes.title.text == "Edited title"


# ------------------------------------------------------- changed-part budget (assertion 3)


def test_double_save_of_one_presentation_is_member_identical():
    """Two saves of the same in-memory package must differ in ZERO members.

    This is load-bearing for every budget assertion: it proves serialization is deterministic
    at member granularity and that save() stamps no volatile content (e.g. core-properties
    timestamps) into member bytes. If this ever fails, budgets are unusable as written.
    """
    prs = _open_minimal()
    diff = diff_zip_members(save_to_bytes(prs), save_to_bytes(prs))
    assert diff.is_empty, diff.describe()


def test_diff_pinpoints_a_single_changed_member():
    prs = _open_minimal()
    before = save_to_bytes(prs)
    prs.slides[0].shapes.title.text_frame.paragraphs[0].runs[0].text = "Edited title"
    after = save_to_bytes(prs)
    diff = diff_zip_members(before, after)
    assert diff.changed == ("ppt/slides/slide1.xml",)
    assert diff.added == ()
    assert diff.removed == ()


def test_diff_reports_added_and_removed_members():
    a = _zip_bytes({"x/common.xml": b"<a/>", "x/removed.xml": b"<b/>"})
    b = _zip_bytes({"x/common.xml": b"<a/>", "x/added.xml": b"<c/>"})
    diff = diff_zip_members(a, b)
    assert diff.added == ("x/added.xml",)
    assert diff.removed == ("x/removed.xml",)
    assert diff.changed == ()


def test_zip_member_map_refuses_duplicate_member_names():
    """OPC forbids duplicate part names; a dict would silently keep one copy and the diff
    would report 'no change' on the clone-forgot-to-rename corruption class."""
    import warnings

    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w") as zipf:
        zipf.writestr("ppt/media/image1.png", b"<orig/>")
        with warnings.catch_warnings():
            # -- suppress zipfile's "Duplicate name" UserWarning: constructing the corrupt
            # -- package is the whole point of this test
            warnings.simplefilter("ignore")
            zipf.writestr("ppt/media/image1.png", b"<CLONE-GARBAGE/>")
    with pytest.raises(AssertionError, match="duplicate zip member names"):
        zip_member_map(buf.getvalue())


def test_assert_changed_parts_accepts_an_exact_budget():
    prs = _open_minimal()
    before = save_to_bytes(prs)
    prs.slides[0].shapes.title.text_frame.paragraphs[0].runs[0].text = "Edited title"
    assert_changed_parts(before, save_to_bytes(prs), expect_changed=["ppt/slides/slide1.xml"])


def test_assert_changed_parts_rejects_unexpected_changes():
    prs = _open_minimal()
    before = save_to_bytes(prs)
    prs.slides[0].shapes.title.text_frame.paragraphs[0].runs[0].text = "Edited title"
    after = save_to_bytes(prs)
    with pytest.raises(AssertionError, match="changed-part budget violated"):
        assert_changed_parts(before, after)  # -- empty budget must fail
    with pytest.raises(AssertionError, match="changed-part budget violated"):
        # -- an expected-but-absent change must fail too
        assert_changed_parts(
            before, after, expect_changed=["ppt/slides/slide1.xml", "ppt/presentation.xml"]
        )


# --------------------------------------------------------- refusal atomicity (assertion 5)


class _FakeRefusal(Exception):
    pass


def test_refusal_atomic_passes_for_a_clean_refusal():
    def clean_refusal(prs):
        raise _FakeRefusal("refused before touching anything")

    raised = assert_refusal_atomic(_open_minimal(), clean_refusal, _FakeRefusal)
    assert "refused before touching" in str(raised)


def test_refusal_atomic_catches_mutate_then_raise():
    def dirty_refusal(prs):
        prs.slides[0].shapes.title.text_frame.paragraphs[0].runs[0].text = "mutated!"
        raise _FakeRefusal("refused after mutating - the forbidden pattern")

    with pytest.raises(AssertionError, match="left the in-memory package dirty"):
        assert_refusal_atomic(_open_minimal(), dirty_refusal, _FakeRefusal)


def test_refusal_atomic_catches_a_missing_refusal():
    with pytest.raises(pytest.fail.Exception):  # -- pytest.raises' DID NOT RAISE failure
        assert_refusal_atomic(_open_minimal(), lambda prs: None, _FakeRefusal)


def test_refusal_atomic_does_not_swallow_wrong_exception_types():
    def wrong_type(prs):
        raise ValueError("a bug, not a refusal")

    with pytest.raises(ValueError, match="a bug"):
        assert_refusal_atomic(_open_minimal(), wrong_type, _FakeRefusal)


def test_file_bytes_unchanged_passes_when_untouched(tmp_path):
    target = tmp_path / "input.pptx"
    target.write_bytes(corpus.fixture_path(MINIMAL).read_bytes())
    with assert_file_bytes_unchanged(target):
        pass


def test_file_bytes_unchanged_catches_a_write(tmp_path):
    target = tmp_path / "input.pptx"
    target.write_bytes(corpus.fixture_path(MINIMAL).read_bytes())
    with pytest.raises(AssertionError, match="was modified by a refused operation"):
        with assert_file_bytes_unchanged(target):
            target.write_bytes(b"clobbered")


def test_file_bytes_unchanged_checks_even_when_the_operation_raises(tmp_path):
    """The natural composition wraps a refused (raising) operation; the byte checks must run
    anyway. The AssertionError supersedes the in-flight refusal - loud is correct."""
    target = tmp_path / "input.pptx"
    target.write_bytes(corpus.fixture_path(MINIMAL).read_bytes())
    with pytest.raises(AssertionError, match="was modified by a refused operation"):
        with assert_file_bytes_unchanged(target):
            target.write_bytes(b"clobbered")
            raise _FakeRefusal("refused after clobbering the input")


def test_file_bytes_unchanged_lets_a_clean_refusal_propagate(tmp_path):
    target = tmp_path / "input.pptx"
    target.write_bytes(corpus.fixture_path(MINIMAL).read_bytes())
    with pytest.raises(_FakeRefusal):
        with assert_file_bytes_unchanged(target):
            raise _FakeRefusal("clean refusal, file untouched")


def test_refusal_atomic_catches_a_content_type_mutation():
    """[Content_Types].xml is generated from part content types at save time, so a mutated
    content type is save-visible state the snapshot must cover."""

    def dirty_refusal(prs):
        part = next(iter(prs.part.package.iter_parts()))
        part._content_type = "application/x-mutated"
        raise _FakeRefusal("refused after flipping a content type")

    with pytest.raises(AssertionError, match="left the in-memory package dirty"):
        assert_refusal_atomic(_open_minimal(), dirty_refusal, _FakeRefusal)


def test_save_reopen_reads_disk_truth_not_live_state(monkeypatch):
    """save_reopen must reflect only what reached the saved bytes: with save() pinned to
    stale bytes, a later in-memory edit must be invisible in the reopened presentation."""
    prs = _open_minimal()
    stale = save_to_bytes(prs)
    monkeypatch.setattr(prs, "save", lambda file: file.write(stale))
    prs.slides[0].shapes.title.text_frame.paragraphs[0].runs[0].text = "never reached disk"
    assert save_reopen(prs).slides[0].shapes.title.text == "Minimal clean fixture"


# -------------------------------------------------------------- relationship-integrity scans


_RELS_XMLNS = 'xmlns="http://schemas.openxmlformats.org/package/2006/relationships"'


def test_relint_detects_a_dangling_relationship_target():
    zip_map = {
        "ppt/slides/slide1.xml": b"<x/>",
        "ppt/slides/_rels/slide1.xml.rels": (
            '<Relationships %s><Relationship Id="rId1" Type="t" Target="../media/gone.png"/>'
            "</Relationships>" % _RELS_XMLNS
        ).encode(),
    }
    assert dangling_relationship_targets(zip_map) == [
        ("ppt/slides/_rels/slide1.xml.rels", "rId1", "ppt/media/gone.png")
    ]


def test_relint_ignores_external_targets():
    zip_map = {
        "ppt/slides/slide1.xml": b"<x/>",
        "ppt/slides/_rels/slide1.xml.rels": (
            '<Relationships %s><Relationship Id="rId1" Type="t" '
            'Target="https://example.com/x" TargetMode="External"/></Relationships>'
            % _RELS_XMLNS
        ).encode(),
    }
    assert dangling_relationship_targets(zip_map) == []


def test_relint_detects_a_missing_relationship_reference():
    r_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    zip_map = {
        "ppt/slides/slide1.xml": (
            '<x xmlns:r="%s"><child r:embed="rId7"/></x>' % r_ns
        ).encode(),
        "ppt/slides/_rels/slide1.xml.rels": (
            "<Relationships %s></Relationships>" % _RELS_XMLNS
        ).encode(),
    }
    assert missing_relationship_references(zip_map) == [("ppt/slides/slide1.xml", "rId7")]


def test_relint_accepts_legal_package_root_relative_targets():
    """Targets like "/ppt/media/x.png" are legal OPC and accepted by python-pptx's loader;
    the scan must not flag them as dangling."""
    zip_map = {
        "ppt/media/image1.png": b"png-bytes",
        "ppt/slides/slide1.xml": b"<x/>",
        "ppt/slides/_rels/slide1.xml.rels": (
            '<Relationships %s><Relationship Id="rId1" Type="t" '
            'Target="/ppt/media/image1.png"/></Relationships>' % _RELS_XMLNS
        ).encode(),
    }
    assert dangling_relationship_targets(zip_map) == []


def test_relint_reports_a_target_less_relationship_instead_of_crashing():
    zip_map = {
        "ppt/slides/slide1.xml": b"<x/>",
        "ppt/slides/_rels/slide1.xml.rels": (
            '<Relationships %s><Relationship Id="rId1" Type="t"/></Relationships>'
            % _RELS_XMLNS
        ).encode(),
    }
    assert dangling_relationship_targets(zip_map) == [
        ("ppt/slides/_rels/slide1.xml.rels", "rId1", "<missing Target attribute>")
    ]


def test_idlists_detects_a_dangling_section_slide_id():
    p = "http://schemas.openxmlformats.org/presentationml/2006/main"
    p14 = "http://schemas.microsoft.com/office/powerpoint/2010/main"
    zip_map = {
        "ppt/presentation.xml": (
            '<p:presentation xmlns:p="%s" xmlns:p14="%s">'
            '<p:sldIdLst><p:sldId id="256"/></p:sldIdLst>'
            '<p:extLst><p:ext uri="{521415D9-36F7-43E2-AB2F-B90AF26B5E84}">'
            '<p14:sectionLst><p14:section name="Ghost" id="{00000000-0000-4000-8000-000000000000}">'
            '<p14:sldIdLst><p14:sldId id="999"/><p14:sldId id="256"/></p14:sldIdLst>'
            "</p14:section></p14:sectionLst></p:ext></p:extLst></p:presentation>" % (p, p14)
        ).encode(),
    }
    assert dangling_section_slide_ids(zip_map) == [("Ghost", "999")]
    assert duplicate_section_slide_ids(zip_map) == []


def test_idlists_detects_an_id_enrolled_in_two_sections():
    p = "http://schemas.openxmlformats.org/presentationml/2006/main"
    p14 = "http://schemas.microsoft.com/office/powerpoint/2010/main"
    section = (
        '<p14:section name="%s" id="%s"><p14:sldIdLst><p14:sldId id="256"/></p14:sldIdLst>'
        "</p14:section>"
    )
    zip_map = {
        "ppt/presentation.xml": (
            '<p:presentation xmlns:p="%s" xmlns:p14="%s">'
            '<p:sldIdLst><p:sldId id="256"/></p:sldIdLst>'
            '<p:extLst><p:ext uri="u"><p14:sectionLst>%s%s</p14:sectionLst></p:ext></p:extLst>'
            "</p:presentation>"
            % (
                p,
                p14,
                section % ("A", "{11111111-1111-4111-8111-111111111111}"),
                section % ("B", "{22222222-2222-4222-8222-222222222222}"),
            )
        ).encode(),
    }
    assert duplicate_section_slide_ids(zip_map) == ["256"]


def test_idlists_pass_on_sectionless_and_presentationless_maps():
    assert dangling_section_slide_ids({"x.xml": b"<x/>"}) == []
    minimal = zip_member_map(corpus.fixture_path(MINIMAL).read_bytes())
    assert dangling_section_slide_ids(minimal) == []
    assert duplicate_section_slide_ids(minimal) == []


# ------------------------------------------------------------------- corpus utility guards


def test_manifest_entries_rejects_duplicate_relpaths():
    duplicated = (
        "aaaa  self_generated/foo.pptx\n"
        "bbbb  self_generated/foo.pptx\n"
    )
    with pytest.raises(ValueError, match="more than once"):
        corpus.manifest_entries(duplicated)


def _valid_sidecar_doc():
    return {
        "fixture": "x.pptx",
        "provenance": {"app": "a", "version": "v", "notes": "n"},
        "features": ["f"],
        "ground_truth": {"k": 1},
        "verified_by": "someone",
        "date": "2026-07-07",
    }


def test_sidecar_validator_accepts_a_valid_doc():
    assert corpus.sidecar_schema_problems(_valid_sidecar_doc(), "bucket/x.pptx") == []


@pytest.mark.parametrize(
    "date", ["2026-13-99", "0000-00-00", "2026-01-02\n", "2026-02-30", "not-a-date"]
)
def test_sidecar_validator_rejects_bogus_dates(date):
    doc = _valid_sidecar_doc()
    doc["date"] = date
    problems = corpus.sidecar_schema_problems(doc, "bucket/x.pptx")
    assert any("date" in p for p in problems)


def test_sidecar_validator_rejects_whitespace_only_features():
    doc = _valid_sidecar_doc()
    doc["features"] = ["   "]
    problems = corpus.sidecar_schema_problems(doc, "bucket/x.pptx")
    assert any("features" in p for p in problems)


# ------------------------------------------------------------------------------ frozen clock


def test_frozen_clock_reports_the_pinned_instant_every_time(frozen_clock):
    assert frozen_clock.now() == PAPER_TEST_INSTANT
    assert frozen_clock.now() == frozen_clock.now()
    assert frozen_clock() == PAPER_TEST_INSTANT  # -- callable-style injection


def test_frozen_clock_accepts_a_custom_instant():
    from datetime import datetime, timezone

    instant = datetime(2030, 6, 7, 8, 9, 10, tzinfo=timezone.utc)
    assert FrozenClock(instant).now() == instant


# --------------------------------------------------------------- lo_smoke oracle (assertion 4)


def test_soffice_probe_returns_none_or_an_existing_executable():
    path = soffice_path()
    if path is not None:
        import os

        assert os.path.isfile(path)


@pytest.mark.lo_smoke
def test_lo_smoke_helper_accepts_a_known_good_file(tmp_path):
    lo_load_smoke(corpus.fixture_path(MINIMAL), tmp_path)


@pytest.mark.lo_smoke
def test_lo_smoke_helper_rejects_non_presentation_bytes(tmp_path):
    """The failing direction of the oracle: without a pinned import filter LibreOffice's
    plain-text fallback 'converts' arbitrary garbage to a valid PDF with exit code 0."""
    garbage = tmp_path / "garbage.pptx"
    garbage.write_bytes(b"this is not a zip file at all " * 100)
    with pytest.raises(AssertionError, match="could not load"):
        lo_load_smoke(garbage, tmp_path)
