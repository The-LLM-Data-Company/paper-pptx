"""Phase 8 contract tests: geometry-preserving image replacement.

Position, size, and crop must survive byte-exactly; the extension-mismatch refusal is atomic;
shared image parts are only orphaned when the last reference goes. The reference's low-res /
natural-size math appears here as test assertions (per the plan: never public API).
"""

from __future__ import annotations

import io

import pytest
from lxml import etree
from PIL import Image as PILImage

from pptx import Presentation
from pptx.errors import PaperRefusal, UnsupportedStructureError

from . import corpus
from .contract import (
    assert_changed_parts,
    assert_refusal_atomic,
    save_to_bytes,
    zip_member_map,
)
from .lo import lo_load_smoke
from .relint import dangling_relationship_targets, missing_relationship_references

GAUNTLET = "self_generated/gauntlet.pptx"
SHARED_MEDIA = "self_generated/shared_media.pptx"


def _open(relpath):
    return Presentation(str(corpus.fixture_path(relpath)))


def _png_bytes(color=(9, 99, 199), size=(32, 32)):
    buf = io.BytesIO()
    PILImage.new("RGB", size, color).save(buf, format="PNG")
    return buf.getvalue()


def _jpg_bytes():
    buf = io.BytesIO()
    PILImage.new("RGB", (32, 32), (5, 5, 5)).save(buf, format="JPEG")
    return buf.getvalue()


def _cropped_picture(prs):
    return next(s for s in prs.slides[2].shapes if s.name == "gauntlet_cropped")


# -------------------------------------------------------------------------------- replacing


def test_replace_preserves_position_size_and_crop_exactly():
    prs = _open(GAUNTLET)
    picture = _cropped_picture(prs)
    geometry_before = (
        picture.left,
        picture.top,
        picture.width,
        picture.height,
        picture.rotation,
        picture.crop_left,
        picture.crop_top,
        picture.crop_right,
        picture.crop_bottom,
    )
    new_bytes = _png_bytes()
    picture.replace_image(io.BytesIO(new_bytes))

    reopened = Presentation(io.BytesIO(save_to_bytes(prs)))
    reopened_picture = _cropped_picture(reopened)
    geometry_after = (
        reopened_picture.left,
        reopened_picture.top,
        reopened_picture.width,
        reopened_picture.height,
        reopened_picture.rotation,
        reopened_picture.crop_left,
        reopened_picture.crop_top,
        reopened_picture.crop_right,
        reopened_picture.crop_bottom,
    )
    assert geometry_after == geometry_before
    assert reopened_picture.image.blob == new_bytes


def test_replace_has_exact_part_budget_and_keeps_shared_original():
    """The gauntlet image part is shared by three pictures; replacing one must not disturb
    the other two, and the budget is exactly: the slide, the new part, the slide rels."""
    prs = _open(GAUNTLET)
    before = save_to_bytes(prs)
    _cropped_picture(prs).replace_image(io.BytesIO(_png_bytes()))
    after = save_to_bytes(prs)
    assert_changed_parts(
        before,
        after,
        expect_changed=["ppt/slides/_rels/slide3.xml.rels", "ppt/slides/slide3.xml"],
        expect_added=["ppt/media/image2.png"],
    )
    zip_map = zip_member_map(after)
    assert dangling_relationship_targets(zip_map) == []
    assert missing_relationship_references(zip_map) == []

    reopened = Presentation(io.BytesIO(after))
    untouched = next(s for s in reopened.slides[1].shapes if s.name == "gauntlet_img_1")
    assert untouched.image.blob != _png_bytes()


def test_old_image_part_is_orphaned_only_when_last_reference_goes():
    prs = _open(SHARED_MEDIA)
    prs.slides[0].shapes[0].replace_image(io.BytesIO(_png_bytes((1, 2, 3))))
    two_parts = zip_member_map(save_to_bytes(prs))
    assert len([n for n in two_parts if n.startswith("ppt/media/")]) == 2

    prs.slides[1].shapes[0].replace_image(io.BytesIO(_png_bytes((1, 2, 3))))
    one_part = zip_member_map(save_to_bytes(prs))
    assert len([n for n in one_part if n.startswith("ppt/media/")]) == 1
    assert dangling_relationship_targets(one_part) == []
    assert missing_relationship_references(one_part) == []


def test_replacing_with_identical_bytes_is_a_complete_noop():
    prs = _open(SHARED_MEDIA)
    original_bytes = prs.slides[0].shapes[0].image.blob
    before = save_to_bytes(prs)
    prs.slides[0].shapes[0].replace_image(io.BytesIO(original_bytes))
    assert_changed_parts(before, save_to_bytes(prs))  # -- empty budget


# ---------------------------------------------------------------------- refusals, atomically


def test_extension_mismatch_refuses_atomically():
    prs = _open(GAUNTLET)
    raised = assert_refusal_atomic(
        prs,
        lambda p: _cropped_picture(p).replace_image(io.BytesIO(_jpg_bytes())),
        UnsupportedStructureError,
    )
    assert "does not match" in str(raised)
    assert isinstance(raised, PaperRefusal)


def test_linked_only_picture_refuses_atomically():
    prs = _open(GAUNTLET)
    picture = _cropped_picture(prs)
    blip = picture._pic.blipFill.blip
    del blip.attrib[
        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
    ]

    def operation(p):
        _cropped_picture(p).replace_image(io.BytesIO(_png_bytes()))

    raised = assert_refusal_atomic(prs, operation, UnsupportedStructureError)
    assert "no embedded image" in str(raised)


def test_dangling_image_relationship_refuses_atomically():
    prs = _open(GAUNTLET)
    picture = _cropped_picture(prs)
    picture._pic.blipFill.blip.rEmbed = "rId99"

    def operation(p):
        _cropped_picture(p).replace_image(io.BytesIO(_png_bytes()))

    raised = assert_refusal_atomic(prs, operation, UnsupportedStructureError)
    assert "does not exist" in str(raised)


def test_unrecognizable_image_bytes_raise_valueerror_atomically():
    prs = _open(GAUNTLET)
    before = save_to_bytes(prs)
    with pytest.raises(ValueError, match="not a recognizable image"):
        _cropped_picture(prs).replace_image(io.BytesIO(b"this is not an image at all"))
    assert_changed_parts(before, save_to_bytes(prs))  # -- empty budget


def test_cross_format_swap_with_allow_format_change(tmp_path):
    """v0.1 2.3: PNG -> JPEG swap; geometry and crop byte-untouched; content types follow."""
    prs = _open(GAUNTLET)
    picture = _cropped_picture(prs)
    spPr_before = etree.tostring(picture._pic.spPr)
    crop_before = (picture.crop_left, picture.crop_top)

    from PIL import Image as PILImage

    jpeg = io.BytesIO()
    PILImage.new("RGB", (32, 32), (200, 100, 50)).save(jpeg, format="JPEG")
    picture.replace_image(io.BytesIO(jpeg.getvalue()), allow_format_change=True)

    saved = save_to_bytes(prs)
    zip_map = zip_member_map(saved)
    assert any(n.startswith("ppt/media/") and n.endswith(".jpg") for n in zip_map)
    assert b"jpeg" in zip_map["[Content_Types].xml"] or b"jpg" in (
        zip_map["[Content_Types].xml"]
    )
    assert dangling_relationship_targets(zip_map) == []
    assert missing_relationship_references(zip_map) == []

    reopened = Presentation(io.BytesIO(saved))
    reopened_pic = _cropped_picture(reopened)
    assert etree.tostring(reopened_pic._pic.spPr) == spPr_before
    assert (reopened_pic.crop_left, reopened_pic.crop_top) == crop_before
    assert reopened_pic.image.ext in ("jpg", "jpeg")


def test_cross_format_swap_still_refuses_by_default():
    prs = _open(GAUNTLET)
    from PIL import Image as PILImage

    jpeg = io.BytesIO()
    PILImage.new("RGB", (16, 16), (1, 2, 3)).save(jpeg, format="JPEG")

    def operation(p):
        _cropped_picture(p).replace_image(io.BytesIO(jpeg.getvalue()))

    raised = assert_refusal_atomic(prs, operation, UnsupportedStructureError)
    assert "allow_format_change" in str(raised)


def test_allow_format_change_rejects_non_bool():
    prs = _open(GAUNTLET)
    with pytest.raises(ValueError):
        _cropped_picture(prs).replace_image(io.BytesIO(_png_bytes()), allow_format_change=1)


def test_replace_keeps_relationship_alive_for_sibling_picture_sharing_it():
    """Regression: two pictures added from identical bytes share ONE relationship; replacing
    one used to drop the rel and leave the sibling's r:embed dangling (unloadable deck)."""
    prs = _open("self_generated/minimal_clean.pptx")
    slide = prs.slides[0]
    same_bytes = _png_bytes((40, 80, 120))
    pic_a = slide.shapes.add_picture(io.BytesIO(same_bytes), 0, 0, 914400)
    pic_b = slide.shapes.add_picture(io.BytesIO(same_bytes), 914400, 0, 914400)
    assert pic_a._pic.blip_rId == pic_b._pic.blip_rId  # -- the shared-rel precondition

    replacement = _png_bytes((200, 10, 10))
    pic_a.replace_image(io.BytesIO(replacement))
    saved = save_to_bytes(prs)
    zip_map = zip_member_map(saved)
    assert dangling_relationship_targets(zip_map) == []
    assert missing_relationship_references(zip_map) == []

    reopened = Presentation(io.BytesIO(saved))
    pictures = [s for s in reopened.slides[0].shapes if s.shape_type.name == "PICTURE"]
    blobs = sorted(p.image.blob for p in pictures)
    assert blobs == sorted([same_bytes, replacement])  # -- sibling still loads its image


# -------------------------------------------- reference low-res math, as assertions only


def test_effective_resolution_math_on_frozen_fixture():
    """The reference's low-res detection: effective DPI = pixel size / displayed inches.

    The shared-media fixture displays a 64px image at 2 inches -> 32 DPI, far below any
    print threshold; the math itself is the asset here, kept as a regression assertion.
    """
    prs = _open(SHARED_MEDIA)
    picture = prs.slides[0].shapes[0]
    pixel_width, pixel_height = picture.image.size
    dpi_x = pixel_width / (picture.width.inches)
    dpi_y = pixel_height / (picture.height.inches)
    assert (pixel_width, pixel_height) == (64, 64)
    assert round(dpi_x) == 32
    assert round(dpi_y) == 32
    assert dpi_x < 144  # -- the reference's default low-res threshold flags this image


# --------------------------------------------------------------------------------- lo_smoke


@pytest.mark.lo_smoke
def test_replaced_image_output_loads_in_libreoffice(tmp_path):
    prs = _open(GAUNTLET)
    _cropped_picture(prs).replace_image(io.BytesIO(_png_bytes()))
    out = tmp_path / "replaced.pptx"
    prs.save(str(out))
    lo_load_smoke(out, tmp_path)
