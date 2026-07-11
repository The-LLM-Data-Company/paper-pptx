"""Contract tests: layout rebind, the template-migration primitive.

The report is the differentiator and is required: the resolver runs before and after, and
every run whose resolved values changed appears with before/after payloads — appearance
never shifts silently. Cross-template rebinds (two-template corpus) are exercised through
the adopt_theme mode, which is built on this machinery; here the shift-reporting
proof uses the gauntlet's branded layout overrides (present on the source layout, absent
on the target).
"""

from __future__ import annotations

import io

import pytest
from lxml import etree

from pptx import Presentation
from pptx.errors import PaperRefusal, UnsupportedStructureError
from pptx.util import Emu

from . import corpus
from .contract import (
    assert_changed_parts,
    assert_refusal_atomic,
    save_reopen,
    save_to_bytes,
)
from .lo import lo_load_smoke
from .relint import dangling_relationship_targets, missing_relationship_references

GAUNTLET = "self_generated/gauntlet.pptx"
MINIMAL = "self_generated/minimal_clean.pptx"
BETA = "self_generated/template_beta.pptx"


def _open(relpath):
    return Presentation(str(corpus.fixture_path(relpath)))


# ------------------------------------------------------------------------------- happy path


def test_rebind_within_template_reports_every_real_shift_and_nothing_else():
    """Gauntlet slide 1 sits on the branded 'Title and Content' (layout title override
    36pt). 'Two Content' has no title override (title re-resolves to the master's 40pt)
    but DOES carry its own body-size lstStyle overrides (28pt/24pt) — all three runs
    genuinely shift and all three must be reported, with exact before/after values."""
    prs = _open(GAUNTLET)
    report = prs.slides[0].rebind_layout(prs.slide_layouts[3])  # -- "Two Content"

    shifts = {s.text: s for s in report.run_shifts}
    assert set(shifts) == {
        "Gauntlet: branded",
        "Inherited body level one",
        "Inherited body level two",
    }
    title = shifts["Gauntlet: branded"]
    assert title.before["size"]["value"] == 457200  # -- 36pt from the layout override
    assert title.after["size"]["value"] == 508000  # -- 40pt from the branded master
    body_one = shifts["Inherited body level one"]
    assert body_one.before["size"]["value"] == 330200  # -- 26pt branded master
    assert body_one.after["size"]["value"] == 355600  # -- 28pt Two Content override
    body_two = shifts["Inherited body level two"]
    assert body_two.before["size"]["value"] == 279400  # -- 22pt branded master
    assert body_two.after["size"]["value"] == 304800  # -- 24pt Two Content override
    assert report.source_layout_name == "Title and Content"
    assert report.target_layout_name == "Two Content"
    assert report.baked_orphans == ()

    reopened = save_reopen(prs)
    assert reopened.slides[0].slide_layout.name == "Two Content"


def test_rebind_has_exact_part_budget_and_relint_clean():
    prs = _open(GAUNTLET)
    before = save_to_bytes(prs)
    prs.slides[0].rebind_layout(prs.slide_layouts[3])
    after = save_to_bytes(prs)
    assert_changed_parts(
        before,
        after,
        expect_changed=["ppt/slides/_rels/slide1.xml.rels"],
    )
    from .contract import zip_member_map

    zip_map = zip_member_map(after)
    assert dangling_relationship_targets(zip_map) == []
    assert missing_relationship_references(zip_map) == []


def test_family_matching_rewrites_placeholder_type_and_idx():
    """minimal_clean's ctrTitle/subTitle must family-match title/object on 'Title and
    Content', with the slide ph elements rewritten so inheritance binds by idx."""
    prs = _open(MINIMAL)
    report = prs.slides[0].rebind_layout(prs.slide_layouts[1])
    assert report.placeholder_map_used == ((0, 0), (1, 1))

    reopened = save_reopen(prs)
    ph_facts = {
        s.element.ph_idx: s.element.ph_type.name
        for s in reopened.slides[0].shapes
        if s.is_placeholder
    }
    assert ph_facts == {0: "TITLE", 1: "OBJECT"}
    assert reopened.slides[0].slide_layout.name == "Title and Content"


def test_rebind_report_is_deterministic():
    prs_a = _open(GAUNTLET)
    prs_b = _open(GAUNTLET)
    report_a = prs_a.slides[0].rebind_layout(prs_a.slide_layouts[3]).to_dict()
    report_b = prs_b.slides[0].rebind_layout(prs_b.slide_layouts[3]).to_dict()
    assert report_a == report_b


# ------------------------------------------------------------------------------ orphans


def test_orphans_refuse_atomically_by_default():
    prs = _open(GAUNTLET)
    raised = assert_refusal_atomic(
        prs,
        lambda p: p.slides[0].rebind_layout(p.slide_layouts[5]),  # -- "Title Only"
        UnsupportedStructureError,
    )
    assert "Content Placeholder 2" in str(raised)
    assert "orphan_policy" in str(raised)
    assert isinstance(raised, PaperRefusal)


def test_bake_orphan_keeps_look_and_becomes_free_shape():
    prs = _open(GAUNTLET)
    body = next(s for s in prs.slides[0].shapes if s.element.ph_idx == 1)
    inherited_geometry = (body.left, body.top, body.width, body.height)

    report = prs.slides[0].rebind_layout(prs.slide_layouts[5], orphan_policy="bake")
    assert report.baked_orphans == ("Content Placeholder 2",)
    # -- baked runs re-resolve to identical values: no shift entries for them
    assert [s.text for s in report.run_shifts] == ["Gauntlet: branded"]

    reopened = save_reopen(prs)
    baked = next(s for s in reopened.slides[0].shapes if s.name == "Content Placeholder 2")
    assert not baked.is_placeholder
    assert (baked.left, baked.top, baked.width, baked.height) == inherited_geometry
    run = baked.text_frame.paragraphs[0].runs[0]
    assert run.font.size == Emu(330200)  # -- the branded master's 26pt, now local
    assert run.font.name == "Trebuchet MS"


def test_bake_geometry_falls_back_through_inheritance():
    """The baked orphan's xfrm is materialized from the resolved inheritance chain."""
    prs = _open(GAUNTLET)
    body = next(s for s in prs.slides[0].shapes if s.element.ph_idx == 1)
    assert body._element.spPr.find(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm"
    ) is None  # -- precondition: geometry lives on the layout, not the slide
    prs.slides[0].rebind_layout(prs.slide_layouts[5], orphan_policy="bake")
    baked = next(s for s in prs.slides[0].shapes if s.name == "Content Placeholder 2")
    assert baked._element.spPr.find(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm"
    ) is not None


def test_bake_refuses_field_bearing_placeholder():
    """Cross-feature: apply_footers writes a slidenum field; force-orphaning that
    placeholder under bake must refuse (a baked field freezes volatile content)."""
    prs = _open(GAUNTLET)
    prs.slides[0].apply_footers(slide_number=True)

    def operation(p):
        p.slides[0].rebind_layout(
            p.slide_layouts[3], placeholder_map={12: None}, orphan_policy="bake"
        )

    raised = assert_refusal_atomic(prs, operation, UnsupportedStructureError)
    assert "a:fld" in str(raised)


def test_bake_localizes_transformed_scheme_color_without_losing_transform():
    prs = _open(GAUNTLET)
    body = next(s for s in prs.slides[0].shapes if s.element.ph_idx == 1)
    run = body.text_frame.paragraphs[0].runs[0]
    rPr = run._r.get_or_add_rPr()
    solidFill = etree.SubElement(
        rPr, "{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill"
    )
    schemeClr = etree.SubElement(
        solidFill, "{http://schemas.openxmlformats.org/drawingml/2006/main}schemeClr"
    )
    schemeClr.set("val", "accent1")
    etree.SubElement(
        schemeClr, "{http://schemas.openxmlformats.org/drawingml/2006/main}lumMod"
    ).set("val", "50000")

    prs.slides[0].rebind_layout(prs.slide_layouts[5], orphan_policy="bake")
    reopened = save_reopen(prs)
    baked = next(
        s for s in reopened.slides[0].shapes if s.name == "Content Placeholder 2"
    )
    color = baked.text_frame.paragraphs[0].runs[0]._r.get_or_add_rPr().find(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill"
    )[0]
    assert color.tag.endswith("srgbClr")
    assert color.find(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}lumMod"
    ).get("val") == "50000"


def test_post_mutation_report_failure_rolls_rebind_back(monkeypatch):
    import pptx.rebind as rebind_module

    prs = _open(GAUNTLET)
    original = rebind_module._resolution_state
    calls = 0

    def fail_after_mutation(slide):
        nonlocal calls
        calls += 1
        if calls == 2:
            raise UnsupportedStructureError("forced report failure")
        return original(slide)

    monkeypatch.setattr(rebind_module, "_resolution_state", fail_after_mutation)
    raised = assert_refusal_atomic(
        prs,
        lambda p: p.slides[0].rebind_layout(p.slide_layouts[3]),
        UnsupportedStructureError,
    )
    assert "forced report failure" in str(raised)


# --------------------------------------------------------------------------- explicit map


def test_explicit_map_overrides_auto():
    prs = _open(GAUNTLET)
    # -- map the body ph onto Two Content's SECOND content slot instead of the first
    report = prs.slides[0].rebind_layout(prs.slide_layouts[3], placeholder_map={1: 2})
    assert (1, 2) in report.placeholder_map_used
    reopened = save_reopen(prs)
    ph_idx_set = {
        s.element.ph_idx for s in reopened.slides[0].shapes if s.is_placeholder
    }
    assert ph_idx_set == {0, 2}


def test_explicit_map_validation():
    prs = _open(GAUNTLET)
    slide = prs.slides[0]
    target = prs.slide_layouts[3]
    with pytest.raises(ValueError, match="not a placeholder on this slide"):
        slide.rebind_layout(target, placeholder_map={99: 1})
    with pytest.raises(ValueError, match="not a placeholder on the target"):
        slide.rebind_layout(target, placeholder_map={1: 99})
    with pytest.raises(ValueError, match="one target"):
        slide.rebind_layout(target, placeholder_map={0: 1, 1: 1})
    with pytest.raises(ValueError, match="placeholder_map"):
        slide.rebind_layout(target, placeholder_map=[("a", "b")])


def test_exact_matches_settle_before_lower_idx_placeholders_can_steal_slots():
    """Regression: auto-match runs exact type+idx as a GLOBAL first
    pass. Previously a lower-idx placeholder fell through to its family fallback and
    claimed a higher-idx placeholder's exact slot before that placeholder was even
    considered - orphaning the rightful owner."""
    import copy

    from pptx.enum.shapes import PP_PLACEHOLDER
    from pptx.rebind import _compute_mapping

    prs = _open(GAUNTLET)
    slide = prs.slides[0]
    target = prs.slide_layouts[3]  # -- "Two Content": TITLE@0, OBJECT@1, OBJECT@2
    # -- make the target's OBJECT@1 slot a CHART slot: OBJECT@2 is now the only content
    # -- slot, and it is ph@2's EXACT match
    target_obj1 = next(
        ph for ph in target.placeholders if ph.element.ph_idx == 1
    )
    target_obj1.element.ph.type = PP_PLACEHOLDER.CHART

    # -- slide: BODY@1 (family-matches OBJECT only) plus an OBJECT@2 copy (exact match)
    body_ph = next(s for s in slide.shapes if s.element.ph_idx == 1)
    body_ph.element.ph.type = PP_PLACEHOLDER.BODY
    duplicate = copy.deepcopy(body_ph._element)
    duplicate.nvSpPr.cNvPr.set("id", "97")
    duplicate.nvSpPr.cNvPr.set("name", "Rightful Owner")
    duplicate.ph.type = PP_PLACEHOLDER.OBJECT
    duplicate.ph.idx = 2
    body_ph._element.addnext(duplicate)

    slide_phs = [s for s in slide.shapes if s.is_placeholder]
    mapping = _compute_mapping(slide_phs, target, "auto")
    # -- ph@2 keeps its exact OBJECT@2 slot; the BODY ph@1 has nowhere to go and
    # -- orphans honestly (previously it stole slot 2 and ph@2 was orphaned instead)
    assert mapping[2] == (PP_PLACEHOLDER.OBJECT, 2)
    assert mapping[1] is None
    assert mapping[0][1] == 0


# ------------------------------------------------------------------------------- refusals


def test_bad_target_arguments():
    prs = _open(GAUNTLET)
    other = _open(BETA)
    slide = prs.slides[0]
    with pytest.raises(ValueError, match="SlideLayout"):
        slide.rebind_layout("Title Only")
    with pytest.raises(ValueError, match="different presentation"):
        slide.rebind_layout(other.slide_layouts[1])
    with pytest.raises(ValueError, match="already this slide's layout"):
        slide.rebind_layout(prs.slides[0].slide_layout)
    with pytest.raises(ValueError, match="orphan_policy"):
        slide.rebind_layout(prs.slide_layouts[3], orphan_policy="improvise")


def test_alternate_content_slide_refuses_atomically():
    prs = _open(GAUNTLET)
    spTree = prs.slides[0].shapes._spTree
    spTree.append(
        etree.SubElement(
            spTree,
            "{http://schemas.openxmlformats.org/markup-compatibility/2006}AlternateContent",
        )
    )
    raised = assert_refusal_atomic(
        prs,
        lambda p: p.slides[0].rebind_layout(p.slide_layouts[3]),
        UnsupportedStructureError,
    )
    assert "AlternateContent" in str(raised)


# --------------------------------------------------------------------------------- lo_smoke


@pytest.mark.lo_smoke
def test_rebound_slide_loads_in_libreoffice(tmp_path):
    prs = _open(GAUNTLET)
    prs.slides[0].rebind_layout(prs.slide_layouts[5], orphan_policy="bake")
    out = tmp_path / "rebound.pptx"
    prs.save(str(out))
    lo_load_smoke(out, tmp_path)


def test_rebound_output_reopens_and_text_intact():
    prs = _open(GAUNTLET)
    prs.slides[0].rebind_layout(prs.slide_layouts[5], orphan_policy="bake")
    reopened = Presentation(io.BytesIO(save_to_bytes(prs)))
    texts = sorted(
        s.text_frame.text for s in reopened.slides[0].shapes if s.has_text_frame
    )
    assert "Gauntlet: branded" in texts
    assert any("Inherited body level one" in t for t in texts)
