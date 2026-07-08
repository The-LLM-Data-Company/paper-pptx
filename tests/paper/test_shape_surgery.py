"""Phase 1.2/1.3 contract tests: shape surgery + generic by-name addressing (SlideShapes)."""

from __future__ import annotations

import io

import pytest

from pptx import Presentation
from pptx.errors import (
    AmbiguousTargetError,
    RelationshipPolicyError,
    TargetNotFoundError,
)
from pptx.util import Emu

from . import corpus
from .contract import (
    assert_changed_parts,
    assert_refusal_atomic,
    save_to_bytes,
    zip_member_map,
)
from .lo import lo_load_smoke
from .relint import dangling_relationship_targets, missing_relationship_references

GAUNTLET = "self_generated/gauntlet.pptx"
CHART_NOTES = "self_generated/chart_notes.pptx"
NESTED = "self_generated/nested_groups.pptx"


def _open(relpath):
    return Presentation(str(corpus.fixture_path(relpath)))


def _reopen(pptx_bytes):
    return Presentation(io.BytesIO(pptx_bytes))


def _assert_relationship_integrity(pptx_bytes):
    zip_map = zip_member_map(pptx_bytes)
    assert dangling_relationship_targets(zip_map) == []
    assert missing_relationship_references(zip_map) == []


# ----------------------------------------------------------------------------------- delete


def test_delete_textbox_round_trips_with_exact_budget():
    prs = _open(GAUNTLET)
    before = save_to_bytes(prs)
    slide3 = prs.slides[2]
    slide3.shapes.delete(slide3.shapes.shape_by_name("real_bullet_box"))
    after = save_to_bytes(prs)
    assert_changed_parts(before, after, expect_changed=["ppt/slides/slide3.xml"])
    reopened = _reopen(after)
    with pytest.raises(TargetNotFoundError):
        reopened.slides[2].shapes.shape_by_name("real_bullet_box")


def test_delete_chart_shape_drops_the_chart_parts_from_the_package():
    prs = _open(CHART_NOTES)
    slide = prs.slides[0]
    slide.shapes.delete(next(s for s in slide.shapes if s.has_chart))
    saved = save_to_bytes(prs)
    zip_map = zip_member_map(saved)
    assert not any(n.startswith("ppt/charts/") for n in zip_map)
    assert not any(n.startswith("ppt/embeddings/") for n in zip_map)
    _assert_relationship_integrity(saved)


def test_delete_keeps_relationship_shared_with_a_sibling_picture():
    """The replace_image bug class: deleting one of two pictures sharing a rel must not
    orphan the survivor's r:embed."""
    prs = _open("self_generated/minimal_clean.pptx")
    slide = prs.slides[0]
    from PIL import Image as PILImage

    png = io.BytesIO()
    PILImage.new("RGB", (8, 8), (5, 5, 5)).save(png, format="PNG")
    pic_a = slide.shapes.add_picture(io.BytesIO(png.getvalue()), 0, 0, Emu(914400))
    pic_b = slide.shapes.add_picture(io.BytesIO(png.getvalue()), Emu(914400), 0, Emu(914400))
    assert pic_a._pic.blip_rId == pic_b._pic.blip_rId

    slide.shapes.delete(pic_a)
    saved = save_to_bytes(prs)
    _assert_relationship_integrity(saved)
    survivor = next(
        s for s in _reopen(saved).slides[0].shapes if s.shape_type.name == "PICTURE"
    )
    assert survivor.image.blob == png.getvalue()


def test_delete_refuses_foreign_and_grouped_shapes_atomically():
    prs = _open(NESTED)
    slide = prs.slides[0]
    grouped_box = slide.shapes.shape_by_name("level1_box")  # -- inside a group

    def delete_grouped(p):
        p.slides[0].shapes.delete(grouped_box)

    raised = assert_refusal_atomic(prs, delete_grouped, TargetNotFoundError)
    assert "group" in str(raised)

    other = _open(GAUNTLET)
    foreign = other.slides[2].shapes.shape_by_name("real_bullet_box")
    assert_refusal_atomic(
        prs, lambda p: p.slides[0].shapes.delete(foreign), TargetNotFoundError
    )


# ------------------------------------------------------------------------------------- move


def test_move_changes_z_order_and_round_trips():
    prs = _open(GAUNTLET)
    slide3 = prs.slides[2]
    names_before = [s.name for s in slide3.shapes]
    top = slide3.shapes.shape_by_name("autofit_none_box")
    before = save_to_bytes(prs)
    slide3.shapes.move(top, len(names_before) - 1)
    after = save_to_bytes(prs)
    assert_changed_parts(before, after, expect_changed=["ppt/slides/slide3.xml"])
    reopened_names = [s.name for s in _reopen(after).slides[2].shapes]
    assert reopened_names[-1] == "autofit_none_box"
    assert sorted(reopened_names) == sorted(names_before)


def test_move_to_front_of_z_order():
    prs = _open(GAUNTLET)
    slide3 = prs.slides[2]
    last = slide3.shapes[-1]
    slide3.shapes.move(last, 0)
    assert slide3.shapes[0].name == last.name


def test_move_rejects_bad_indices_and_foreign_shapes():
    prs = _open(GAUNTLET)
    slide3 = prs.slides[2]
    shape = slide3.shapes[0]
    with pytest.raises(ValueError):
        slide3.shapes.move(shape, 99)
    with pytest.raises(ValueError):
        slide3.shapes.move(shape, -1)
    with pytest.raises(ValueError):
        slide3.shapes.move(shape, True)
    other = _open(NESTED)
    with pytest.raises(TargetNotFoundError):
        slide3.shapes.move(other.slides[0].shapes[0], 0)


# --------------------------------------------------------------------------------- add_copy


def test_add_copy_of_chart_shape_deep_copies_chart_and_workbook():
    prs = _open(CHART_NOTES)
    source = next(s for s in prs.slides[0].shapes if s.has_chart)
    before = save_to_bytes(prs)
    copy = prs.slides[0].shapes.add_copy(source)
    assert copy.shape_id != source.shape_id
    saved = save_to_bytes(prs)
    zip_map = zip_member_map(saved)
    assert "ppt/charts/chart2.xml" in zip_map
    assert len([n for n in zip_map if n.startswith("ppt/embeddings/")]) == 2
    _assert_relationship_integrity(saved)

    # -- cross-contamination: editing the copy leaves the original chart XML untouched
    reopened = _reopen(saved)
    charts = [s for s in reopened.slides[0].shapes if s.has_chart]
    assert len(charts) == 2
    original_chart_xml = zip_map["ppt/charts/chart1.xml"]
    charts[1].chart.replace_data_safe(["X"], [("New", (9.0,))])
    assert zip_member_map(save_to_bytes(reopened))["ppt/charts/chart1.xml"] == (
        original_chart_xml
    )
    del before  # -- budget for add_copy is covered by the member assertions above


def test_add_copy_shares_image_parts():
    prs = _open(GAUNTLET)
    slide2 = prs.slides[1]
    picture = slide2.shapes.picture_by_name("gauntlet_img_1")
    prs.slides[3].shapes.add_copy(picture)
    saved = save_to_bytes(prs)
    zip_map = zip_member_map(saved)
    assert [n for n in zip_map if n.startswith("ppt/media/")] == ["ppt/media/image1.png"]
    _assert_relationship_integrity(saved)


def test_add_copy_copies_external_hyperlinks():
    prs = _open(GAUNTLET)
    slide4 = prs.slides[3]
    link_box = slide4.shapes.shape_by_name("hyperlink_box")
    copy = prs.slides[0].shapes.add_copy(link_box)
    address = copy.text_frame.paragraphs[0].runs[0].hyperlink.address
    assert address == "https://example.com/paper"
    _assert_relationship_integrity(save_to_bytes(prs))


def test_add_copy_refuses_unsupported_relationship_types_atomically():
    prs = _open(GAUNTLET)
    slide3 = prs.slides[2]
    box = slide3.shapes.shape_by_name("autofit_none_box")
    # -- give the shape a reference to an unsupported (diagram-style) relationship
    image_pic = prs.slides[1].shapes.picture_by_name("gauntlet_img_1")
    target_part = prs.slides[1].part.related_part(image_pic._pic.blip_rId)
    rId = slide3.part.relate_to(target_part, "http://example.com/relationships/diagramData")
    box._element.spPr.set(
        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}dm", rId
    )

    def do_copy(p):
        p.slides[0].shapes.add_copy(box)

    raised = assert_refusal_atomic(prs, do_copy, RelationshipPolicyError)
    assert "does not support" in str(raised)


def test_add_copy_refuses_shapes_from_another_presentation_atomically():
    prs = _open(GAUNTLET)
    other = _open(CHART_NOTES)
    foreign = next(s for s in other.slides[0].shapes if s.has_chart)

    raised = assert_refusal_atomic(
        prs, lambda p: p.slides[0].shapes.add_copy(foreign), TargetNotFoundError
    )
    assert "different presentation" in str(raised)


def test_add_copy_refuses_chart_with_unsupported_child_relationship():
    """Regression (review, CRITICAL): add_copy used to skip the chart child-relationship
    validation Slides.clone performs, silently leaf-copying a chart whose child rels the
    copy cannot honor — dangling rIds in the output."""
    prs = _open(CHART_NOTES)
    chart_shape = next(s for s in prs.slides[0].shapes if s.has_chart)
    chart_part = chart_shape.chart.part
    image_part = prs.slides[0].part.package.get_or_add_image_part(
        io.BytesIO(
            Presentation(str(corpus.fixture_path(GAUNTLET)))
            .slides[1]
            .shapes.picture_by_name("gauntlet_img_1")
            .image.blob
        )
    )
    chart_part.relate_to(image_part, "http://example.com/relationships/bogus")

    raised = assert_refusal_atomic(
        prs, lambda p: p.slides[0].shapes.add_copy(chart_shape), RelationshipPolicyError
    )
    assert "chart part" in str(raised)


# ------------------------------------------------------------------------ by-name addressing


def test_by_name_finds_shapes_inside_nested_groups():
    prs = _open(NESTED)
    shapes = prs.slides[0].shapes
    assert shapes.shape_by_name("level3_box").name == "level3_box"


def test_by_name_refusals():
    prs = _open(GAUNTLET)
    shapes3 = prs.slides[2].shapes
    with pytest.raises(TargetNotFoundError, match="no shape named"):
        shapes3.shape_by_name("does_not_exist")
    with pytest.raises(TargetNotFoundError, match="not a picture"):
        shapes3.picture_by_name("real_bullet_box")
    with pytest.raises(TargetNotFoundError, match="not a table"):
        shapes3.table_by_name("real_bullet_box")

    box_a = prs.slides[2].shapes.shape_by_name("autofit_none_box")
    box_a.name = "duplicate_name"
    box_b = prs.slides[2].shapes.shape_by_name("autofit_shape_box")
    box_b.name = "duplicate_name"
    with pytest.raises(AmbiguousTargetError):
        shapes3.shape_by_name("duplicate_name")


def test_table_by_name_returns_the_table_even_inside_a_group():
    prs = _open("self_generated/tables_in_group.pptx")
    table = prs.slides[0].shapes.table_by_name("grouped_table")
    assert table.cell(0, 0).text == "cell r0c0"


def test_chart_by_name_is_group_aware_since_v01():
    """chart_by_name gained group traversal with the 1.3 wave (ledgered in PAPER.md)."""
    prs = _open(CHART_NOTES)
    chart_frame = next(s for s in prs.slides[0].shapes if s.has_chart)
    group = prs.slides[0].shapes.add_group_shape()
    group._element.append(chart_frame._element)  # -- move the chart into the group
    chart = prs.slides[0].shapes.chart_by_name("clone_fixture_chart")
    assert [series.name for series in chart.series] == ["Q1", "Q2"]


# ------------------------------------------------------------------------------- lo_smoke


@pytest.mark.lo_smoke
def test_surgery_output_loads_in_libreoffice(tmp_path):
    prs = _open(GAUNTLET)
    slide3 = prs.slides[2]
    slide3.shapes.delete(slide3.shapes.shape_by_name("real_bullet_box"))
    slide3.shapes.move(slide3.shapes[0], len(list(slide3.shapes)) - 1)
    picture = prs.slides[1].shapes.picture_by_name("gauntlet_img_1")
    prs.slides[3].shapes.add_copy(picture)
    out = tmp_path / "surgery.pptx"
    prs.save(str(out))
    lo_load_smoke(out, tmp_path)
