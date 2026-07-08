"""Phase 5 contract tests: the package kernel (`pptx.package` additive functions).

Every CONVENTIONS §7 pinned requirement is exercised here: the meaningful-whitespace trap
(via the frozen trailing-space fixture pair), no-op byte-identity, exact single-part budgets,
zip determinism (pinned entry order and epoch timestamps), and mid-write crash atomicity.
"""

from __future__ import annotations

import zipfile

import pytest

from pptx import Presentation
from pptx.errors import PaperRefusal, UnsupportedStructureError
from pptx.package import PackageDiff, diff_package, patch_save, xml_equivalent

from . import corpus
from .contract import snapshot_parts, zip_member_map
from .lo import lo_load_smoke

MINIMAL = "self_generated/minimal_clean.pptx"
GAUNTLET = "self_generated/gauntlet.pptx"
PAIR_A = "self_generated/whitespace_trailing_a.pptx"
PAIR_B = "self_generated/whitespace_trailing_b.pptx"
LO_MINIMAL = "libreoffice_export/lo_minimal_clean.pptx"


def _fixture(relpath):
    return str(corpus.fixture_path(relpath))


def _slide1_xml(relpath):
    return zip_member_map(corpus.fixture_path(relpath).read_bytes())["ppt/slides/slide1.xml"]


# -------------------------------------------------------------------------- xml_equivalent


def test_trailing_space_in_a_t_is_never_equivalent():
    """THE §7 trap test, on the frozen fixture pair differing by one trailing space."""
    assert not xml_equivalent(_slide1_xml(PAIR_A), _slide1_xml(PAIR_B))


def test_whitespace_only_text_of_childless_element_is_preserved():
    assert not xml_equivalent("<a><t> </t></a>", "<a><t/></a>")
    assert not xml_equivalent("<a><t> </t></a>", "<a><t>  </t></a>")


def test_structural_indentation_is_equivalent():
    """Pretty-printed vs compact (the LibreOffice-vs-python-pptx serialization difference)."""
    assert xml_equivalent('<a>\n  <b x="1">text </b>\n</a>', '<a><b x="1">text </b></a>')


def test_attribute_order_and_prefix_spelling_are_equivalent():
    assert xml_equivalent(
        '<r:x xmlns:r="urn:n" a="1" b="2">t </r:x>',
        '<q:x xmlns:q="urn:n" b="2" a="1">t </q:x>',
    )


def test_element_order_is_significant():
    assert not xml_equivalent("<a><b/><c/></a>", "<a><c/><b/></a>")


def test_malformed_xml_raises_valueerror():
    with pytest.raises(ValueError):
        xml_equivalent("not xml at all", "<a/>")


def test_identical_files_compare_equivalent_part_by_part():
    members = zip_member_map(corpus.fixture_path(MINIMAL).read_bytes())
    for name, blob in members.items():
        if name.endswith(".xml") or name.endswith(".rels"):
            assert xml_equivalent(blob, blob)


# --------------------------------------------------------------------------- diff_package


def test_diff_reports_exactly_the_trailing_space_part():
    diff = diff_package(_fixture(PAIR_A), _fixture(PAIR_B))
    assert [d.partname for d in diff.deltas] == ["/ppt/slides/slide1.xml"]
    assert diff.deltas[0].kind == "xml"
    assert diff.deltas[0].change == "changed"


def test_diff_of_a_package_with_itself_is_empty():
    diff = diff_package(_fixture(GAUNTLET), _fixture(GAUNTLET))
    assert diff.is_empty


def test_diff_reports_added_removed_and_binary_changes(tmp_path):
    source = zip_member_map(corpus.fixture_path(MINIMAL).read_bytes())
    modified = dict(source)
    del modified["docProps/thumbnail.jpeg"]
    modified["ppt/media/new.png"] = b"png-bytes"
    modified["docProps/app.xml"] = source["docProps/app.xml"].replace(
        b"</Properties>", b"\n  </Properties>"
    )

    original_path, modified_path = tmp_path / "a.pptx", tmp_path / "b.pptx"
    for path, members in ((original_path, source), (modified_path, modified)):
        with zipfile.ZipFile(str(path), "w") as zipf:
            for name, data in members.items():
                zipf.writestr(name, data)

    diff = diff_package(str(original_path), str(modified_path))
    by_part = {d.partname: d for d in diff.deltas}
    assert by_part["/docProps/thumbnail.jpeg"].change == "removed"
    assert by_part["/ppt/media/new.png"].change == "added"
    assert by_part["/ppt/media/new.png"].kind == "binary"
    # -- the app.xml tweak only adds structural whitespace: semantically equivalent
    assert "/docProps/app.xml" not in by_part
    assert len(diff.deltas) == 2


def test_diff_to_dict_carries_pinned_schema_and_is_deterministic():
    diff = diff_package(_fixture(PAIR_A), _fixture(PAIR_B))
    payload = diff.to_dict()
    assert payload["schema"] == "paper-package-diff"
    assert payload["version"] == 1
    assert payload == diff_package(_fixture(PAIR_A), _fixture(PAIR_B)).to_dict()


# ------------------------------------------------------------------------------ patch_save


@pytest.mark.parametrize("relpath", [MINIMAL, GAUNTLET])
def test_noop_round_trip_is_byte_identical(relpath, tmp_path):
    out = tmp_path / "noop.pptx"
    diff = patch_save(_fixture(relpath), Presentation(_fixture(relpath)), str(out))
    assert diff.is_empty
    assert out.read_bytes() == corpus.fixture_path(relpath).read_bytes()


def test_noop_on_a_libreoffice_file_restores_every_part_but_content_types(tmp_path):
    """Real-third-party round trip: all parts restore; the one residual delta is genuine
    (LibreOffice declares Default content types for extensions with no parts — fntdata,
    jpeg, png — and regeneration from live parts drops them)."""
    out = tmp_path / "lo_noop.pptx"
    diff = patch_save(_fixture(LO_MINIMAL), Presentation(_fixture(LO_MINIMAL)), str(out))
    assert [d.partname for d in diff.deltas] == ["/[Content_Types].xml"]


def test_single_part_edit_changes_exactly_that_part(tmp_path):
    prs = Presentation(_fixture(MINIMAL))
    prs.slides[0].shapes.title.text_frame.paragraphs[0].runs[0].text = "Edited title"
    out = tmp_path / "edit.pptx"
    diff = patch_save(_fixture(MINIMAL), prs, str(out))
    assert [d.partname for d in diff.deltas] == ["/ppt/slides/slide1.xml"]
    assert Presentation(str(out)).slides[0].shapes.title.text == "Edited title"


def test_patch_save_never_restores_over_a_trailing_space_edit(tmp_path):
    """§3's corruption-inside-the-safety-tooling scenario: removing a meaningful trailing
    space IS the edit; restore logic must not 'fix' it back to the original bytes."""
    prs = Presentation(_fixture(PAIR_A))
    box = next(s for s in prs.slides[0].shapes if s.name == "whitespace_box")
    box.text_frame.paragraphs[0].runs[0].text = "Trailing space"  # -- space removed
    out = tmp_path / "ws_edit.pptx"
    diff = patch_save(_fixture(PAIR_A), prs, str(out))
    assert [d.partname for d in diff.deltas] == ["/ppt/slides/slide1.xml"]

    reopened_box = next(
        s for s in Presentation(str(out)).slides[0].shapes if s.name == "whitespace_box"
    )
    assert reopened_box.text_frame.paragraphs[0].runs[0].text == "Trailing space"


def test_patch_save_output_is_deterministic_with_pinned_order_and_timestamps(tmp_path):
    def build():
        prs = Presentation(_fixture(MINIMAL))
        prs.slides[0].shapes.title.text_frame.paragraphs[0].runs[0].text = "Deterministic"
        return prs

    out1, out2 = tmp_path / "det1.pptx", tmp_path / "det2.pptx"
    patch_save(_fixture(MINIMAL), build(), str(out1))
    patch_save(_fixture(MINIMAL), build(), str(out2))
    assert out1.read_bytes() == out2.read_bytes()

    with zipfile.ZipFile(str(out1)) as zipf:
        names = zipf.namelist()
        assert names[0] == "[Content_Types].xml"
        assert names[1] == "_rels/.rels"
        assert names[2:] == sorted(names[2:])
        assert all(info.date_time == (1980, 1, 1, 0, 0, 0) for info in zipf.infolist())


def test_mid_write_crash_leaves_existing_output_intact(tmp_path, monkeypatch):
    """§7 failure-injection: the original survives a crash halfway through writing."""
    out = tmp_path / "out.pptx"
    precious = b"precious existing bytes"
    out.write_bytes(precious)

    prs = Presentation(_fixture(MINIMAL))
    prs.slides[0].shapes.title.text_frame.paragraphs[0].runs[0].text = "Crash"

    calls = {"count": 0}
    real_writestr = zipfile.ZipFile.writestr

    def exploding_writestr(self, *args, **kwargs):
        calls["count"] += 1
        if calls["count"] >= 5:
            raise OSError("simulated disk failure")
        return real_writestr(self, *args, **kwargs)

    monkeypatch.setattr(zipfile.ZipFile, "writestr", exploding_writestr)
    with pytest.raises(OSError, match="simulated disk failure"):
        patch_save(_fixture(MINIMAL), prs, str(out))
    monkeypatch.undo()

    assert out.read_bytes() == precious
    assert list(tmp_path.glob("*.partial")) == []  # -- no temp litter


def test_patch_save_refuses_a_non_zip_original_before_touching_anything(tmp_path):
    bogus = tmp_path / "notazip.pptx"
    bogus.write_bytes(b"this is not a zip " * 10)
    out = tmp_path / "never.pptx"
    prs = Presentation(_fixture(MINIMAL))
    before = snapshot_parts(prs)

    with pytest.raises(UnsupportedStructureError) as exc_info:
        patch_save(str(bogus), prs, str(out))
    assert isinstance(exc_info.value, PaperRefusal)
    assert not out.exists()
    assert snapshot_parts(prs) == before


def test_patch_save_rejects_a_document_that_cannot_save(tmp_path):
    with pytest.raises(ValueError):
        patch_save(_fixture(MINIMAL), object(), str(tmp_path / "x.pptx"))


def test_packagediff_type_shape():
    diff = diff_package(_fixture(MINIMAL), _fixture(MINIMAL))
    assert isinstance(diff, PackageDiff)
    assert diff.is_empty
    assert diff.to_dict()["deltas"] == []


@pytest.mark.lo_smoke
def test_patch_saved_output_loads_in_libreoffice(tmp_path):
    prs = Presentation(_fixture(MINIMAL))
    prs.slides[0].shapes.title.text_frame.paragraphs[0].runs[0].text = "LO check"
    out = tmp_path / "patched.pptx"
    patch_save(_fixture(MINIMAL), prs, str(out))
    lo_load_smoke(out, tmp_path)
