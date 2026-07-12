"""Contract tests: real text fields (a:fld) and header/footer flags (p:hf)."""

from __future__ import annotations

import io
import re

import pytest
from lxml import etree

from pptx import Presentation
from pptx.errors import TargetNotFoundError
from pptx.inspect import inspect_text
from pptx.util import Emu

from . import corpus
from .contract import assert_changed_parts, save_to_bytes
from .fragval import _validation_errors
from .lo import lo_load_smoke

MINIMAL = "self_generated/minimal_clean.pptx"
_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
_GUID_RE = re.compile(r"^\{[0-9A-F]{8}-[0-9A-F]{4}-[0-9A-F]{4}-[0-9A-F]{4}-[0-9A-F]{12}\}$")


def _open(relpath):
    return Presentation(str(corpus.fixture_path(relpath)))


def _reopen(pptx_bytes):
    return Presentation(io.BytesIO(pptx_bytes))


def _footer_box(prs):
    box = prs.slides[0].shapes.add_textbox(0, 0, Emu(914400 * 3), Emu(914400 // 2))
    box.name = "footer_box"
    return box


def test_header_footer_flags_refuse_after_layout_removal():
    prs = Presentation()
    layout = prs.slide_layouts[0]
    master = layout.slide_master
    flags = layout.header_footers
    master.slide_layouts.remove(layout)

    with pytest.raises(TargetNotFoundError, match="stale"):
        flags.footer_visible = False


# ------------------------------------------------------------------------------ a:fld fields


def test_field_setter_refuses_a_deleted_paragraph_proxy():
    prs = _open(MINIMAL)
    slide = prs.slides[0]
    shape = _footer_box(prs)
    paragraph = shape.text_frame.paragraphs[0]
    slide.shapes.delete(shape)

    with pytest.raises(TargetNotFoundError, match="paragraph is stale"):
        paragraph.add_slide_number_field()


def test_field_setter_refuses_a_paragraph_on_a_removed_layout():
    prs = Presentation()
    master = prs.slide_masters[0]
    layout = master.slide_layouts[1]
    paragraph = next(
        shape for shape in layout.shapes if shape.has_text_frame
    ).text_frame.paragraphs[0]
    master.slide_layouts.remove(layout)

    with pytest.raises(TargetNotFoundError, match="paragraph is stale"):
        paragraph.add_slide_number_field()


def test_slide_number_field_round_trips_with_exact_budget():
    prs = _open(MINIMAL)
    before = save_to_bytes(prs)
    paragraph = _footer_box(prs).text_frame.paragraphs[0]
    paragraph.add_run().text = "Page "
    paragraph.add_slide_number_field()
    after = save_to_bytes(prs)
    assert_changed_parts(before, after, expect_changed=["ppt/slides/slide1.xml"])

    reopened = _reopen(after)
    box = reopened.slides[0].shapes.shape_by_name("footer_box")
    fld = box.text_frame.paragraphs[0]._p.find("{%s}fld" % _A)
    assert fld is not None
    assert fld.get("type") == "slidenum"
    assert _GUID_RE.match(fld.get("id"))
    assert fld.find("{%s}t" % _A).text == "1"


def test_datetime_field_formats_validate_before_mutation():
    prs = _open(MINIMAL)
    paragraph = _footer_box(prs).text_frame.paragraphs[0]
    before = save_to_bytes(prs)
    with pytest.raises(ValueError, match="format_code"):
        paragraph.add_datetime_field("datetime99")
    with pytest.raises(ValueError, match="format_code"):
        paragraph.add_datetime_field("slidenum")
    assert_changed_parts(before, save_to_bytes(prs))  # -- empty budget

    paragraph.add_datetime_field("datetime3")
    fld = paragraph._p.find("{%s}fld" % _A)
    assert fld.get("type") == "datetime3"


def test_field_lands_before_end_paragraph_properties():
    prs = _open(MINIMAL)
    paragraph = _footer_box(prs).text_frame.paragraphs[0]
    endParaRPr = paragraph._p.makeelement("{%s}endParaRPr" % _A, {})
    paragraph._p.append(endParaRPr)
    paragraph.add_slide_number_field()
    children = [etree.QName(child).localname for child in paragraph._p]
    assert children.index("fld") < children.index("endParaRPr")


def test_field_paragraph_passes_the_fragment_oracle():
    prs = _open(MINIMAL)
    paragraph = _footer_box(prs).text_frame.paragraphs[0]
    paragraph.add_run().text = "Page "
    paragraph.add_slide_number_field()
    errors = _validation_errors(paragraph._p, "p", "CT_TextParagraph")
    assert not errors, "\n".join(errors)


def test_inspect_reports_fields_without_hashing_their_volatile_text():
    prs = _open(MINIMAL)
    paragraph = _footer_box(prs).text_frame.paragraphs[0]
    paragraph.add_run().text = "Page "
    text_only_hash = next(
        b.anchor.content_hash
        for b in inspect_text(prs.slides[0]).blocks
        if b.shape_name == "footer_box"
    )
    paragraph.add_slide_number_field()

    block = next(
        b for b in inspect_text(prs.slides[0]).blocks if b.shape_name == "footer_box"
    )
    assert block.fields == ("slidenum",)
    assert block.text == "Page "  # -- field display text excluded from text...
    assert block.anchor.content_hash == text_only_hash  # -- ...so the anchor is hash-stable
    assert block.to_dict()["fields"] == ["slidenum"]


def test_normalize_autofit_still_covers_created_fields():
    """The organ interplay: a created field participates in autofit freezing."""
    from pptx.enum.text import MSO_AUTO_SIZE
    from pptx.util import Pt

    prs = _open(MINIMAL)
    box = _footer_box(prs)
    paragraph = box.text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = "Page "
    run.font.size = Pt(10)  # -- run is sized; only the field lacks a size
    paragraph.add_slide_number_field()
    box.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    normAutofit = box.text_frame._txBody.find("{%s}bodyPr" % _A).find(
        "{%s}normAutofit" % _A
    )
    normAutofit.set("fontScale", "50000")

    from pptx.errors import UnsupportedStructureError

    with pytest.raises(UnsupportedStructureError, match="field"):
        box.text_frame.normalize_autofit()  # -- unsized field refuses, as pinned in 0.4


# --------------------------------------------------------------------------------- p:hf flags


def test_layout_and_master_header_footer_flags_round_trip():
    prs = _open(MINIMAL)
    before = save_to_bytes(prs)
    layout = prs.slides[0].slide_layout
    layout.header_footers.slide_number_visible = False
    layout.header_footers.footer_visible = False
    prs.slide_masters[0].header_footers.date_visible = False
    after = save_to_bytes(prs)
    assert_changed_parts(
        before,
        after,
        expect_changed=[
            "ppt/slideLayouts/slideLayout1.xml",
            "ppt/slideMasters/slideMaster1.xml",
        ],
    )

    reopened = _reopen(after)
    reopened_hf = reopened.slides[0].slide_layout.header_footers
    assert reopened_hf.slide_number_visible is False
    assert reopened_hf.footer_visible is False
    assert reopened_hf.date_visible is None  # -- untouched: inherit
    assert reopened.slide_masters[0].header_footers.date_visible is False


def test_assigning_none_removes_the_flag():
    prs = _open(MINIMAL)
    hf = prs.slides[0].slide_layout.header_footers
    hf.slide_number_visible = False
    assert hf.slide_number_visible is False
    hf.slide_number_visible = None
    assert hf.slide_number_visible is None


def test_flag_rejects_non_bool():
    prs = _open(MINIMAL)
    with pytest.raises(ValueError):
        prs.slides[0].slide_layout.header_footers.footer_visible = "yes"


# ------------------------------------------------------------------------------- lo_smoke


@pytest.mark.lo_smoke
def test_field_output_loads_in_libreoffice(tmp_path):
    prs = _open(MINIMAL)
    paragraph = _footer_box(prs).text_frame.paragraphs[0]
    paragraph.add_run().text = "Page "
    paragraph.add_slide_number_field()
    out = tmp_path / "fields.pptx"
    prs.save(str(out))
    lo_load_smoke(out, tmp_path)
