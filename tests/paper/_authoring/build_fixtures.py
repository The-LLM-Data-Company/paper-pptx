"""One-time authoring script for the bootstrap fixture corpus in `tests/paper/fixtures/`.

This script is development tooling, not package code and not test code. It exists so every
self-generated fixture's construction is reviewable and reproducible in *content* (regenerated
files are semantically identical but not byte-identical: zip entry timestamps differ). The
fixtures themselves are FROZEN: tests verify them against `MANIFEST.sha256` and never call this
script. Regenerating a fixture requires updating the manifest and sidecar in a reviewed PR.

Provenance produced here is honestly "self-generated" (python-pptx, this file). The
LibreOffice-export bucket is produced separately by round-tripping these files through headless
LibreOffice (see `tests/paper/fixtures/README.md`); real-PowerPoint and Google-export fixtures
require a human and are specified in `FIXTURE-REQUESTS.md`.

Usage:  python tests/paper/_authoring/build_fixtures.py
"""

from __future__ import annotations

import io
import zipfile
from pathlib import Path

from lxml import etree
from PIL import Image

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.oxml.ns import qn
from pptx.util import Inches

OUT_DIR = Path(__file__).resolve().parent.parent / "fixtures" / "self_generated"

# -- Distinctive, non-default style values so inheritance in ground truth is unambiguous.
# -- Default-template values are asserted before each edit so template drift fails loudly.
BRAND_TITLE_MASTER_SZ = "4000"  # default 4400
BRAND_TITLE_LAYOUT_SZ = "3600"  # layout-level override, no default counterpart
BRAND_BODY_L1_SZ = "2600"  # default 3200
BRAND_BODY_L2_SZ = "2200"  # default 2800
BRAND_BODY_L1_FONT = "Trebuchet MS"  # default "+mn-lt"


def _expect(condition: bool, message: str) -> None:
    """Fail loudly if a precondition about the default template no longer holds."""
    if not condition:
        raise AssertionError("fixture-authoring precondition failed: %s" % message)


def _save(prs, filename: str) -> Path:
    path = OUT_DIR / filename
    prs.save(str(path))
    return path


def _png_bytes(palette: int = 0) -> bytes:
    """Return a small deterministic 64x64 quadrant-pattern PNG.

    `palette` selects a color rotation so callers can produce distinct-but-deterministic
    image bytes (default 0 preserves the original v0 fixture palette exactly).
    """
    img = Image.new("RGB", (64, 64))
    colors = [(200, 30, 30), (30, 160, 60), (30, 60, 200), (220, 180, 40)]
    colors = colors[palette % 4 :] + colors[: palette % 4]
    quadrant_colors = {(0, 0): colors[0], (1, 0): colors[1], (0, 1): colors[2], (1, 1): colors[3]}
    for x in range(64):
        for y in range(64):
            img.putpixel((x, y), quadrant_colors[(x // 32, y // 32)])
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


def _brand_master(prs) -> None:
    """Give the master's txStyles distinctive sizes/fonts (see module constants)."""
    txStyles = prs.slide_masters[0].element.find(qn("p:txStyles"))
    _expect(txStyles is not None, "master has no p:txStyles")

    title_defRPr = txStyles.find(qn("p:titleStyle")).find(qn("a:lvl1pPr")).find(qn("a:defRPr"))
    _expect(title_defRPr.get("sz") == "4400", "default master titleStyle lvl1 sz is not 4400")
    title_defRPr.set("sz", BRAND_TITLE_MASTER_SZ)
    title_latin = title_defRPr.find(qn("a:latin"))
    _expect(
        title_latin is not None and title_latin.get("typeface") == "+mj-lt",
        "default master titleStyle font is not the +mj-lt theme reference",
    )  # -- left as "+mj-lt" deliberately: title font must resolve through the theme

    bodyStyle = txStyles.find(qn("p:bodyStyle"))
    l1_defRPr = bodyStyle.find(qn("a:lvl1pPr")).find(qn("a:defRPr"))
    _expect(l1_defRPr.get("sz") == "3200", "default master bodyStyle lvl1 sz is not 3200")
    l1_defRPr.set("sz", BRAND_BODY_L1_SZ)
    l1_latin = l1_defRPr.find(qn("a:latin"))
    _expect(l1_latin is not None, "master bodyStyle lvl1 has no a:latin")
    l1_latin.set("typeface", BRAND_BODY_L1_FONT)

    l2_defRPr = bodyStyle.find(qn("a:lvl2pPr")).find(qn("a:defRPr"))
    _expect(l2_defRPr.get("sz") == "2800", "default master bodyStyle lvl2 sz is not 2800")
    l2_defRPr.set("sz", BRAND_BODY_L2_SZ)
    l2_latin = l2_defRPr.find(qn("a:latin"))
    _expect(
        l2_latin is not None and l2_latin.get("typeface") == "+mn-lt",
        "default master bodyStyle lvl2 font is not the +mn-lt theme reference",
    )  # -- lvl2 font left as "+mn-lt" so ground truth exercises master->theme resolution


def _override_layout_title_size(layout, sz: str) -> None:
    """Add an `a:lvl1pPr/a:defRPr@sz` override to the layout title placeholder's lstStyle."""
    title_ph = next(p for p in layout.placeholders if p.placeholder_format.idx == 0)
    lstStyle = title_ph.text_frame._txBody.find(qn("a:lstStyle"))
    _expect(lstStyle is not None, "layout title placeholder has no a:lstStyle")
    _expect(len(lstStyle) == 0, "layout title placeholder lstStyle is not empty")
    lvl1pPr = etree.SubElement(lstStyle, qn("a:lvl1pPr"))
    defRPr = etree.SubElement(lvl1pPr, qn("a:defRPr"))
    defRPr.set("sz", sz)


def _remap_clrmap(prs) -> None:
    """Swap accent1<->accent2 and invert bg1/tx1 in the master's p:clrMap."""
    clrMap = prs.slide_masters[0].element.find(qn("p:clrMap"))
    defaults = {"bg1": "lt1", "tx1": "dk1", "accent1": "accent1", "accent2": "accent2"}
    for attr, expected in defaults.items():
        _expect(clrMap.get(attr) == expected, "default clrMap %s is not %s" % (attr, expected))
    clrMap.set("bg1", "dk1")
    clrMap.set("tx1", "lt1")
    clrMap.set("accent1", "accent2")
    clrMap.set("accent2", "accent1")


def _add_named_textbox(slide, name, left, top, width, height, text):
    box = slide.shapes.add_textbox(left, top, width, height)
    box.name = name
    tf = box.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].add_run().text = text
    return box


def _add_real_bullet(paragraph, char="•", font="Arial") -> None:
    """Write a real a:buFont/a:buChar pair into the paragraph's a:pPr.

    Raw-XML on purpose: bullets are exactly the gap the fork fills in Phase 2; the fixture must
    exist before the API does. a:buFont precedes a:buChar per the a:pPr child sequence.
    """
    pPr = paragraph._p.get_or_add_pPr()
    _expect(len(pPr) == 0, "paragraph pPr is not empty; bullet insert order not handled here")
    buFont = etree.SubElement(pPr, qn("a:buFont"))
    buFont.set("typeface", font)
    buChar = etree.SubElement(pPr, qn("a:buChar"))
    buChar.set("char", char)


def _set_norm_autofit_details(text_frame, font_scale="62500", ln_spc_reduction="20000") -> None:
    """Set fontScale/lnSpcReduction on an existing a:normAutofit (raw: oxml lacks the attrs)."""
    normAutofit = text_frame._txBody.find(qn("a:bodyPr")).find(qn("a:normAutofit"))
    _expect(normAutofit is not None, "text frame has no a:normAutofit element")
    normAutofit.set("fontScale", font_scale)
    normAutofit.set("lnSpcReduction", ln_spc_reduction)


# ------------------------------------------------------------------------- fixtures


def build_minimal_clean() -> Path:
    """One title slide from the default template. The minimal-clean corpus anchor."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text_frame.paragraphs[0].add_run().text = "Minimal clean fixture"
    slide.placeholders[1].text_frame.paragraphs[0].add_run().text = "One slide, no extras."
    return _save(prs, "minimal_clean.pptx")


def build_branded_template() -> Path:
    """Feature: placeholder text inheriting size/font through layout/master/theme.

    Title run: no local size/font; layout lstStyle sz=3600 overrides master titleStyle sz=4000;
    font resolves master "+mj-lt" -> theme majorFont (Calibri).
    Body run lvl0: master bodyStyle lvl1 sz=2600, explicit "Trebuchet MS".
    Body run lvl1: master bodyStyle lvl2 sz=2200; font "+mn-lt" -> theme minorFont (Calibri).
    """
    prs = Presentation()
    _brand_master(prs)
    layout = prs.slide_layouts[1]  # -- "Title and Content"
    _override_layout_title_size(layout, BRAND_TITLE_LAYOUT_SZ)

    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text_frame.paragraphs[0].add_run().text = "Branded Title"
    body_tf = slide.placeholders[1].text_frame
    body_tf.paragraphs[0].add_run().text = "Body level one"
    level_two = body_tf.add_paragraph()
    level_two.level = 1
    level_two.add_run().text = "Body level two"

    for paragraph in [slide.shapes.title.text_frame.paragraphs[0]] + list(body_tf.paragraphs):
        for run in paragraph.runs:
            _expect(run.font.size is None, "fixture run must not carry a local font size")
            _expect(run.font.name is None, "fixture run must not carry a local font name")
    return _save(prs, "branded_template.pptx")


def build_clrmap_remap() -> Path:
    """Feature: master remaps theme colors via p:clrMap.

    accent1<->accent2 swapped and bg1/tx1 inverted. A rectangle filled with scheme color
    "accent1" must resolve to the theme's accent2 RGB (C0504D); text colored "tx1" must resolve
    to theme lt1 (FFFFFF).
    """
    prs = Presentation()
    _remap_clrmap(prs)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # -- "Blank"

    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(3), Inches(1.5)
    )
    rect.name = "accent1_box"
    rect.fill.solid()
    rect.fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_1

    box = _add_named_textbox(
        slide, "tx1_text", Inches(1), Inches(3), Inches(4), Inches(1), "Mapped text color"
    )
    run = box.text_frame.paragraphs[0].runs[0]
    run.font.color.theme_color = MSO_THEME_COLOR.TEXT_1

    fill_val = rect.fill.fore_color._color._xClr.get("val")
    _expect(fill_val == "accent1", "rectangle fill schemeClr val is %r" % fill_val)
    text_val = run.font.color._color._xClr.get("val")
    _expect(text_val == "tx1", "run color schemeClr val is %r" % text_val)
    return _save(prs, "clrmap_remap.pptx")


def build_chart_notes() -> Path:
    """The clone fixture: native chart WITH embedded workbook, speaker notes on same slide."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # -- "Title Only"
    slide.shapes.title.text_frame.paragraphs[0].add_run().text = "Chart and notes"

    chart_data = CategoryChartData()
    chart_data.categories = ["East", "West", "Midwest"]
    chart_data.add_series("Q1", (19.2, 21.4, 16.7))
    chart_data.add_series("Q2", (22.3, 28.6, 15.2))
    frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1.5), Inches(8), Inches(5), chart_data
    )
    frame.name = "clone_fixture_chart"

    slide.notes_slide.notes_text_frame.text = "Speaker notes for the clone fixture."
    return _save(prs, "chart_notes.pptx")


def build_shared_media() -> Path:
    """Two slides that deliberately share one image part (python-pptx dedups identical bytes)."""
    prs = Presentation()
    png = _png_bytes()
    for index in range(2):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        picture = slide.shapes.add_picture(io.BytesIO(png), Inches(1), Inches(1), Inches(2))
        picture.name = "shared_image_%d" % (index + 1)
    image_parts = set()
    for slide in prs.slides:
        rId = slide.shapes[0]._element.blip_rId
        _expect(rId is not None, "picture has no blip rId")
        image_parts.add(slide.part.related_part(rId).partname)
    _expect(len(image_parts) == 1, "slides do not share a single image part: %r" % image_parts)
    return _save(prs, "shared_media.pptx")


def _autofit_deck(mode_setter) -> "Presentation":
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = _add_named_textbox(
        slide,
        "autofit_box",
        Inches(1),
        Inches(1),
        Inches(4),
        Inches(2),
        "Autofit fixture text that is long enough to exercise fitting behavior.",
    )
    mode_setter(box.text_frame)
    return prs


def build_autofit_none() -> Path:
    """Feature: explicit a:noAutofit."""

    def setter(tf):
        tf.auto_size = MSO_AUTO_SIZE.NONE

    return _save(_autofit_deck(setter), "autofit_none.pptx")


def build_autofit_normal() -> Path:
    """Feature: a:normAutofit with explicit fontScale=62500 and lnSpcReduction=20000."""

    def setter(tf):
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        _set_norm_autofit_details(tf)

    return _save(_autofit_deck(setter), "autofit_normal.pptx")


def build_autofit_shape() -> Path:
    """Feature: a:spAutoFit (shape grows to fit text)."""

    def setter(tf):
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

    return _save(_autofit_deck(setter), "autofit_shape.pptx")


def _whitespace_deck(text) -> "Presentation":
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_named_textbox(slide, "whitespace_box", Inches(1), Inches(1), Inches(4), Inches(1), text)
    return prs


def build_whitespace_pair() -> "tuple[Path, Path]":
    """Two decks identical except for one trailing space inside an a:t text node.

    The Phase 5 kernel trap: any comparison that trims meaningful whitespace will call these
    equivalent. They are not.
    """
    path_a = _save(_whitespace_deck("Trailing space "), "whitespace_trailing_a.pptx")
    path_b = _save(_whitespace_deck("Trailing space"), "whitespace_trailing_b.pptx")
    return path_a, path_b


def build_gauntlet() -> Path:
    """Everything ugly combined: branding + clrMap remap + chart/notes + shared media +
    all three autofit modes + table + real bullet + cropped picture + hyperlink + empty
    placeholder."""
    prs = Presentation()
    _brand_master(prs)
    _remap_clrmap(prs)
    _override_layout_title_size(prs.slide_layouts[1], BRAND_TITLE_LAYOUT_SZ)
    png = _png_bytes()

    # -- slide 1: branded placeholders ------------------------------------------------
    s1 = prs.slides.add_slide(prs.slide_layouts[1])
    s1.shapes.title.text_frame.paragraphs[0].add_run().text = "Gauntlet: branded"
    body_tf = s1.placeholders[1].text_frame
    body_tf.paragraphs[0].add_run().text = "Inherited body level one"
    p2 = body_tf.add_paragraph()
    p2.level = 1
    p2.add_run().text = "Inherited body level two"

    # -- slide 2: chart + embedded workbook + notes + shared image --------------------
    s2 = prs.slides.add_slide(prs.slide_layouts[5])
    s2.shapes.title.text_frame.paragraphs[0].add_run().text = "Gauntlet: chart and notes"
    chart_data = CategoryChartData()
    chart_data.categories = ["Alpha", "Beta"]
    chart_data.add_series("S1", (3.5, 4.25))
    frame = s2.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1.5), Inches(5), Inches(4), chart_data
    )
    frame.name = "gauntlet_chart"
    s2.notes_slide.notes_text_frame.text = "Gauntlet speaker notes."
    s2.shapes.add_picture(io.BytesIO(png), Inches(7), Inches(2), Inches(2)).name = "gauntlet_img_1"

    # -- slide 3: autofit trio + table + real bullet + cropped picture ----------------
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    none_box = _add_named_textbox(
        s3, "autofit_none_box", Inches(0.5), Inches(0.5), Inches(3), Inches(1.2), "No autofit."
    )
    none_box.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    norm_box = _add_named_textbox(
        s3, "autofit_normal_box", Inches(0.5), Inches(2), Inches(3), Inches(1.2), "Norm autofit."
    )
    norm_box.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    _set_norm_autofit_details(norm_box.text_frame)
    shape_box = _add_named_textbox(
        s3, "autofit_shape_box", Inches(0.5), Inches(3.5), Inches(3), Inches(1.2), "Sp autofit."
    )
    shape_box.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

    bullet_box = _add_named_textbox(
        s3, "real_bullet_box", Inches(4), Inches(0.5), Inches(3), Inches(1.5), "A real bullet"
    )
    # -- explicit autofit so only the three autofit_* boxes carry their signature elements
    # -- (add_textbox's default bodyPr contains a:spAutoFit, which would muddy ground truth)
    bullet_box.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    _add_real_bullet(bullet_box.text_frame.paragraphs[0])

    table_shape = s3.shapes.add_table(3, 3, Inches(4), Inches(2.5), Inches(4), Inches(2))
    table_shape.name = "gauntlet_table"
    for row in range(3):
        for col in range(3):
            cell_tf = table_shape.table.cell(row, col).text_frame
            cell_tf.paragraphs[0].add_run().text = "r%dc%d" % (row, col)

    cropped = s3.shapes.add_picture(io.BytesIO(png), Inches(8.5), Inches(0.5), Inches(1.5))
    cropped.name = "gauntlet_cropped"
    cropped.crop_left = 0.25
    cropped.crop_top = 0.1

    # -- slide 4: shared media again + hyperlink + empty body placeholder -------------
    s4 = prs.slides.add_slide(prs.slide_layouts[1])
    s4.shapes.title.text_frame.paragraphs[0].add_run().text = "Gauntlet: end"
    # -- body placeholder (idx 1) left deliberately empty --
    s4.shapes.add_picture(io.BytesIO(png), Inches(6), Inches(4), Inches(2)).name = "gauntlet_img_2"
    link_box = _add_named_textbox(
        s4, "hyperlink_box", Inches(1), Inches(4), Inches(4), Inches(0.8), "External link"
    )
    link_box.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    link_box.text_frame.paragraphs[0].runs[0].hyperlink.address = "https://example.com/paper"

    return _save(prs, "gauntlet.pptx")


def build_corrupt_dangling_sldid() -> Path:
    """Corrupt-by-construction: p:sldIdLst entry referencing a relationship that does not exist.

    Opens under python-pptx (parts load lazily) but any traversal touching the second sldId
    raises KeyError. Negative tests only.
    """
    buf = io.BytesIO()
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text_frame.paragraphs[0].add_run().text = "Corrupt fixture"
    prs.save(buf)

    source = zipfile.ZipFile(buf)
    presentation_xml = etree.fromstring(source.read("ppt/presentation.xml"))
    sldIdLst = presentation_xml.find(qn("p:sldIdLst"))
    _expect(sldIdLst is not None and len(sldIdLst) == 1, "expected exactly one p:sldId")
    dangling = etree.SubElement(sldIdLst, qn("p:sldId"))
    dangling.set("id", "9999")
    dangling.set(qn("r:id"), "rId99")

    path = OUT_DIR / "corrupt_dangling_sldid.pptx"
    with zipfile.ZipFile(str(path), "w", zipfile.ZIP_DEFLATED) as out:
        for info in source.infolist():
            if info.filename == "ppt/presentation.xml":
                out.writestr(
                    info,
                    etree.tostring(
                        presentation_xml, xml_declaration=True, encoding="UTF-8", standalone=True
                    ),
                )
            else:
                out.writestr(info, source.read(info.filename))
    return path


def build_sections() -> Path:
    """Five-slide deck with a p14:sectionLst (3 sections) and a p:custShowLst (1 custom show).

    Sections/custom shows cannot be authored via python-pptx, so presentation.xml is rewritten
    by zip surgery after a normal build — same technique as the corrupt fixture, honestly
    self-generated. Section GUIDs are fixed constants for determinism. A real-PowerPoint
    version is requested as FIXTURE-REQUESTS.md R8 (LibreOffice does not faithfully author
    the p14 extension list).
    """
    P14 = "http://schemas.microsoft.com/office/powerpoint/2010/main"
    R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    SECTION_EXT_URI = "{521415D9-36F7-43E2-AB2F-B90AF26B5E84}"

    buf = io.BytesIO()
    prs = Presentation()
    for index in range(5):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text_frame.paragraphs[0].add_run().text = "Section deck slide %d" % (
            index + 1
        )
    prs.save(buf)

    source = zipfile.ZipFile(buf)
    presentation_xml = etree.fromstring(source.read("ppt/presentation.xml"))
    sldIdLst = presentation_xml.find(qn("p:sldIdLst"))
    slide_entries = [(s.get("id"), s.get("{%s}id" % R)) for s in sldIdLst]
    _expect(len(slide_entries) == 5, "expected five sldId entries")

    # -- p:custShowLst goes after p:notesSz per the CT_Presentation child sequence
    custShowLst = etree.SubElement(presentation_xml, qn("p:custShowLst"))
    notesSz = presentation_xml.find(qn("p:notesSz"))
    notesSz.addnext(custShowLst)
    custShow = etree.SubElement(custShowLst, qn("p:custShow"))
    custShow.set("name", "Focus")
    custShow.set("id", "0")
    sldLst = etree.SubElement(custShow, qn("p:sldLst"))
    for slide_index in (1, 3):  # -- slides 2 and 4, by relationship id
        sld = etree.SubElement(sldLst, qn("p:sld"))
        sld.set("{%s}id" % R, slide_entries[slide_index][1])

    # -- p:extLst last; p14:sectionLst inside the Microsoft section extension
    extLst = etree.SubElement(presentation_xml, qn("p:extLst"))
    ext = etree.SubElement(extLst, qn("p:ext"))
    ext.set("uri", SECTION_EXT_URI)
    sectionLst = etree.SubElement(ext, "{%s}sectionLst" % P14)
    section_plan = [
        ("Intro", "{11111111-1111-4111-8111-111111111111}", [0]),
        ("Body", "{22222222-2222-4222-8222-222222222222}", [1, 2, 3]),
        ("Close", "{33333333-3333-4333-8333-333333333333}", [4]),
    ]
    for name, guid, slide_indices in section_plan:
        section = etree.SubElement(sectionLst, "{%s}section" % P14)
        section.set("name", name)
        section.set("id", guid)
        p14_sldIdLst = etree.SubElement(section, "{%s}sldIdLst" % P14)
        for slide_index in slide_indices:
            p14_sldId = etree.SubElement(p14_sldIdLst, "{%s}sldId" % P14)
            p14_sldId.set("id", slide_entries[slide_index][0])

    path = OUT_DIR / "sections.pptx"
    with zipfile.ZipFile(str(path), "w", zipfile.ZIP_DEFLATED) as out:
        for info in source.infolist():
            if info.filename == "ppt/presentation.xml":
                out.writestr(
                    info,
                    etree.tostring(
                        presentation_xml, xml_declaration=True, encoding="UTF-8", standalone=True
                    ),
                )
            else:
                out.writestr(info, source.read(info.filename))
    return path


def build_tables_in_group() -> Path:
    """A table (a:graphicFrame) inside a p:grpSp, plus an in-group textbox and a top-level one.

    PowerPoint's UI does not group tables but the schema allows it and third-party producers
    emit it; the graphicFrame is moved into the group element directly (python-pptx has no
    group-level add_table).
    """
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_named_textbox(
        slide, "top_level_box", Inches(0.5), Inches(0.5), Inches(3), Inches(0.8), "Top-level text"
    )
    group = slide.shapes.add_group_shape()
    group.name = "outer_group"
    in_group_box = group.shapes.add_textbox(Inches(1), Inches(2), Inches(3), Inches(0.8))
    in_group_box.name = "in_group_box"
    in_group_box.text_frame.paragraphs[0].add_run().text = "In-group text"

    table_shape = slide.shapes.add_table(2, 2, Inches(1), Inches(3.5), Inches(4), Inches(1.5))
    table_shape.name = "grouped_table"
    for row in range(2):
        for col in range(2):
            cell_tf = table_shape.table.cell(row, col).text_frame
            cell_tf.paragraphs[0].add_run().text = "cell r%dc%d" % (row, col)
    # -- move the table's graphicFrame element into the group element
    group._element.append(table_shape._element)
    return _save(prs, "tables_in_group.pptx")


def build_nested_groups() -> Path:
    """Groups nested three deep, a named textbox at every level (recursion + depth fixture)."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_named_textbox(
        slide, "level0_box", Inches(0.5), Inches(0.5), Inches(3), Inches(0.8), "Level zero"
    )
    outer = slide.shapes.add_group_shape()
    outer.name = "group_level1"
    box1 = outer.shapes.add_textbox(Inches(1), Inches(1.5), Inches(3), Inches(0.6))
    box1.name = "level1_box"
    box1.text_frame.paragraphs[0].add_run().text = "Level one"
    middle = outer.shapes.add_group_shape()
    middle.name = "group_level2"
    box2 = middle.shapes.add_textbox(Inches(1.2), Inches(2.5), Inches(3), Inches(0.6))
    box2.name = "level2_box"
    box2.text_frame.paragraphs[0].add_run().text = "Level two"
    inner = middle.shapes.add_group_shape()
    inner.name = "group_level3"
    box3 = inner.shapes.add_textbox(Inches(1.4), Inches(3.5), Inches(3), Inches(0.6))
    box3.name = "level3_box"
    box3.text_frame.paragraphs[0].add_run().text = "Level three"
    return _save(prs, "nested_groups.pptx")


def build_autofit_inherited() -> Path:
    """normAutofit placeholder whose font sizes resolve ONLY through the master's txStyles.

    The PLAN-v0.1 Phase 0.4 fixture: normalize_autofit() must refuse without resolution and
    succeed with `resolve=True` (sizes come from the branded master: lvl1 2600, lvl2 2200).
    fontScale only — no lnSpcReduction — because spacing resolution stays a refusal in v0.1.
    """
    prs = Presentation()
    _brand_master(prs)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text_frame.paragraphs[0].add_run().text = "Inherited autofit"
    body_tf = slide.placeholders[1].text_frame
    body_tf.paragraphs[0].add_run().text = "Level one inherits 26pt"
    level_two = body_tf.add_paragraph()
    level_two.level = 1
    level_two.add_run().text = "Level two inherits 22pt"
    for paragraph in body_tf.paragraphs:
        for run in paragraph.runs:
            _expect(run.font.size is None, "runs must not carry local sizes")

    body_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    normAutofit = body_tf._txBody.find(qn("a:bodyPr")).find(qn("a:normAutofit"))
    normAutofit.set("fontScale", "62500")
    return _save(prs, "autofit_inherited.pptx")


def build_hf_flags() -> Path:
    """Master and first layout carrying explicit p:hf visibility flags.

    Injected by zip surgery (v0 python-pptx could not author p:hf), so the corpus
    round-trip and the HeaderFooters read path exercise authored-elsewhere flags — not just
    flags this package wrote itself. Real-PowerPoint version tracked in FIXTURE-REQUESTS.md.
    """
    buf = io.BytesIO()
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text_frame.paragraphs[0].add_run().text = "Header/footer flags"
    prs.save(buf)

    source = zipfile.ZipFile(buf)

    def with_hf(member_name, attrs, successors_first):
        root = etree.fromstring(source.read(member_name))
        hf = root.makeelement(qn("p:hf"), attrs)
        anchor = None
        for tag in successors_first:
            anchor = root.find(qn(tag))
            if anchor is not None:
                break
        if anchor is not None:
            anchor.addprevious(hf)
        else:
            root.append(hf)
        return etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)

    rewrites = {
        "ppt/slideMasters/slideMaster1.xml": with_hf(
            "ppt/slideMasters/slideMaster1.xml",
            {"sldNum": "0", "ftr": "0", "dt": "1"},
            ("p:txStyles", "p:extLst"),
        ),
        "ppt/slideLayouts/slideLayout1.xml": with_hf(
            "ppt/slideLayouts/slideLayout1.xml", {"sldNum": "1"}, ("p:extLst",)
        ),
    }
    path = OUT_DIR / "hf_flags.pptx"
    with zipfile.ZipFile(str(path), "w", zipfile.ZIP_DEFLATED) as out:
        for info in source.infolist():
            out.writestr(info, rewrites.get(info.filename, source.read(info.filename)))
    return path


def build_large_smoke() -> Path:
    """Large deck for perf smoke: 120 text slides, a picture every 10th slide."""
    prs = Presentation()
    png = _png_bytes()
    for index in range(120):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text_frame.paragraphs[0].add_run().text = "Slide %d" % (index + 1)
        body_tf = slide.placeholders[1].text_frame
        body_tf.paragraphs[0].add_run().text = "First line of body content on slide %d." % (
            index + 1
        )
        body_tf.add_paragraph().add_run().text = "Second line of body content."
        if index % 10 == 0:
            slide.shapes.add_picture(io.BytesIO(png), Inches(7), Inches(4), Inches(2))
    return _save(prs, "large_smoke.pptx")


# ---------------------------------------------------------------- v0.11 fixtures
#
# Added 2026-07-08 for PLAN-v0.11 Phase 0. Same rules as above: honest self-generated
# provenance, zip surgery only where python-pptx cannot author the structure, fixed
# GUIDs/ids for determinism, preconditions asserted so template drift fails loudly.
# Real-PowerPoint equivalents are FIXTURE-REQUESTS.md R9-R14.

_RELS_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
_CTYPES_NS = "http://schemas.openxmlformats.org/package/2006/content-types"


def _xml_bytes(root) -> bytes:
    return etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)


def _rezip(source: zipfile.ZipFile, path: Path, rewrites: dict, additions: dict) -> Path:
    """Write `source`'s members to `path`, replacing per `rewrites`, appending `additions`."""
    with zipfile.ZipFile(str(path), "w", zipfile.ZIP_DEFLATED) as out:
        for info in source.infolist():
            data = rewrites.get(info.filename)
            out.writestr(info, source.read(info.filename) if data is None else data)
        for name, data in additions.items():
            out.writestr(name, data)
    return path


def _with_rel(rels_xml: bytes, rel_type: str, target: str) -> "tuple[bytes, str]":
    """Return (rels_xml with one appended Relationship at the next free rId, that rId)."""
    root = etree.fromstring(rels_xml)
    used = [int(r.get("Id")[3:]) for r in root if r.get("Id", "").startswith("rId")]
    rId = "rId%d" % (max(used, default=0) + 1)
    rel = etree.SubElement(root, "{%s}Relationship" % _RELS_NS)
    rel.set("Id", rId)
    rel.set("Type", rel_type)
    rel.set("Target", target)
    return _xml_bytes(root), rId


def _with_content_types(ct_xml: bytes, defaults: dict, overrides: dict) -> bytes:
    root = etree.fromstring(ct_xml)
    have = {el.get("Extension") for el in root if el.tag.endswith("Default")}
    for extension, content_type in defaults.items():
        if extension in have:
            continue
        el = etree.SubElement(root, "{%s}Default" % _CTYPES_NS)
        el.set("Extension", extension)
        el.set("ContentType", content_type)
    for partname, content_type in overrides.items():
        el = etree.SubElement(root, "{%s}Override" % _CTYPES_NS)
        el.set("PartName", partname)
        el.set("ContentType", content_type)
    return _xml_bytes(root)


def build_merged_tables() -> Path:
    """5x4 table: header row merged across all four columns, a 2-row vertical merge in
    column one, everything else regular with distinct text. Written with upstream's own
    `_Cell.merge` so the gridSpan/rowSpan/hMerge/vMerge attributes are exactly what
    python-pptx produces; the real-PowerPoint counterpart is FIXTURE-REQUESTS.md R11."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_table(5, 4, Inches(0.5), Inches(0.7), Inches(9), Inches(5))
    shape.name = "merged_table"
    table = shape.table
    table.cell(0, 0).merge(table.cell(0, 3))  # -- header row: gridSpan=4 + 3 hMerge cells
    table.cell(2, 0).merge(table.cell(3, 0))  # -- column one rows 3-4: rowSpan=2 + vMerge
    table.cell(0, 0).text_frame.paragraphs[0].add_run().text = "Merged header"
    table.cell(2, 0).text_frame.paragraphs[0].add_run().text = "Merged rows"
    for row in range(1, 5):
        for col in range(4):
            cell = table.cell(row, col)
            if cell.is_merge_origin or cell.is_spanned:
                continue
            cell.text_frame.paragraphs[0].add_run().text = "r%dc%d" % (row, col)

    _expect(table.cell(0, 0)._tc.gridSpan == 4, "header origin gridSpan is not 4")
    _expect(table.cell(0, 1)._tc.hMerge, "header continuation cell is not hMerge")
    _expect(table.cell(2, 0)._tc.rowSpan == 2, "column-one origin rowSpan is not 2")
    _expect(table.cell(3, 0)._tc.vMerge, "column-one continuation cell is not vMerge")
    _expect(
        all(len(tr.findall(qn("a:tc"))) == 4 for tr in table._tbl.findall(qn("a:tr"))),
        "every row must hold exactly one a:tc per gridCol (continuations included)",
    )
    return _save(prs, "merged_tables.pptx")


_HF_LAYOUT_IDX = {"dt": 10, "ftr": 11, "sldNum": 12}  # -- default-template layout furniture


def _next_shape_id(slide) -> int:
    ids = [
        int(el.get("id"))
        for el in slide.shapes._spTree.iter(qn("p:cNvPr"))
        if el.get("id", "").isdigit()
    ]
    return max(ids, default=0) + 1


def _materialize_dialog_placeholder(slide, ph_type: str, name: str, body_xml: str) -> None:
    """Append the minimal placeholder `p:sp` PowerPoint's Header & Footer dialog persists.

    Raw-XML on purpose (fixture-first: this is the structure the v0.11 Phase 2 API must
    produce, so the fixture cannot be authored with that API). Geometry and formatting are
    deliberately absent - they inherit from the layout's matching-idx placeholder.
    """
    from pptx.oxml import parse_xml
    from pptx.oxml.ns import nsdecls

    sp_xml = (
        "<p:sp %s>"
        "<p:nvSpPr>"
        '<p:cNvPr id="%d" name="%s"/>'
        '<p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>'
        '<p:nvPr><p:ph type="%s" sz="quarter" idx="%d"/></p:nvPr>'
        "</p:nvSpPr>"
        "<p:spPr/>"
        "<p:txBody><a:bodyPr/><a:lstStyle/><a:p>%s</a:p></p:txBody>"
        "</p:sp>"
    ) % (
        nsdecls("p", "a"),
        _next_shape_id(slide),
        name,
        ph_type,
        _HF_LAYOUT_IDX[ph_type],
        body_xml,
    )
    slide.shapes._spTree.append(parse_xml(sp_xml))


def _fld_xml(guid_serial: int, fld_type: str, cached_text: str) -> str:
    """A deterministic-GUID a:fld with cached text, per ISO 29500-1 21.1.2.2.4."""
    guid = "{00000000-1111-4222-8333-%012d}" % guid_serial
    return (
        '<a:fld id="%s" type="%s"><a:rPr lang="en-US" smtClean="0"/><a:t>%s</a:t></a:fld>'
        % (guid, fld_type, cached_text)
    )


def build_footers_applied() -> Path:
    """Dialog-applied footer furniture, reproduced from the mechanism probe (v0.11 Phase 0).

    Five slides. Each carries the three minimal placeholder shapes PowerPoint's
    Insert > Header & Footer dialog materializes on "Apply to All": dt (a:fld
    type="datetime", cached text), ftr (literal run), sldNum (a:fld type="slidenum",
    cached number) - inheriting geometry from the layout's idx-10/11/12 furniture.
    Slide 3 omits the footer placeholder (the per-slide uncheck); slide 5 is hidden
    (`show="0"` on p:sld). p:hf is deliberately absent everywhere: all four attributes
    default true, which is exactly what PowerPoint leaves in this state.
    Real-PowerPoint provenance for the same mechanism is FIXTURE-REQUESTS.md R9.
    """
    prs = Presentation()
    layout = prs.slide_layouts[1]
    layout_idx = {
        ph.get("type"): int(ph.get("idx", "0"))
        for ph in layout._element.iter(qn("p:ph"))
        if ph.get("type") in ("dt", "ftr", "sldNum")
    }
    _expect(
        layout_idx == _HF_LAYOUT_IDX,
        "default-template layout furniture idx map changed: %r" % layout_idx,
    )

    guid_serial = 0
    for n in range(1, 6):
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text_frame.paragraphs[0].add_run().text = "Footered slide %d" % n
        body_tf = slide.placeholders[1].text_frame
        body_tf.paragraphs[0].add_run().text = "Body content on slide %d." % n

        guid_serial += 1
        _materialize_dialog_placeholder(
            slide, "dt", "Date Placeholder %d" % n, _fld_xml(guid_serial, "datetime", "7/8/2026")
        )
        if n != 3:  # -- slide 3: footer unchecked for this slide only
            _materialize_dialog_placeholder(
                slide,
                "ftr",
                "Footer Placeholder %d" % n,
                '<a:r><a:rPr lang="en-US"/><a:t>Paper Fixture Footer</a:t></a:r>',
            )
        guid_serial += 1
        _materialize_dialog_placeholder(
            slide,
            "sldNum",
            "Slide Number Placeholder %d" % n,
            _fld_xml(guid_serial, "slidenum", str(n)),
        )

    prs.slides[4]._element.set("show", "0")  # -- hidden slide
    return _save(prs, "footers_applied.pptx")


def build_scrub_gauntlet() -> Path:
    """Everything scrub must remove, next to everything it must keep (v0.11 Phase 3).

    Live content: 4 slides (slide 4 hidden), notes on slides 1 and 3, a picture on
    slide 1 (media reachable from a live slide - must survive every scrub), core-props
    personal info. Removal targets added by zip surgery, honestly self-generated:
    classic comments (`p:cmLst` part on slide 1 + `p:cmAuthorLst` on the presentation),
    a custom-properties part, an embedded font (`p:embeddedFontLst` + fntdata part), and
    a picture reachable ONLY from unused slideLayout4 (unused-layout media - removed
    exactly when unused layouts are). Real-PowerPoint counterparts: R12/R14.
    """
    buf = io.BytesIO()
    prs = Presentation()
    layouts = prs.slide_layouts
    s1 = prs.slides.add_slide(layouts[1])
    s1.shapes.title.text_frame.paragraphs[0].add_run().text = "Scrub target one"
    s1.placeholders[1].text_frame.paragraphs[0].add_run().text = "Slide one body."
    s1.shapes.add_picture(io.BytesIO(_png_bytes()), Inches(6), Inches(4), Inches(2)).name = (
        "live_media_pic"
    )
    s1.notes_slide.notes_text_frame.text = "Notes on slide one."
    s2 = prs.slides.add_slide(layouts[5])
    s2.shapes.title.text_frame.paragraphs[0].add_run().text = "Scrub target two"
    s3 = prs.slides.add_slide(layouts[1])
    s3.shapes.title.text_frame.paragraphs[0].add_run().text = "Scrub target three"
    s3.placeholders[1].text_frame.paragraphs[0].add_run().text = "Slide three body."
    s3.notes_slide.notes_text_frame.text = "Notes on slide three."
    s4 = prs.slides.add_slide(layouts[6])
    _add_named_textbox(
        s4, "hidden_slide_box", Inches(1), Inches(1), Inches(4), Inches(1), "Hidden slide text"
    )
    s4._element.set("show", "0")

    core = prs.core_properties
    core.author = "Paper Fixture Author"
    core.last_modified_by = "Paper Fixture Editor"
    core.comments = "Deck-level metadata comment."
    prs.save(buf)

    source = zipfile.ZipFile(buf)
    P = "http://schemas.openxmlformats.org/presentationml/2006/main"
    RT_BASE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

    comments_xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n'
        '<p:cmLst xmlns:p="%s">'
        '<p:cm authorId="0" dt="2026-07-08T09:00:00.000" idx="1">'
        '<p:pos x="10" y="10"/><p:text>First fixture comment.</p:text></p:cm>'
        '<p:cm authorId="0" dt="2026-07-08T09:05:00.000" idx="2">'
        '<p:pos x="20" y="20"/><p:text>Second fixture comment.</p:text></p:cm>'
        "</p:cmLst>" % P
    ).encode("utf-8")
    authors_xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n'
        '<p:cmAuthorLst xmlns:p="%s">'
        '<p:cmAuthor id="0" name="Paper Reviewer" initials="PR" lastIdx="2" clrIdx="0"/>'
        "</p:cmAuthorLst>" % P
    ).encode("utf-8")
    custom_props_xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n'
        "<Properties"
        ' xmlns="http://schemas.openxmlformats.org/officeDocument/2006/custom-properties"'
        ' xmlns:vt="http://schemas.openxmlformats.org/officeDocument/2006/docPropsVTypes">'
        '<property fmtid="{D5CDD505-2E9C-101B-9397-08002B2CF9AE}" pid="2" name="Department">'
        "<vt:lpwstr>Fixtures</vt:lpwstr></property>"
        "</Properties>"
    ).encode("utf-8")
    fntdata = b"PAPER-FIXTURE-FNTDATA\x00" * 64  # -- structural stand-in, not a real font

    # -- slide 1 -> comments rel
    slide1_rels, _ = _with_rel(
        source.read("ppt/slides/_rels/slide1.xml.rels"),
        RT_BASE + "/comments",
        "../comments/comment1.xml",
    )
    # -- presentation -> commentAuthors + font rels
    pres_rels, _ = _with_rel(
        source.read("ppt/_rels/presentation.xml.rels"),
        RT_BASE + "/commentAuthors",
        "commentAuthors.xml",
    )
    pres_rels, font_rId = _with_rel(pres_rels, RT_BASE + "/font", "fonts/font1.fntdata")
    # -- package -> custom properties rel
    pkg_rels, _ = _with_rel(
        source.read("_rels/.rels"), RT_BASE + "/custom-properties", "docProps/custom.xml"
    )
    # -- presentation.xml: p:embeddedFontLst between p:notesSz and any p:custShowLst
    pres_root = etree.fromstring(source.read("ppt/presentation.xml"))
    _expect(pres_root.find(qn("p:embeddedFontLst")) is None, "embeddedFontLst already present")
    embeddedFontLst = pres_root.makeelement(qn("p:embeddedFontLst"), {})
    embeddedFont = etree.SubElement(embeddedFontLst, qn("p:embeddedFont"))
    font = etree.SubElement(embeddedFont, qn("p:font"))
    font.set("typeface", "Paper Fixture Font")
    regular = etree.SubElement(embeddedFont, qn("p:regular"))
    regular.set(qn("r:id"), font_rId)
    pres_root.find(qn("p:notesSz")).addnext(embeddedFontLst)

    # -- unused slideLayout4 ("Two Content" in the default template): picture + rel
    layout_rels, layout_img_rId = _with_rel(
        source.read("ppt/slideLayouts/_rels/slideLayout4.xml.rels"),
        RT_BASE + "/image",
        "../media/image_unused_layout.png",
    )
    layout_root = etree.fromstring(source.read("ppt/slideLayouts/slideLayout4.xml"))
    spTree = layout_root.find(qn("p:cSld")).find(qn("p:spTree"))
    _expect(spTree is not None, "layout4 has no spTree")
    A = "http://schemas.openxmlformats.org/drawingml/2006/main"
    pic_xml = (
        '<p:pic xmlns:p="%s" xmlns:a="%s" xmlns:r="%s">'
        '<p:nvPicPr><p:cNvPr id="990" name="unused_layout_pic"/><p:cNvPicPr/><p:nvPr/></p:nvPicPr>'
        '<p:blipFill><a:blip r:embed="%s"/><a:stretch><a:fillRect/></a:stretch></p:blipFill>'
        '<p:spPr><a:xfrm><a:off x="457200" y="457200"/><a:ext cx="914400" cy="914400"/></a:xfrm>'
        '<a:prstGeom prst="rect"><a:avLst/></a:prstGeom></p:spPr>'
        "</p:pic>" % (P, A, RT_BASE, layout_img_rId)
    )
    spTree.append(etree.fromstring(pic_xml))

    content_types = _with_content_types(
        source.read("[Content_Types].xml"),
        defaults={"fntdata": "application/x-fontdata"},
        overrides={
            "/ppt/comments/comment1.xml": (
                "application/vnd.openxmlformats-officedocument.presentationml.comments+xml"
            ),
            "/ppt/commentAuthors.xml": (
                "application/vnd.openxmlformats-officedocument.presentationml.commentAuthors+xml"
            ),
            "/docProps/custom.xml": (
                "application/vnd.openxmlformats-officedocument.custom-properties+xml"
            ),
        },
    )

    return _rezip(
        source,
        OUT_DIR / "scrub_gauntlet.pptx",
        rewrites={
            "[Content_Types].xml": content_types,
            "_rels/.rels": pkg_rels,
            "ppt/_rels/presentation.xml.rels": pres_rels,
            "ppt/presentation.xml": _xml_bytes(pres_root),
            "ppt/slides/_rels/slide1.xml.rels": slide1_rels,
            "ppt/slideLayouts/_rels/slideLayout4.xml.rels": layout_rels,
            "ppt/slideLayouts/slideLayout4.xml": _xml_bytes(layout_root),
        },
        additions={
            "ppt/comments/comment1.xml": comments_xml,
            "ppt/commentAuthors.xml": authors_xml,
            "docProps/custom.xml": custom_props_xml,
            "ppt/fonts/font1.fntdata": fntdata,
            "ppt/media/image_unused_layout.png": _png_bytes(palette=1),
        },
    )


_THEME_FONT_DEFAULTS = ("Calibri", "Calibri")  # -- default-template major/minor latin
_THEME_ACCENT1_DEFAULT = "4F81BD"


def _retheme(source: zipfile.ZipFile, theme_name: str, major: str, minor: str, accent1: str):
    """Return theme1.xml bytes with renamed theme, new major/minor latin fonts, new accent1."""
    A = "http://schemas.openxmlformats.org/drawingml/2006/main"
    root = etree.fromstring(source.read("ppt/theme/theme1.xml"))
    _expect(root.get("name") == "Office Theme", "theme name is not the expected default")
    root.set("name", theme_name)
    fontScheme = root.find("{%s}themeElements/{%s}fontScheme" % (A, A))
    major_latin = fontScheme.find("{%s}majorFont/{%s}latin" % (A, A))
    minor_latin = fontScheme.find("{%s}minorFont/{%s}latin" % (A, A))
    _expect(
        (major_latin.get("typeface"), minor_latin.get("typeface")) == _THEME_FONT_DEFAULTS,
        "theme font defaults changed",
    )
    major_latin.set("typeface", major)
    minor_latin.set("typeface", minor)
    accent1_el = root.find("{%s}themeElements/{%s}clrScheme/{%s}accent1/{%s}srgbClr" % (A, A, A, A))
    _expect(accent1_el.get("val") == _THEME_ACCENT1_DEFAULT, "accent1 default changed")
    accent1_el.set("val", accent1)
    return _xml_bytes(root)


def build_template_alpha() -> Path:
    """Import/rebind corpus, house style A (v0.11 Phases 4-5).

    Default-template deck rethemed by zip surgery ("Paper Alpha": Georgia/Verdana, accent1
    AA3311). Three slides: title, title-and-content with two bullet levels, blank with a
    named picture and textbox. Layout names keep the default-template set, so every name
    collides with template beta's (the import-collision requirement); real-PowerPoint
    counterpart is FIXTURE-REQUESTS.md R10.
    """
    buf = io.BytesIO()
    prs = Presentation()
    s1 = prs.slides.add_slide(prs.slide_layouts[0])
    s1.shapes.title.text_frame.paragraphs[0].add_run().text = "Alpha Overview"
    s1.placeholders[1].text_frame.paragraphs[0].add_run().text = "Paper template alpha"
    s2 = prs.slides.add_slide(prs.slide_layouts[1])
    s2.shapes.title.text_frame.paragraphs[0].add_run().text = "Alpha Content"
    body_tf = s2.placeholders[1].text_frame
    body_tf.paragraphs[0].add_run().text = "Alpha point one"
    p2 = body_tf.add_paragraph()
    p2.level = 1
    p2.add_run().text = "Alpha point two"
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    picture = s3.shapes.add_picture(
        io.BytesIO(_png_bytes(palette=2)), Inches(1), Inches(1), Inches(2)
    )
    picture.name = "alpha_pic"
    _add_named_textbox(
        s3, "alpha_box", Inches(4), Inches(2), Inches(4), Inches(1), "Alpha floating text"
    )
    prs.save(buf)
    source = zipfile.ZipFile(buf)
    return _rezip(
        source,
        OUT_DIR / "template_alpha.pptx",
        rewrites={
            "ppt/theme/theme1.xml": _retheme(
                source, "Paper Alpha", "Georgia", "Verdana", "AA3311"
            )
        },
        additions={},
    )


def build_template_beta() -> Path:
    """Import/rebind corpus, house style B ("Paper Beta": Courier New/Times New Roman,
    accent1 1166BB). Four slides: title, title-and-content, a chart (with embedded
    workbook) on the layout renamed "Beta Special" (a layout name alpha does NOT have),
    and a picture slide. All other layout names collide with alpha's by construction."""
    buf = io.BytesIO()
    prs = Presentation()
    beta_special = prs.slide_layouts[5]  # -- "Title Only" in the default template
    _expect(beta_special.name == "Title Only", "layout 5 is not 'Title Only'")
    beta_special.name = "Beta Special"

    s1 = prs.slides.add_slide(prs.slide_layouts[0])
    s1.shapes.title.text_frame.paragraphs[0].add_run().text = "Beta Overview"
    s1.placeholders[1].text_frame.paragraphs[0].add_run().text = "Paper template beta"
    s2 = prs.slides.add_slide(prs.slide_layouts[1])
    s2.shapes.title.text_frame.paragraphs[0].add_run().text = "Beta Content"
    body_tf = s2.placeholders[1].text_frame
    body_tf.paragraphs[0].add_run().text = "Beta point one"
    p2 = body_tf.add_paragraph()
    p2.level = 1
    p2.add_run().text = "Beta point two"
    s3 = prs.slides.add_slide(beta_special)
    s3.shapes.title.text_frame.paragraphs[0].add_run().text = "Beta Chart"
    chart_data = CategoryChartData()
    chart_data.categories = ["North", "South"]
    chart_data.add_series("FY26", (12.5, 8.75))
    frame = s3.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1.5), Inches(6), Inches(4), chart_data
    )
    frame.name = "beta_chart"
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    picture = s4.shapes.add_picture(
        io.BytesIO(_png_bytes(palette=3)), Inches(2), Inches(2), Inches(3)
    )
    picture.name = "beta_pic"
    prs.save(buf)
    source = zipfile.ZipFile(buf)
    return _rezip(
        source,
        OUT_DIR / "template_beta.pptx",
        rewrites={
            "ppt/theme/theme1.xml": _retheme(
                source, "Paper Beta", "Courier New", "Times New Roman", "1166BB"
            )
        },
        additions={},
    )


def build_lineage_trio() -> "tuple[Path, Path, Path]":
    """Diff ground-truth corpus (v0.11 Phase 6): v1, v2 saved FROM v1 with known edits,
    and a reorder-only variant. v2/reorder are built by loading v1's frozen bytes and
    applying shipped paper-pptx v0/v0.1 APIs, so permanent slide ids persist exactly as
    they do in real save-a-copy lineage. The exact edit list lives in the sidecars; the
    real-PowerPoint counterpart is FIXTURE-REQUESTS.md R13.
    """
    from pptx.edit import replace_text

    prs = Presentation()
    s1 = prs.slides.add_slide(prs.slide_layouts[1])
    s1.shapes.title.text_frame.paragraphs[0].add_run().text = "Lineage slide one"
    s1.placeholders[1].text_frame.paragraphs[0].add_run().text = "Constant body text."
    s2 = prs.slides.add_slide(prs.slide_layouts[1])
    s2.shapes.title.text_frame.paragraphs[0].add_run().text = "Lineage slide two"
    s2.placeholders[1].text_frame.paragraphs[0].add_run().text = "Second slide body."
    s2.notes_slide.notes_text_frame.text = "Original notes for slide two."
    s3 = prs.slides.add_slide(prs.slide_layouts[5])
    s3.shapes.title.text_frame.paragraphs[0].add_run().text = "Lineage slide three"
    chart_data = CategoryChartData()
    chart_data.categories = ["North", "South"]
    chart_data.add_series("FY", (10.0, 20.0))
    frame = s3.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1.5), Inches(6), Inches(4), chart_data
    )
    frame.name = "lineage_chart"
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    s4.shapes.add_picture(io.BytesIO(_png_bytes()), Inches(1), Inches(1), Inches(2)).name = (
        "lineage_pic"
    )
    _add_named_textbox(
        s4, "lineage_box", Inches(4), Inches(5), Inches(3), Inches(1), "Anchored box"
    )
    s5 = prs.slides.add_slide(prs.slide_layouts[1])
    s5.shapes.title.text_frame.paragraphs[0].add_run().text = "Lineage slide five"
    s5.placeholders[1].text_frame.paragraphs[0].add_run().text = "Fifth slide body."
    v1_path = _save(prs, "lineage_v1.pptx")

    # -- v2: load v1's bytes and apply the documented edit list with shipped APIs --------
    prs2 = Presentation(str(v1_path))
    replace_text(prs2, "Lineage slide one", "Lineage slide one, retitled")  # 1. text edit
    chart = prs2.slides[2].shapes.chart_by_name("lineage_chart")
    chart.replace_data_safe(["North", "South"], [("FY", (10.0, 25.0))])  # 2. chart data
    prs2.slides[1].replace_notes_text("Updated notes for slide two.")  # 3. notes edit
    box = prs2.slides[3].shapes.shape_by_name("lineage_box")
    box.left = Inches(5)  # 4. geometry change
    prs2.slides[3].shapes.picture_by_name("lineage_pic").replace_image(
        io.BytesIO(_png_bytes(palette=1))
    )  # 5. image replacement
    # -- add BEFORE delete: upstream allocates slide ids as max+1, so deleting the max id
    # -- (260, "five") first would hand the new slide that same id - id reuse would make
    # -- id-based diff matching read delete+add as one edited slide. That hazard is real
    # -- and documented for the diff organ, but this fixture pins the clean lineage case.
    s6 = prs2.slides.add_slide(prs2.slide_layouts[1])  # 6. add a new slide (id 261)
    s6.shapes.title.text_frame.paragraphs[0].add_run().text = "Lineage slide six, new"
    prs2.slides.delete(4)  # 7. delete "Lineage slide five" (id 260)
    prs2.slides.move(1, 0)  # 8. move "two" to the front
    v2_path = _save(prs2, "lineage_v2.pptx")

    # -- reorder-only variant: identical content, one move -------------------------------
    prs3 = Presentation(str(v1_path))
    prs3.slides.move(4, 0)
    reorder_path = _save(prs3, "lineage_reorder.pptx")
    return v1_path, v2_path, reorder_path


def main() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    builders = [
        build_minimal_clean,
        build_branded_template,
        build_clrmap_remap,
        build_chart_notes,
        build_shared_media,
        build_autofit_none,
        build_autofit_normal,
        build_autofit_shape,
        build_whitespace_pair,
        build_gauntlet,
        build_corrupt_dangling_sldid,
        build_large_smoke,
    ]
    for builder in builders:
        result = builder()
        for path in result if isinstance(result, tuple) else (result,):
            with zipfile.ZipFile(str(path)) as z:
                names = z.namelist()
            print("%s  (%d bytes, %d parts)" % (path.name, path.stat().st_size, len(names)))


if __name__ == "__main__":
    main()
