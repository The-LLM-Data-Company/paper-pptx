"""THE explicit golden-update command (CONVENTIONS §4: goldens update only via this).

Regenerates every golden inspection JSON under `tests/paper/goldens/` from the frozen fixture
corpus. Golden diffs are human-reviewed in the PR that lands them; tests never call this.

Usage:  python tests/paper/_authoring/update_goldens.py
"""

from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.inspect import inspect_deck, inspect_text

GOLDENS_DIR = Path(__file__).resolve().parent.parent / "goldens"
FIXTURES_DIR = Path(__file__).resolve().parent.parent / "fixtures"

#: (golden filename, fixture relpath, slide index) — one golden per entry
GOLDENS = (
    ("branded_template.inspect.json", "self_generated/branded_template.pptx", 0),
    ("clrmap_remap.inspect.json", "self_generated/clrmap_remap.pptx", 0),
    ("gauntlet_slide1.inspect.json", "self_generated/gauntlet.pptx", 0),
)

#: (golden filename, fixture relpath) — whole-deck structural manifests (Phase 2.1)
MANIFEST_GOLDENS = (
    ("gauntlet.manifest.json", "self_generated/gauntlet.pptx"),
    ("tables_in_group.manifest.json", "self_generated/tables_in_group.pptx"),
)


def golden_json(fixture_relpath: str, slide_index: int) -> str:
    """Return the canonical golden serialization for one slide's inspection."""
    prs = Presentation(str(FIXTURES_DIR / fixture_relpath))
    payload = inspect_text(prs.slides[slide_index]).to_dict()
    return json.dumps(payload, indent=2, ensure_ascii=False) + "\n"


def manifest_golden_json(fixture_relpath: str) -> str:
    """Return the canonical golden serialization for one deck's structural manifest."""
    prs = Presentation(str(FIXTURES_DIR / fixture_relpath))
    payload = inspect_deck(prs).to_dict()
    return json.dumps(payload, indent=2, ensure_ascii=False) + "\n"


def scrub_golden_json() -> str:
    """Return the canonical golden for the full scrub of the scrub_gauntlet fixture.

    v0.11 Phase 3: the report is the operation's evidence object; the golden pins its
    exact shape and the exact member budget of scrub-everything on the frozen fixture.
    """
    prs = Presentation(str(FIXTURES_DIR / "self_generated/scrub_gauntlet.pptx"))
    report = prs.scrub(
        notes=True,
        comments=True,
        metadata=True,
        hidden_slides=True,
        unused_layouts=True,
        unused_masters=True,
        unreachable_media=True,
        embedded_fonts=True,
    )
    return json.dumps(report.to_dict(), indent=2, ensure_ascii=False) + "\n"


def import_golden_json() -> str:
    """Return the canonical golden for one keep_appearance import (v0.11 Phase 5)."""
    dest = Presentation(str(FIXTURES_DIR / "self_generated/template_alpha.pptx"))
    source = Presentation(str(FIXTURES_DIR / "self_generated/template_beta.pptx"))
    report = dest.import_slide(source, 0, mode="keep_appearance")
    return json.dumps(report.to_dict(), indent=2, ensure_ascii=False) + "\n"


def diff_golden_json() -> str:
    """Return the canonical golden for the lineage v1->v2 diff (v0.11 Phase 6)."""
    from pptx.diff import diff_decks

    report = diff_decks(
        str(FIXTURES_DIR / "self_generated/lineage_v1.pptx"),
        str(FIXTURES_DIR / "self_generated/lineage_v2.pptx"),
        detail="text",
    )
    return json.dumps(report.to_dict(), indent=2, ensure_ascii=False) + "\n"


def main() -> None:
    GOLDENS_DIR.mkdir(parents=True, exist_ok=True)
    for golden_name, fixture_relpath, slide_index in GOLDENS:
        out = GOLDENS_DIR / golden_name
        out.write_text(golden_json(fixture_relpath, slide_index), encoding="utf-8")
        print("wrote", out)
    for golden_name, fixture_relpath in MANIFEST_GOLDENS:
        out = GOLDENS_DIR / golden_name
        out.write_text(manifest_golden_json(fixture_relpath), encoding="utf-8")
        print("wrote", out)
    out = GOLDENS_DIR / "scrub_gauntlet.scrub.json"
    out.write_text(scrub_golden_json(), encoding="utf-8")
    print("wrote", out)
    out = GOLDENS_DIR / "import_beta_keep.import.json"
    out.write_text(import_golden_json(), encoding="utf-8")
    print("wrote", out)
    out = GOLDENS_DIR / "lineage_v1_v2.diff.json"
    out.write_text(diff_golden_json(), encoding="utf-8")
    print("wrote", out)


if __name__ == "__main__":
    main()
