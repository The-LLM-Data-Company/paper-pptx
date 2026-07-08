"""Phase 7 contract tests: slide clone / delete / reorder / move.

This is where "opens in python-pptx but not in PowerPoint" corruption gets manufactured, so
the tests lean hardest on the oracles: exact changed-part budgets, relationship-integrity
scans on every output, LibreOffice smoke on every operation class, and cross-contamination
proofs (mutating a clone's chart leaves the original chart XML byte-identical).
"""

from __future__ import annotations

import io

import pytest
from lxml import etree

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.errors import PaperRefusal, RelationshipPolicyError, TargetNotFoundError
from pptx.slide import SlideClonePolicy
from pptx.util import Inches

from . import corpus
from .contract import (
    assert_changed_parts,
    assert_refusal_atomic,
    save_to_bytes,
    zip_member_map,
)
from .lo import lo_load_smoke
from .relint import dangling_relationship_targets, missing_relationship_references

CHART_NOTES = "self_generated/chart_notes.pptx"
SHARED_MEDIA = "self_generated/shared_media.pptx"
GAUNTLET = "self_generated/gauntlet.pptx"
MINIMAL = "self_generated/minimal_clean.pptx"
LO_CHART_NOTES = "libreoffice_export/lo_chart_notes.pptx"
_RELS_NS = "{http://schemas.openxmlformats.org/package/2006/relationships}"


def _open(relpath):
    return Presentation(str(corpus.fixture_path(relpath)))


def _assert_relationship_integrity(pptx_bytes):
    zip_map = zip_member_map(pptx_bytes)
    assert dangling_relationship_targets(zip_map) == []
    assert missing_relationship_references(zip_map) == []


def _reopen(pptx_bytes):
    return Presentation(io.BytesIO(pptx_bytes))


# ------------------------------------------------------------------------------------ clone


def test_clone_deep_copies_chart_workbook_and_notes_with_exact_budget():
    prs = _open(CHART_NOTES)
    before = save_to_bytes(prs)
    clone = prs.slides.clone(0)
    assert prs.slides.index(clone) == 1
    after = save_to_bytes(prs)

    assert_changed_parts(
        before,
        after,
        expect_changed=[
            "[Content_Types].xml",  # -- Override entries for the new parts
            "ppt/_rels/presentation.xml.rels",
            "ppt/presentation.xml",
        ],
        expect_added=[
            "ppt/charts/_rels/chart2.xml.rels",
            "ppt/charts/chart2.xml",
            "ppt/embeddings/Microsoft_Excel_Sheet2.xlsx",
            "ppt/notesSlides/_rels/notesSlide2.xml.rels",
            "ppt/notesSlides/notesSlide2.xml",
            "ppt/slides/_rels/slide2.xml.rels",
            "ppt/slides/slide2.xml",
        ],
    )
    _assert_relationship_integrity(after)


def test_mutating_the_clones_chart_leaves_the_original_chart_byte_identical():
    """THE cross-contamination test: the corruption class this rewrite makes impossible."""
    prs = _open(CHART_NOTES)
    prs.slides.clone(0)
    reopened = _reopen(save_to_bytes(prs))
    original_chart_xml = zip_member_map(save_to_bytes(reopened))["ppt/charts/chart1.xml"]

    clone_chart = next(s for s in reopened.slides[1].shapes if s.has_chart).chart
    chart_data = CategoryChartData()
    chart_data.categories = ["X", "Y", "Z"]
    chart_data.add_series("Mutated", (1.0, 2.0, 3.0))
    clone_chart.replace_data(chart_data)

    after_map = zip_member_map(save_to_bytes(reopened))
    assert after_map["ppt/charts/chart1.xml"] == original_chart_xml
    original_chart = next(s for s in reopened.slides[0].shapes if s.has_chart).chart
    assert [series.name for series in original_chart.series] == ["Q1", "Q2"]


def test_clone_notes_are_neither_dropped_nor_cross_linked():
    prs = _open(CHART_NOTES)
    prs.slides.clone(0)
    saved = save_to_bytes(prs)
    zip_map = zip_member_map(saved)
    for rels_member, expected_slide in (
        ("ppt/notesSlides/_rels/notesSlide1.xml.rels", "slide1.xml"),
        ("ppt/notesSlides/_rels/notesSlide2.xml.rels", "slide2.xml"),
    ):
        rels = etree.fromstring(zip_map[rels_member])
        slide_targets = [
            rel.get("Target")
            for rel in rels.iter(_RELS_NS + "Relationship")
            if rel.get("Type").endswith("/slide")
        ]
        assert len(slide_targets) == 1
        assert slide_targets[0].endswith(expected_slide), rels_member

    reopened = _reopen(saved)
    reopened.slides[1].replace_notes_text("clone-only notes")
    assert reopened.slides[0].read_notes_text() == "Speaker notes for the clone fixture."
    assert reopened.slides[1].read_notes_text() == "clone-only notes"


def test_clone_shares_media_by_default_and_copies_on_request():
    prs = _open(SHARED_MEDIA)
    prs.slides.clone(0)
    shared_map = zip_member_map(save_to_bytes(prs))
    assert [n for n in shared_map if n.startswith("ppt/media/")] == ["ppt/media/image1.png"]

    prs = _open(SHARED_MEDIA)
    prs.slides.clone(0, policy=SlideClonePolicy(share_media=False))
    copied_map = zip_member_map(save_to_bytes(prs))
    assert sorted(n for n in copied_map if n.startswith("ppt/media/")) == [
        "ppt/media/image1.png",
        "ppt/media/image2.png",
    ]
    _assert_relationship_integrity(save_to_bytes(prs))


def test_clone_can_drop_notes_by_policy_without_touching_the_source():
    prs = _open(CHART_NOTES)
    clone = prs.slides.clone(0, policy=SlideClonePolicy(deep_copy_notes=False))
    assert clone.has_notes_slide is False
    assert prs.slides[0].has_notes_slide is True
    _assert_relationship_integrity(save_to_bytes(prs))


def test_clone_copies_external_hyperlink_relationships():
    prs = _open(GAUNTLET)
    prs.slides.clone(3)  # -- the hyperlink slide
    zip_map = zip_member_map(save_to_bytes(prs))
    rels = etree.fromstring(zip_map["ppt/slides/_rels/slide5.xml.rels"])
    external = [
        rel.get("Target")
        for rel in rels.iter(_RELS_NS + "Relationship")
        if rel.get("TargetMode") == "External"
    ]
    assert external == ["https://example.com/paper"]
    _assert_relationship_integrity(save_to_bytes(prs))


def test_clone_of_a_libreoffice_chart_without_workbook_copies_style_parts():
    """LO charts carry colors/style parts and no embedded workbook; clone must cope."""
    prs = _open(LO_CHART_NOTES)
    before = save_to_bytes(prs)
    prs.slides.clone(0)
    saved = save_to_bytes(prs)
    assert_changed_parts(
        before,
        saved,
        expect_changed=[
            "[Content_Types].xml",
            "ppt/_rels/presentation.xml.rels",
            "ppt/presentation.xml",
        ],
        expect_added=[
            "ppt/charts/_rels/chart2.xml.rels",
            "ppt/charts/chart2.xml",
            "ppt/charts/colors2.xml",
            "ppt/charts/style2.xml",
            "ppt/notesSlides/_rels/notesSlide2.xml.rels",
            "ppt/notesSlides/notesSlide2.xml",
            "ppt/slides/_rels/slide2.xml.rels",
            "ppt/slides/slide2.xml",
        ],
    )
    assert not any(n.startswith("ppt/embeddings/") for n in zip_member_map(saved))
    _assert_relationship_integrity(saved)


def test_clone_after_parameter_positions_the_copy():
    prs = _open(GAUNTLET)
    clone = prs.slides.clone(0, after=2)
    clone_id = clone.slide_id
    reopened = _reopen(save_to_bytes(prs))
    assert reopened.slides[3].slide_id == clone_id

    prs = _open(GAUNTLET)
    clone = prs.slides.clone(0, after=prs.slides[1])  # -- Slide-typed after
    assert prs.slides.index(clone) == 2


def test_clone_with_two_charts_allocates_distinct_partnames():
    """Regression: parts created mid-clone are invisible to package partname allocation, so
    two deep-copied charts used to receive the SAME partname — duplicate zip members and one
    chart's data silently clobbered."""
    prs = _open(CHART_NOTES)
    chart_data = CategoryChartData()
    chart_data.categories = ["a", "b"]
    chart_data.add_series("Second", (7.0, 8.0))
    frame = prs.slides[0].shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1), Inches(3), Inches(2), chart_data
    )
    frame.name = "second_chart"

    prs.slides.clone(0)
    saved = save_to_bytes(prs)
    zip_map = zip_member_map(saved)  # -- asserts no duplicate member names by itself
    charts = sorted(
        n for n in zip_map if n.startswith("ppt/charts/chart") and n.endswith(".xml")
    )
    workbooks = sorted(n for n in zip_map if n.startswith("ppt/embeddings/"))
    assert len(charts) == 4
    assert len(workbooks) == 4
    _assert_relationship_integrity(saved)

    reopened = _reopen(saved)
    clone_series = sorted(
        series.name
        for shape in reopened.slides[1].shapes
        if shape.has_chart
        for series in shape.chart.series
    )
    assert clone_series == ["Q1", "Q2", "Second"]  # -- neither chart clobbered the other


def test_clone_with_two_unshared_images_allocates_distinct_partnames():
    import io as io_module

    from PIL import Image as PILImage

    def png(color):
        buf = io_module.BytesIO()
        PILImage.new("RGB", (16, 16), color).save(buf, format="PNG")
        return buf.getvalue()

    prs = _open(MINIMAL)
    slide = prs.slides[0]
    slide.shapes.add_picture(io_module.BytesIO(png((250, 0, 0))), 0, 0, 914400)
    slide.shapes.add_picture(io_module.BytesIO(png((0, 250, 0))), 914400, 0, 914400)

    prs.slides.clone(0, policy=SlideClonePolicy(share_media=False))
    saved = save_to_bytes(prs)
    zip_map = zip_member_map(saved)
    media = sorted(n for n in zip_map if n.startswith("ppt/media/"))
    assert len(media) == 4
    assert len({zip_map[n] for n in media}) == 2  # -- two distinct images, each twice
    _assert_relationship_integrity(saved)

    reopened = _reopen(saved)
    clone_blobs = {
        s.image.blob
        for s in reopened.slides[1].shapes
        if s.shape_type.name == "PICTURE"
    }
    assert clone_blobs == {png((250, 0, 0)), png((0, 250, 0))}


def test_cloning_a_clone_and_repeated_clones_stay_consistent():
    prs = _open(CHART_NOTES)
    first_clone = prs.slides.clone(0)
    prs.slides.clone(prs.slides.index(first_clone))  # -- clone the clone
    prs.slides.clone(0)
    saved = save_to_bytes(prs)
    zip_map = zip_member_map(saved)
    assert len([n for n in zip_map if n.startswith("ppt/charts/chart")]) >= 4
    _assert_relationship_integrity(saved)
    assert len(_reopen(saved).slides) == 4


def test_clone_refuses_chart_with_unsupported_child_relationship():
    """A chart part related to something outside the allowed child set refuses atomically."""
    prs = _open(CHART_NOTES)
    chart_part = next(s for s in prs.slides[0].shapes if s.has_chart).chart.part
    image_part = prs.slides[0].part.package.get_or_add_image_part(
        io.BytesIO(
            Presentation(str(corpus.fixture_path(SHARED_MEDIA)))
            .slides[0]
            .shapes[0]
            .image.blob
        )
    )
    chart_part.relate_to(image_part, "http://example.com/relationships/bogus")

    raised = assert_refusal_atomic(prs, lambda p: p.slides.clone(0), RelationshipPolicyError)
    assert "chart part" in str(raised)


def test_clone_refuses_notes_with_unsupported_child_relationship():
    prs = _open(CHART_NOTES)
    notes_part = prs.slides[0].part.part_related_by(
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide"
    )
    other_slide_part = prs.slides[0].part
    notes_part.relate_to(other_slide_part, "http://example.com/relationships/bogus")

    raised = assert_refusal_atomic(prs, lambda p: p.slides.clone(0), RelationshipPolicyError)
    assert "notes slide" in str(raised)


def test_clone_refuses_to_share_charts():
    prs = _open(CHART_NOTES)
    raised = assert_refusal_atomic(
        prs,
        lambda p: p.slides.clone(0, policy=SlideClonePolicy(deep_copy_charts=False)),
        RelationshipPolicyError,
    )
    assert "cross-contamination" in str(raised)
    assert isinstance(raised, PaperRefusal)


def test_clone_refuses_unsupported_relationship_types_atomically():
    prs = _open(MINIMAL)
    prs.slides[0].part.add_embedded_ole_object_part("Excel.Sheet.12", io.BytesIO(b"fake-ole"))
    raised = assert_refusal_atomic(prs, lambda p: p.slides.clone(0), RelationshipPolicyError)
    assert "does not support" in str(raised)


def test_clone_rejects_foreign_slides_and_bad_policy():
    prs = _open(GAUNTLET)
    other = _open(MINIMAL)
    with pytest.raises(TargetNotFoundError):
        prs.slides.clone(other.slides[0])
    with pytest.raises(ValueError):
        prs.slides.clone(0, policy="deep")


# ----------------------------------------------------------------------------------- delete


def test_delete_removes_slide_and_its_unshared_parts_with_exact_budget():
    prs = _open(CHART_NOTES)
    before = save_to_bytes(prs)
    prs.slides.delete(0)
    after = save_to_bytes(prs)
    assert_changed_parts(
        before,
        after,
        expect_changed=[
            "[Content_Types].xml",
            "ppt/_rels/presentation.xml.rels",
            "ppt/presentation.xml",
        ],
        expect_removed=[
            "ppt/charts/_rels/chart1.xml.rels",
            "ppt/charts/chart1.xml",
            "ppt/embeddings/Microsoft_Excel_Sheet1.xlsx",
            "ppt/notesSlides/_rels/notesSlide1.xml.rels",
            "ppt/notesSlides/notesSlide1.xml",
            "ppt/slides/_rels/slide1.xml.rels",
            "ppt/slides/slide1.xml",
        ],
    )
    _assert_relationship_integrity(after)
    assert len(_reopen(after).slides) == 0


def test_delete_keeps_media_shared_with_surviving_slides():
    prs = _open(SHARED_MEDIA)
    prs.slides.delete(0)
    zip_map = zip_member_map(save_to_bytes(prs))
    assert [n for n in zip_map if n.startswith("ppt/media/")] == ["ppt/media/image1.png"]
    _assert_relationship_integrity(save_to_bytes(prs))


def test_global_relationship_scan_after_every_gauntlet_delete():
    """Delete each gauntlet slide in turn; no output may carry a dangling reference."""
    for index in range(4):
        prs = _open(GAUNTLET)
        prs.slides.delete(index)
        saved = save_to_bytes(prs)
        _assert_relationship_integrity(saved)
        assert len(_reopen(saved).slides) == 3


def test_deleting_the_last_slide_leaves_a_valid_empty_deck():
    prs = _open(MINIMAL)
    prs.slides.delete(0)
    assert len(_reopen(save_to_bytes(prs)).slides) == 0


# --------------------------------------------------------------------------- reorder / move


def test_reorder_permutes_slides_and_round_trips():
    prs = _open(GAUNTLET)
    titles_before = [s.shapes.title.text if s.shapes.title else None for s in prs.slides]
    before = save_to_bytes(prs)
    prs.slides.reorder([2, 0, 3, 1])
    after = save_to_bytes(prs)
    assert_changed_parts(before, after, expect_changed=["ppt/presentation.xml"])

    titles_after = [
        s.shapes.title.text if s.shapes.title else None for s in _reopen(after).slides
    ]
    assert titles_after == [titles_before[i] for i in [2, 0, 3, 1]]


@pytest.mark.parametrize("bad_order", [[0, 1, 2], [0, 1, 2, 2], [0, 1, 2, 4], []])
def test_reorder_rejects_non_permutations_atomically(bad_order):
    prs = _open(GAUNTLET)
    before = save_to_bytes(prs)
    with pytest.raises(ValueError):
        prs.slides.reorder(bad_order)
    assert_changed_parts(before, save_to_bytes(prs))


def test_move_repositions_a_single_slide():
    prs = _open(GAUNTLET)
    last = prs.slides[3]
    moved_id = last.slide_id
    before = save_to_bytes(prs)
    prs.slides.move(last, 0)
    after = save_to_bytes(prs)
    assert_changed_parts(before, after, expect_changed=["ppt/presentation.xml"])
    assert _reopen(after).slides[0].slide_id == moved_id
    with pytest.raises(ValueError):
        prs.slides.move(0, 99)
    with pytest.raises(ValueError):
        prs.slides.move(0, -1)


def test_delete_error_paths_leave_the_deck_untouched():
    prs = _open(GAUNTLET)
    before = save_to_bytes(prs)
    with pytest.raises(IndexError):
        prs.slides.delete(99)
    other = _open(MINIMAL)
    with pytest.raises(TargetNotFoundError):
        prs.slides.delete(other.slides[0])
    assert_changed_parts(before, save_to_bytes(prs))  # -- empty budget


# --------------------------------------------------------------------------------- lo_smoke


@pytest.mark.lo_smoke
@pytest.mark.parametrize("operation", ["clone", "delete", "reorder"])
def test_slide_operation_outputs_load_in_libreoffice(operation, tmp_path):
    prs = _open(GAUNTLET)
    if operation == "clone":
        prs.slides.clone(1)  # -- the chart+notes slide, the hardest case
    elif operation == "delete":
        prs.slides.delete(1)
    else:
        prs.slides.reorder([3, 2, 1, 0])
    out = tmp_path / ("%s.pptx" % operation)
    prs.save(str(out))
    lo_load_smoke(out, tmp_path)
