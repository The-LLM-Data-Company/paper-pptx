"""Contract tests: the structural deck manifest (pptx.inspect.inspect_deck)."""

from __future__ import annotations

import json

import pytest

from pptx import Presentation
from pptx.inspect import DeckManifest, inspect_deck

from . import corpus
from .contract import snapshot_parts

GAUNTLET = "self_generated/gauntlet.pptx"


def _open(relpath):
    return Presentation(str(corpus.fixture_path(relpath)))


def test_manifest_surveys_the_gauntlet_against_its_sidecar():
    manifest = inspect_deck(_open(GAUNTLET))
    sidecar = corpus.load_sidecar(GAUNTLET)["ground_truth"]
    assert isinstance(manifest, DeckManifest)
    assert manifest.slide_count == sidecar["slide_count"]
    assert manifest.slides[0].layout_name == "Title and Content"
    assert manifest.slides[1].has_notes is True

    all_shapes = [shape for slide in manifest.slides for shape in slide.shapes]
    chart = next(s.chart for s in all_shapes if s.chart)
    assert chart["chart_type"] == "COLUMN_CLUSTERED"
    table = next(s.table for s in all_shapes if s.table)
    assert table == {
        "rows": sidecar["table"]["rows"],
        "cols": sidecar["table"]["cols"],
    }
    image = next(s.image for s in all_shapes if s.image)
    assert image["ext"] == "png"
    assert image["natural_size_px"] == [64, 64]
    autofit = next(
        s.autofit for s in all_shapes if s.autofit and s.autofit["font_scale"] is not None
    )
    assert autofit == {
        "mode": "TEXT_TO_FIT_SHAPE",
        "font_scale": 62.5,
        "line_space_reduction": 20.0,
    }


def test_manifest_nests_group_children_and_reports_masters():
    manifest = inspect_deck(_open("self_generated/nested_groups.pptx"))
    group = next(s for s in manifest.slides[0].shapes if s.kind == "GROUP")
    child_names = [child.name for child in group.children]
    assert "level1_box" in child_names
    inner = next(c for c in group.children if c.kind == "GROUP")
    assert [c.name for c in inner.children if c.kind != "GROUP"] == ["level2_box"]
    assert manifest.masters[0]["layouts"][0] == "Title Slide"


def test_manifest_reports_placeholder_geometry_via_upstream_inheritance():
    """Placeholder position/size is the one inheritance upstream already resolves — the
    manifest reports the resolved value, matching the layout placeholder's geometry."""
    prs = _open("self_generated/branded_template.pptx")
    manifest = inspect_deck(prs)
    title = next(s for s in manifest.slides[0].shapes if s.placeholder_type == "TITLE")
    layout_title = next(
        ph for ph in prs.slides[0].slide_layout.placeholders
        if ph.placeholder_format.idx == 0
    )
    assert title.x == int(layout_title.left)
    assert title.cx == int(layout_title.width)


@pytest.mark.parametrize(
    ("golden_name", "fixture_relpath"),
    [
        ("gauntlet.manifest.json", GAUNTLET),
        ("tables_in_group.manifest.json", "self_generated/tables_in_group.pptx"),
    ],
)
def test_manifest_matches_frozen_golden(golden_name, fixture_relpath):
    golden_path = corpus.FIXTURES_DIR.parent / "goldens" / golden_name
    payload = inspect_deck(_open(fixture_relpath)).to_dict()
    actual = (json.dumps(payload, indent=2, ensure_ascii=False) + "\n").encode("utf-8")
    assert actual == golden_path.read_bytes(), (
        "manifest drifted from golden %s; regenerate ONLY via"
        " tests/paper/_authoring/update_goldens.py and review the diff" % golden_name
    )


def test_manifest_is_deterministic_and_read_only():
    prs = _open(GAUNTLET)
    before = snapshot_parts(prs)
    first = json.dumps(inspect_deck(prs).to_dict())
    second = json.dumps(inspect_deck(prs).to_dict())
    assert first == second
    assert snapshot_parts(prs) == before  # -- inspection never mutates


def test_manifest_payload_carries_pinned_schema():
    payload = inspect_deck(_open(GAUNTLET)).to_dict()
    assert payload["schema"] == "paper-deck-manifest"
    assert payload["version"] == 1
    assert payload["slide_width"] == 9144000
