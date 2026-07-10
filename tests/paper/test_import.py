"""Contract tests: slide import and deck merge.

The plan's required list: cross-contamination (source byte-identical after edits to the
import), dedupe (three slides from one source -> one transplanted master), every mode
against the two-template corpus, relint + section scan + LO smoke on all outputs, and
determinism goldens on the import report.
"""

from __future__ import annotations

import io
import json

import pytest

from pptx import Presentation
from pptx.errors import (
    PaperRefusal,
    RelationshipPolicyError,
    TargetNotFoundError,
    UnsupportedStructureError,
)

from . import corpus
from .contract import assert_changed_parts, save_reopen, save_to_bytes, zip_member_map
from .idlists import dangling_section_slide_ids, duplicate_section_slide_ids
from .lo import lo_load_smoke
from .relint import dangling_relationship_targets, missing_relationship_references

ALPHA = "self_generated/template_alpha.pptx"
BETA = "self_generated/template_beta.pptx"
LO_ALPHA = "libreoffice_export/lo_template_alpha.pptx"
SECTIONS = "self_generated/sections.pptx"
GAUNTLET = "self_generated/gauntlet.pptx"


def _open(relpath):
    return Presentation(str(corpus.fixture_path(relpath)))


def _assert_clean(saved_bytes):
    zip_map = zip_member_map(saved_bytes)
    assert dangling_relationship_targets(zip_map) == []
    assert missing_relationship_references(zip_map) == []
    assert dangling_section_slide_ids(zip_map) == []
    assert duplicate_section_slide_ids(zip_map) == []


# --------------------------------------------------------------------------- mode contracts


def test_keep_appearance_transplants_chain_with_zero_shifts():
    """The keep-appearance invariant: identical chain, identical resolution - no shifts."""
    dest = _open(ALPHA)
    source = _open(BETA)
    report = dest.import_slide(source, 0, mode="keep_appearance")
    assert report.run_shifts == ()
    assert report.layout_binding_method == "transplant"

    saved = save_to_bytes(dest)
    _assert_clean(saved)
    reopened = Presentation(io.BytesIO(saved))
    assert len(reopened.slide_masters) == 2
    imported = reopened.slides[3]
    assert imported.slide_layout.slide_master is not reopened.slide_masters[0]
    # -- the beta theme travelled: title resolves to beta's Courier New in the dest
    title_font = imported.shapes.title.text_frame.paragraphs[0].runs[0].effective_font()
    assert title_font.name.value == "Courier New"


def test_keep_appearance_dedupes_master_across_three_imports():
    """Three imports from one source: ONE new master and ONE new theme, reused twice."""
    dest = _open(ALPHA)
    source = _open(BETA)
    reports = [dest.import_slide(source, i, mode="keep_appearance") for i in range(3)]
    assert reports[0].parts_reused == ()
    for later in reports[1:]:
        assert "/ppt/slideMasters/slideMaster2.xml" in later.parts_reused

    saved = save_to_bytes(dest)
    zip_map = zip_member_map(saved)
    masters = [
        m
        for m in zip_map
        if m.startswith("ppt/slideMasters/slideMaster") and m.endswith(".xml")
    ]
    themes = [m for m in zip_map if m.startswith("ppt/theme/") and m.endswith(".xml")]
    assert sorted(masters) == [
        "ppt/slideMasters/slideMaster1.xml",
        "ppt/slideMasters/slideMaster2.xml",
    ]
    assert sorted(themes) == ["ppt/theme/theme1.xml", "ppt/theme/theme2.xml"]
    _assert_clean(saved)
    # -- the transplanted master accumulated all three layouts
    reopened = Presentation(io.BytesIO(saved))
    assert len(reopened.slide_masters[1].slide_layouts) == 3


def test_adopt_theme_takes_house_style_and_reports_shifts():
    dest = _open(ALPHA)
    report = dest.import_slide(_open(BETA), 0, mode="adopt_theme")
    assert report.layout_binding_method == "name-match"  # -- "Title Slide" collides
    assert report.layout_binding == "/ppt/slideLayouts/slideLayout1.xml"
    shifted_names = {
        (s.before["name"]["value"], s.after["name"]["value"]) for s in report.run_shifts
    }
    assert ("Courier New", "Georgia") in shifted_names  # -- title: beta major -> alpha major
    assert ("Times New Roman", "Verdana") in shifted_names

    saved = save_to_bytes(dest)
    _assert_clean(saved)
    reopened = Presentation(io.BytesIO(saved))
    assert len(reopened.slide_masters) == 1  # -- nothing transplanted
    imported = reopened.slides[3]
    title_font = imported.shapes.title.text_frame.paragraphs[0].runs[0].effective_font()
    assert title_font.name.value == "Georgia"  # -- the house look


def test_adopt_theme_falls_back_to_type_match_for_renamed_layout():
    """Beta's chart slide sits on 'Beta Special' - a name alpha lacks - but the layout
    kept type="titleOnly", so auto-selection falls back to alpha's Title Only."""
    dest = _open(ALPHA)
    report = dest.import_slide(_open(BETA), 2, mode="adopt_theme")
    assert report.layout_binding_method == "type-match"
    assert report.layout_binding == "/ppt/slideLayouts/slideLayout6.xml"


def test_adopt_theme_unmatched_layout_refuses_and_explicit_target_recovers():
    dest = _open(ALPHA)
    source = _open(BETA)
    # -- force a truly unmatchable source layout: alien name AND no type token
    source_layout = source.slides[2].slide_layout
    source_layout._element.attrib.pop("type", None)

    before = save_to_bytes(dest)
    with pytest.raises(UnsupportedStructureError, match="target_layout"):
        dest.import_slide(source, 2, mode="adopt_theme")
    assert_changed_parts(before, save_to_bytes(dest))  # -- empty budget

    report = dest.import_slide(
        source, 2, mode="adopt_theme", target_layout=dest.slide_layouts[5]
    )
    assert report.layout_binding_method == "explicit"
    saved = save_to_bytes(dest)
    _assert_clean(saved)
    reopened = Presentation(io.BytesIO(saved))
    assert reopened.slides[3].shapes.chart_by_name("beta_chart") is not None


def test_bake_freezes_look_without_importing_masters():
    dest = _open(ALPHA)
    report = dest.import_slide(_open(BETA), 1, mode="bake")
    assert report.run_shifts == ()  # -- baked: resolution cannot shift
    assert set(report.baked_shapes) == {"Title 1", "Content Placeholder 2"}

    saved = save_to_bytes(dest)
    _assert_clean(saved)
    reopened = Presentation(io.BytesIO(saved))
    assert len(reopened.slide_masters) == 1
    imported = reopened.slides[3]
    assert all(not s.is_placeholder for s in imported.shapes)  # -- all free shapes now
    title = next(s for s in imported.shapes if s.name == "Title 1")
    run = title.text_frame.paragraphs[0].runs[0]
    assert run.font.name == "Courier New"  # -- beta's look, made local
    assert run.font.size is not None


def test_bake_drops_furniture_placeholders():
    source = _open(GAUNTLET)
    source.apply_footers(footer="Travelling footer", slide_number=True)
    dest = _open(ALPHA)
    report = dest.import_slide(source, 0, mode="bake")
    assert len(report.dropped_placeholders) == 2  # -- ftr + sldNum
    # -- dropping shapes must not slip the shift keys: surviving runs are baked and
    # -- re-resolve identically, so there are NO shifts, phantom or real
    assert report.run_shifts == ()
    reopened = save_reopen(dest)
    imported = reopened.slides[3]
    texts = [s.text_frame.text for s in imported.shapes if s.has_text_frame]
    assert "Travelling footer" not in texts


# ------------------------------------------------------------------- cross-contamination


def test_source_is_never_mutated_and_imported_chart_is_independent():
    """Edit the imported chart; the source presentation stays byte-identical."""
    dest = _open(ALPHA)
    source = _open(BETA)
    source_before = save_to_bytes(source)
    dest.import_slide(source, 2, mode="keep_appearance")

    chart = dest.slides[3].shapes.chart_by_name("beta_chart")
    chart.replace_data_safe(["North", "South"], [("FY26", (99.0, 1.0))])

    assert save_to_bytes(source) == source_before
    reopened_source = Presentation(io.BytesIO(save_to_bytes(source)))
    source_chart = reopened_source.slides[2].shapes.chart_by_name("beta_chart")
    values = [pt for series in source_chart.plots[0].series for pt in series.values]
    assert values == [12.5, 8.75]  # -- source data untouched


def test_media_always_copies_never_shared_across_packages():
    dest = _open(ALPHA)
    source = _open(BETA)
    report = dest.import_slide(source, 3, mode="adopt_theme")  # -- beta picture slide
    assert any("/ppt/media/" in part for part in report.parts_added)
    saved = save_to_bytes(dest)
    reopened = Presentation(io.BytesIO(saved))
    imported_pic = reopened.slides[3].shapes.picture_by_name("beta_pic")
    source_pic = source.slides[3].shapes.picture_by_name("beta_pic")
    assert imported_pic.image.blob == source_pic.image.blob  # -- same bytes, copied part


def test_notes_policy():
    source = _open(GAUNTLET)  # -- slide 2 has speaker notes
    dest = _open(ALPHA)
    with_notes = dest.import_slide(source, 1, mode="adopt_theme", notes=True)
    assert with_notes.notes_copied is True
    reopened = save_reopen(dest)
    assert reopened.slides[3].read_notes_text() == "Gauntlet speaker notes."

    dest2 = _open(ALPHA)
    without = dest2.import_slide(source, 1, mode="adopt_theme", notes=False)
    assert without.notes_copied is False
    reopened2 = save_reopen(dest2)
    assert not reopened2.slides[3].has_notes_slide


# ------------------------------------------------------------------ position and sections


def test_position_inserts_at_index():
    dest = _open(ALPHA)
    report = dest.import_slide(_open(BETA), 0, mode="adopt_theme", position=0)
    assert report.position == 0
    reopened = save_reopen(dest)
    assert reopened.slides[0].shapes.title.text_frame.text == "Beta Overview"
    assert len(reopened.slides) == 4


def test_section_enrollment_named_and_adjacent():
    dest = _open(SECTIONS)
    source = _open(BETA)
    report = dest.import_slide(source, 0, mode="adopt_theme", section="Close")
    assert report.section == "Close"
    saved = save_to_bytes(dest)
    _assert_clean(saved)  # -- the id-list scan proves every section id resolves

    # -- adjacent enrollment: inserting at position 1 lands in "Intro" (slide 1's
    # -- section), DIRECTLY AFTER slide 1's entry, and the report says which section
    dest2 = _open(SECTIONS)
    report2 = dest2.import_slide(source, 0, mode="adopt_theme", position=1)
    assert report2.section == "Intro"  # -- actual enrollment, not the (None) argument
    saved2 = save_to_bytes(dest2)
    _assert_clean(saved2)
    from lxml import etree

    presentation = etree.fromstring(zip_member_map(saved2)["ppt/presentation.xml"])
    P14 = "http://schemas.microsoft.com/office/powerpoint/2010/main"
    intro = next(
        s
        for s in presentation.findall(".//{%s}section" % P14)
        if s.get("name") == "Intro"
    )
    intro_ids = [e.get("id") for e in intro.findall(".//{%s}sldId" % P14)]
    P = "http://schemas.openxmlformats.org/presentationml/2006/main"
    deck_ids = [
        e.get("id") for e in presentation.findall(".//{%s}sldIdLst/{%s}sldId" % (P, P))
    ]
    # -- section order mirrors deck order: [slide 1's id, the imported slide's id]
    assert intro_ids == [deck_ids[0], deck_ids[1]]


def test_missing_section_refuses_atomically():
    dest = _open(SECTIONS)
    before = save_to_bytes(dest)
    with pytest.raises(TargetNotFoundError):
        dest.import_slide(_open(BETA), 0, mode="adopt_theme", section="No Such Section")
    assert_changed_parts(before, save_to_bytes(dest))  # -- empty budget


# --------------------------------------------------------------------------- append_deck


def test_append_deck_imports_all_slides_in_order():
    dest = _open(ALPHA)
    reports = dest.append_deck(_open(BETA), mode="keep_appearance")
    assert len(reports) == 4
    saved = save_to_bytes(dest)
    _assert_clean(saved)
    reopened = Presentation(io.BytesIO(saved))
    assert len(reopened.slides) == 7
    assert len(reopened.slide_masters) == 2  # -- dedupe held across the whole merge
    titles = [
        s.shapes.title.text_frame.text if s.shapes.title is not None else None
        for s in reopened.slides
    ]
    assert titles[3] == "Beta Overview"
    assert titles[5] == "Beta Chart"


def test_append_deck_validates_whole_source_before_first_write():
    """Poison the LAST source slide; the destination must stay untouched."""
    dest = _open(ALPHA)
    source = _open(BETA)
    source.slides[3].part.relate_to(
        source.slides[0].part,
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide",
    )
    before = save_to_bytes(dest)
    with pytest.raises(RelationshipPolicyError):
        dest.append_deck(source, mode="keep_appearance")
    assert_changed_parts(before, save_to_bytes(dest))  # -- empty budget


# ------------------------------------------------------------------------------- refusals


def test_refusal_ledger_unsupported_relationship():
    dest = _open(ALPHA)
    source = _open(BETA)
    source.slides[0].part.relate_to(
        source.slides[1].part,
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide",
    )
    before = save_to_bytes(dest)
    with pytest.raises(RelationshipPolicyError) as excinfo:
        dest.import_slide(source, 0, mode="keep_appearance")
    assert "refusal ledger" in str(excinfo.value)
    assert isinstance(excinfo.value, PaperRefusal)
    assert_changed_parts(before, save_to_bytes(dest))  # -- empty budget


def test_argument_validation():
    dest = _open(ALPHA)
    source = _open(BETA)
    with pytest.raises(ValueError, match="mode"):
        dest.import_slide(source, 0, mode="magic")
    with pytest.raises(ValueError, match="same presentation"):
        dest.import_slide(dest, 0, mode="bake")
    with pytest.raises(ValueError, match="out of range"):
        dest.import_slide(source, 99, mode="bake")
    with pytest.raises(ValueError, match="position"):
        dest.import_slide(source, 0, mode="bake", position=99)
    with pytest.raises(ValueError, match="does not belong"):
        dest.import_slide(source, dest.slides[0], mode="bake")
    with pytest.raises(ValueError, match="target_layout does not apply"):
        dest.import_slide(
            source, 0, mode="keep_appearance", target_layout=dest.slide_layouts[0]
        )
    with pytest.raises(ValueError, match="destination"):
        dest.import_slide(
            source, 0, mode="adopt_theme", target_layout=source.slide_layouts[0]
        )


def test_alternate_content_refuses_for_reconciling_modes_only():
    from lxml import etree

    dest = _open(ALPHA)
    source = _open(BETA)
    spTree = source.slides[0].shapes._spTree
    etree.SubElement(
        spTree,
        "{http://schemas.openxmlformats.org/markup-compatibility/2006}AlternateContent",
    )
    with pytest.raises(UnsupportedStructureError, match="AlternateContent"):
        dest.import_slide(source, 0, mode="adopt_theme")
    with pytest.raises(UnsupportedStructureError, match="AlternateContent"):
        dest.import_slide(source, 0, mode="bake")
    report = dest.import_slide(source, 0, mode="keep_appearance")  # -- opaque: allowed
    assert report.mode == "keep_appearance"
    _assert_clean(save_to_bytes(dest))


# ---------------------------------------------------------------- reports and determinism


def test_import_report_matches_frozen_golden():
    """Deterministic report, byte-identical to the reviewed golden."""
    dest = _open(ALPHA)
    report = dest.import_slide(_open(BETA), 0, mode="keep_appearance")
    actual = (
        json.dumps(report.to_dict(), indent=2, ensure_ascii=False) + "\n"
    ).encode("utf-8")
    golden_path = corpus.FIXTURES_DIR.parent / "goldens" / "import_beta_keep.import.json"
    assert actual == golden_path.read_bytes()


def test_import_report_is_deterministic_across_runs():
    first = _open(ALPHA).import_slide(_open(BETA), 1, mode="bake").to_dict()
    second = _open(ALPHA).import_slide(_open(BETA), 1, mode="bake").to_dict()
    assert first == second


def test_import_from_libreoffice_authored_source():
    """Producer diversity: the source deck's final bytes were written by LibreOffice."""
    dest = _open(BETA)
    source = _open(LO_ALPHA)
    report = dest.import_slide(source, 1, mode="keep_appearance")
    assert report.run_shifts == ()
    saved = save_to_bytes(dest)
    _assert_clean(saved)
    reopened = Presentation(io.BytesIO(saved))
    assert len(reopened.slide_masters) == 2


def test_import_delete_scrub_reimport_never_duplicates_partnames():
    """Regression: the fingerprint-dedupe cache must not
    resurrect parts that scrub removed - a ghost hit re-relates a part whose freed
    partname a later import reallocated, producing duplicate zip members with different
    content. The cycle below must yield a clean, fully-registered package."""
    import warnings
    import zipfile

    dest = Presentation()
    dest.slides.add_slide(dest.slide_layouts[6])
    source_alpha = _open(ALPHA)
    source_beta = _open(BETA)
    dest.import_slide(source_alpha, 0, mode="keep_appearance")
    dest.slides.delete(len(dest.slides) - 1)
    dest.scrub(unused_layouts=True, unused_masters=True)
    dest.import_slide(source_beta, 0, mode="keep_appearance")
    dest.import_slide(source_alpha, 0, mode="keep_appearance")

    buf = io.BytesIO()
    with warnings.catch_warnings():
        warnings.simplefilter("error")  # -- zipfile's 'Duplicate name' warning = failure
        dest.save(buf)
    saved = buf.getvalue()
    names = zipfile.ZipFile(io.BytesIO(saved)).namelist()
    assert len(names) == len(set(names))
    partnames = [str(p.partname) for p in dest.part.package.iter_parts()]
    assert len(partnames) == len(set(partnames))
    _assert_clean(saved)
    reopened = Presentation(io.BytesIO(saved))
    reachable_masters = {
        str(p.partname)
        for p in reopened.part.package.iter_parts()
        if "slideMasters/slideMaster" in str(p.partname) and str(p.partname).endswith(".xml")
    }
    assert len(reachable_masters) == len(reopened.slide_masters)


def test_import_placeholder_picture_slide_under_reconciling_modes():
    """Regression: a placeholder PICTURE has no text frame; bake and
    adopt_theme must not crash on it."""
    source = Presentation()
    layout = source.slide_layouts[8]  # -- "Picture with Caption"
    slide = source.slides.add_slide(layout)
    slide.shapes.title.text_frame.paragraphs[0].add_run().text = "Pic ph source"
    picture_ph = next(
        ph for ph in slide.placeholders if ph.placeholder_format.type.name == "PICTURE"
    )
    from PIL import Image as PILImage

    png = io.BytesIO()
    PILImage.new("RGB", (16, 16), (5, 50, 100)).save(png, format="PNG")
    picture_ph.insert_picture(io.BytesIO(png.getvalue()))

    from pptx.shapes.picture import Picture

    for mode in ("bake", "adopt_theme"):
        dest = _open(ALPHA)
        report = dest.import_slide(source, 0, mode=mode)
        assert report.mode == mode
        saved = save_to_bytes(dest)
        _assert_clean(saved)
        reopened = Presentation(io.BytesIO(saved))
        blobs = [
            shape.image.blob
            for shape in reopened.slides[3].shapes
            if isinstance(shape, Picture) or (
                shape.is_placeholder and hasattr(shape, "image")
            )
        ]
        assert any(blob == png.getvalue() for blob in blobs)


def test_append_deck_corrupt_source_refuses_typed():
    dest = _open(ALPHA)
    corrupt = Presentation(
        str(corpus.fixture_path("self_generated/corrupt_dangling_sldid.pptx"))
    )
    before = save_to_bytes(dest)
    with pytest.raises(UnsupportedStructureError, match="relationship graph is broken"):
        dest.append_deck(corrupt, mode="bake")
    assert_changed_parts(before, save_to_bytes(dest))  # -- empty budget


def test_notes_import_enrolls_destination_notes_master():
    """Regression: a destination without a notes master gets one created
    on notes import; it must be enrolled in p:notesMasterIdLst, not just related."""
    dest = _open(ALPHA)  # -- alpha has no notes, hence no notes master
    assert dest._element.notesMasterIdLst is None
    source = _open(GAUNTLET)
    dest.import_slide(source, 1, mode="adopt_theme", notes=True)
    reopened = save_reopen(dest)
    notesMasterIdLst = reopened._element.notesMasterIdLst
    assert notesMasterIdLst is not None
    entry = notesMasterIdLst.notesMasterId
    assert entry is not None
    target = reopened.part.related_part(entry.rId)
    assert "notesMaster" in str(target.partname)


def test_scrub_removes_transplanted_master_once_its_slides_go():
    """Cross-feature: keep_appearance adds a second master; deleting the imported slide
    leaves that master unused, and scrub(unused_masters=True) removes the whole chain."""
    dest = _open(ALPHA)
    dest.import_slide(_open(BETA), 0, mode="keep_appearance")
    dest.slides.delete(3)
    report = dest.scrub(unused_masters=True)
    assert report.unused_masters_removed == ("/ppt/slideMasters/slideMaster2.xml",)
    saved = save_to_bytes(dest)
    _assert_clean(saved)
    zip_map = zip_member_map(saved)
    assert "ppt/theme/theme2.xml" not in zip_map  # -- the chain went with it
    assert "ppt/slideLayouts/slideLayout12.xml" not in zip_map


# --------------------------------------------------------------------------------- lo_smoke


@pytest.mark.lo_smoke
@pytest.mark.parametrize("mode", ["adopt_theme", "keep_appearance", "bake"])
def test_imported_deck_loads_in_libreoffice(mode, tmp_path):
    dest = _open(ALPHA)
    source = _open(BETA)
    for index in range(2):
        dest.import_slide(source, index, mode=mode)
    out = tmp_path / ("import_%s.pptx" % mode)
    dest.save(str(out))
    lo_load_smoke(out, tmp_path)
