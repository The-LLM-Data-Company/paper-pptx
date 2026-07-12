"""Focused text-inspection regressions for promised visible content."""

from __future__ import annotations

from pptx import Presentation
from pptx.inspect import inspect_text

from . import corpus


def _open(relpath):
    return Presentation(str(corpus.fixture_path(relpath)))


def test_chart_frame_is_reported_as_a_typed_blind_region():
    slide = _open("self_generated/chart_notes.pptx").slides[0]
    chart_shape = next(shape for shape in slide.shapes if shape.has_chart)

    marker = next(
        block for block in inspect_text(slide).blocks if block.container == "chart"
    )

    assert marker.shape_id == chart_shape.shape_id
    assert marker.shape_name == chart_shape.name
    assert marker.blind is True
    assert marker.text == ""
    assert marker.runs == ()


def test_fields_are_reported_in_document_position_without_cached_display_text():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paragraph = slide.shapes.add_textbox(
        0, 0, 914400, 914400
    ).text_frame.paragraphs[0]
    paragraph.add_run().text = "Page "
    paragraph.add_slide_number_field()

    block = inspect_text(slide).blocks[0]

    assert block.text == "Page "
    assert block.fields == ("slidenum",)
    assert [run.field_type for run in block.runs] == [None, "slidenum"]
    assert block.to_dict()["runs"][1]["field_type"] == "slidenum"
