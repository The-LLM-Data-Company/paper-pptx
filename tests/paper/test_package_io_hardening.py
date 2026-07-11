"""Adversarial coverage for bounded package reads and atomic ordinary saves."""

from __future__ import annotations

import io
import os
import stat
import warnings
import zipfile

import pytest
from lxml import etree

from pptx import Presentation
from pptx.errors import PackageLimitError, PaperRefusal
from pptx.opc import serialized

from . import corpus


def _minimal_path():
    return corpus.fixture_path("self_generated/minimal_clean.pptx")


def test_normal_open_refuses_duplicate_members(tmp_path):
    target = tmp_path / "duplicate.pptx"
    source = _minimal_path()
    with zipfile.ZipFile(source) as incoming, zipfile.ZipFile(target, "w") as outgoing:
        for info in incoming.infolist():
            outgoing.writestr(info, incoming.read(info.filename))
        slide = incoming.read("ppt/slides/slide1.xml")
        with warnings.catch_warnings():
            warnings.simplefilter("ignore", UserWarning)
            outgoing.writestr("ppt/slides/slide1.xml", slide)

    assert issubclass(PackageLimitError, PaperRefusal)
    with pytest.raises(PackageLimitError, match="duplicate member"):
        Presentation(target)


def test_normal_open_refuses_noncanonical_member_names(tmp_path):
    target = tmp_path / "noncanonical.pptx"
    with zipfile.ZipFile(target, "w") as archive:
        archive.writestr("../ppt/presentation.xml", b"<presentation/>")

    with pytest.raises(PackageLimitError, match="noncanonical"):
        Presentation(target)


def test_saved_repetitive_deck_reopens(tmp_path):
    """Regression: the save -> reopen covenant must hold for this package's OWN output.

    Machine-generated decks (thousands of near-identical paragraphs) legitimately
    exceed any expanded-to-compressed ratio a zip bomb would need. A ratio guard once
    refused such files at reopen; the absolute member and package byte limits are the
    safety envelope instead.
    """
    from pptx.util import Inches

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    frame = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(4)).text_frame
    for index in range(12_000):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = "Quarterly revenue synergy alignment placeholder row"
    target = tmp_path / "repetitive.pptx"
    presentation.save(target)

    with zipfile.ZipFile(target) as archive:
        worst_ratio = max(
            info.file_size / info.compress_size
            for info in archive.infolist()
            if info.compress_size
        )
    assert worst_ratio > 100, "fixture lost its bite; deck no longer compresses past 100:1"

    reopened = Presentation(target)
    assert len(reopened.slides) == 1


def _rewrite_package(source, target, transform):
    with zipfile.ZipFile(source) as incoming, zipfile.ZipFile(target, "w") as outgoing:
        for info in incoming.infolist():
            replacement = transform(info.filename, incoming.read(info.filename))
            if replacement is not None:
                outgoing.writestr(info, replacement)


def test_normal_open_refuses_a_missing_relationship_target(tmp_path):
    target = tmp_path / "missing-target.pptx"

    def transform(name, blob):
        if name == "_rels/.rels":
            return blob.replace(b"ppt/presentation.xml", b"ppt/missing-part.xml")
        return blob

    _rewrite_package(_minimal_path(), target, transform)

    with pytest.raises(PackageLimitError, match="targets missing part"):
        Presentation(target)


def test_normal_open_refuses_an_unreachable_part(tmp_path):
    target = tmp_path / "unreachable.pptx"
    _rewrite_package(_minimal_path(), target, lambda _name, blob: blob)
    with zipfile.ZipFile(target, "a") as archive:
        archive.writestr("ppt/orphan.bin", b"would be dropped on save")

    with pytest.raises(PackageLimitError, match="unreachable parts"):
        Presentation(target)


def test_normal_open_refuses_duplicate_relationship_ids(tmp_path):
    target = tmp_path / "duplicate-rid.pptx"

    def transform(name, blob):
        if name != "_rels/.rels":
            return blob
        root = etree.fromstring(blob)
        root.append(etree.fromstring(etree.tostring(root[0])))
        return etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)

    _rewrite_package(_minimal_path(), target, transform)

    with pytest.raises(PackageLimitError, match="duplicate ids"):
        Presentation(target)


def test_path_save_failure_preserves_existing_file_and_mode(tmp_path, monkeypatch):
    presentation = Presentation(_minimal_path())
    destination = tmp_path / "existing.pptx"
    destination.write_bytes(b"known-good destination")
    destination.chmod(0o640)
    original_write = serialized._ZipPkgWriter.write
    writes = 0

    def fail_during_write(self, pack_uri, blob):
        nonlocal writes
        writes += 1
        if writes == 2:
            raise OSError("forced ZIP write failure")
        return original_write(self, pack_uri, blob)

    monkeypatch.setattr(serialized._ZipPkgWriter, "write", fail_during_write)
    with pytest.raises(OSError, match="forced ZIP write failure"):
        presentation.save(destination)

    assert destination.read_bytes() == b"known-good destination"
    assert stat.S_IMODE(destination.stat().st_mode) == 0o640
    assert not list(tmp_path.glob(".existing.pptx.*.partial"))


def test_successful_path_save_preserves_existing_mode(tmp_path):
    presentation = Presentation(_minimal_path())
    destination = tmp_path / "existing.pptx"
    destination.write_bytes(b"old")
    destination.chmod(0o604)

    presentation.save(destination)

    assert stat.S_IMODE(destination.stat().st_mode) == 0o604
    assert len(Presentation(io.BytesIO(destination.read_bytes())).slides) == 1


def test_new_path_save_honors_umask(tmp_path):
    """Regression: a save to a NEW path must not keep mkstemp's private 0600 mode."""
    presentation = Presentation(_minimal_path())
    destination = tmp_path / "brand-new.pptx"
    previous_umask = os.umask(0o027)
    try:
        presentation.save(destination)
    finally:
        os.umask(previous_umask)

    assert stat.S_IMODE(destination.stat().st_mode) == 0o666 & ~0o027
    assert len(Presentation(io.BytesIO(destination.read_bytes())).slides) == 1


def test_patch_save_honors_umask_and_existing_mode(tmp_path):
    """`patch_save` output modes follow the same contract as ordinary path saves."""
    from pptx.package import patch_save

    source = _minimal_path()
    destination = tmp_path / "patched.pptx"
    previous_umask = os.umask(0o027)
    try:
        patch_save(str(source), Presentation(source), str(destination))
    finally:
        os.umask(previous_umask)
    assert stat.S_IMODE(destination.stat().st_mode) == 0o666 & ~0o027

    destination.chmod(0o604)
    patch_save(str(source), Presentation(source), str(destination))
    assert stat.S_IMODE(destination.stat().st_mode) == 0o604
