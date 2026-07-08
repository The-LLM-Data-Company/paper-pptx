"""PR-0 stub tests (CONVENTIONS §8): one strict xfail per unimplemented organ.

Each stub imports the names API-PROPOSAL.md pins for its organ. `strict=True` means the stub
FAILS the suite the moment the organ lands (XPASS), forcing the landing phase to replace its
stub with real contract tests in the same change — and keeping this file an accurate ledger of
what remains unimplemented.
"""

from __future__ import annotations

import pytest

from pptx import Presentation


@pytest.mark.xfail(strict=True, reason="PR-0 stub - lands with Phase 2 (bullets + errors)")
def test_pr0_errors_and_bullets_api():
    from pptx.enum.text import PP_BULLET_TYPE  # noqa: F401
    from pptx.errors import (  # noqa: F401
        AmbiguousTargetError,
        BoundaryViolationError,
        PaperRefusal,
        RelationshipPolicyError,
        TargetNotFoundError,
        UnsupportedStructureError,
    )
    from pptx.text.bullet import BulletFormat

    slide = Presentation().slides.add_slide(Presentation().slide_layouts[6])
    paragraph = slide.shapes.add_textbox(0, 0, 100, 100).text_frame.paragraphs[0]
    assert isinstance(paragraph.bullet, BulletFormat)
    for name in ("type", "char", "number_scheme", "start_at", "font_name", "size_percent"):
        assert hasattr(BulletFormat, name)
    for name in ("set_character", "set_numbered", "set_none"):
        assert callable(getattr(BulletFormat, name))


@pytest.mark.xfail(strict=True, reason="PR-0 stub - lands with Phase 3 (autofit)")
def test_pr0_autofit_api():
    from pptx.text.text import TextFrame

    for name in ("font_scale", "line_space_reduction"):
        assert isinstance(getattr(TextFrame, name), property)
    assert callable(TextFrame.normalize_autofit)


@pytest.mark.xfail(strict=True, reason="PR-0 stub - lands with Phase 4 (effective inspection)")
def test_pr0_inspect_api():
    from pptx.inspect import (  # noqa: F401
        BlockAnchor,
        EffectiveFont,
        EffectiveValue,
        ProvenanceStep,
        TextInspection,
        effective_font,
        inspect_text,
    )
    from pptx.text.text import _Run

    assert callable(_Run.effective_font)


@pytest.mark.xfail(strict=True, reason="PR-0 stub - lands with Phase 5 (package kernel)")
def test_pr0_package_kernel_api():
    from pptx.package import (  # noqa: F401
        PackageDiff,
        PartDelta,
        diff_package,
        patch_save,
        xml_equivalent,
    )


@pytest.mark.xfail(strict=True, reason="PR-0 stub - lands with Phase 6 (speaker notes)")
def test_pr0_notes_api():
    from pptx.slide import Slide

    assert callable(Slide.read_notes_text)
    assert callable(Slide.replace_notes_text)


@pytest.mark.xfail(strict=True, reason="PR-0 stub - lands with Phase 7 (slide operations)")
def test_pr0_slide_ops_api():
    from pptx.slide import SlideClonePolicy, Slides  # noqa: F401

    for name in ("clone", "delete", "reorder", "move"):
        assert callable(getattr(Slides, name))


@pytest.mark.xfail(strict=True, reason="PR-0 stub - lands with Phase 8 (image replacement)")
def test_pr0_image_replacement_api():
    from pptx.shapes.picture import Picture

    assert callable(Picture.replace_image)


@pytest.mark.xfail(strict=True, reason="PR-0 stub - lands with Phase 9 (chart data routing)")
def test_pr0_chart_routing_api():
    from pptx.chart.chart import Chart
    from pptx.shapes.shapetree import SlideShapes

    assert callable(SlideShapes.chart_by_name)
    assert callable(Chart.replace_data_safe)
