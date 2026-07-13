"""The QBR walkthrough eval: the canonical template job as a permanent test.

This module makes a real-job simulation permanent: `test_walkthrough_end_to_end` executes the
canonical template job — build a QBR deck from the gauntlet "corporate template" — using only
shipped public API, asserting through the contract harness at every step. Capabilities not
shipped yet are strict-xfail step tests: the
suite FAILS the moment an organ lands without its walkthrough step flipping, and the suite is not
done while any such step here is an xfail.
"""

from __future__ import annotations

import io

import pytest

from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.package import patch_save
from pptx.util import Pt

from . import corpus
from .contract import zip_member_map
from .idlists import dangling_section_slide_ids, duplicate_section_slide_ids
from .lo import lo_load_smoke
from .relint import dangling_relationship_targets, missing_relationship_references

GAUNTLET = "self_generated/gauntlet.pptx"


def _build_qbr_deck(tmp_path):
    """Run the canonical template job with shipped API; return (deck_path, presentation)."""
    template_path = tmp_path / "template.pptx"
    template_path.write_bytes(corpus.fixture_path(GAUNTLET).read_bytes())
    prs = Presentation(str(template_path))

    # -- structural pass first (template-editing doctrine): 3 content sections from the
    # -- branded slide, chart slide kept, autofit playground dropped, closing slide last
    for _ in range(2):
        prs.slides.clone(0)
    prs.slides.delete(prs.slides.index(next(s for s in prs.slides if any(
        shape.name == "autofit_normal_box" for shape in s.shapes
    ))))
    assert len(prs.slides) == 5  # -- 3 branded + chart + closing

    # -- content pass: anchored retitles — each branded clone gets its own title
    # -- via a hash-checked anchor; then one deck-wide format-preserving language pass
    from pptx.edit import replace_text, replace_text_at
    from pptx.inspect import inspect_text

    section_titles = ["Q4 Business Review", "Wins", "Risks"]
    branded = [s for s in prs.slides if s.shapes.title is not None
               and s.shapes.title.text == "Gauntlet: branded"]
    assert len(branded) == 3
    for slide, title in zip(branded, section_titles):
        anchor = next(
            b.anchor for b in inspect_text(slide).blocks if b.text == "Gauntlet: branded"
        )
        result = replace_text_at(prs, anchor, "Gauntlet: branded", title)
        assert result.replacements == 1
    language_pass = replace_text(prs, "Inherited body level", "Section content line")
    assert language_pass.replacements == 6  # -- two lines on each of the three sections
    closing_title = prs.slides[-1].shapes.title
    closing_title.text_frame.paragraphs[0].runs[0].text = "Next steps"
    # -- capture an anchor now; later structural passes will shift block indices and the
    # -- refind recovery path brings it back
    closing_title_anchor = next(
        b.anchor for b in inspect_text(prs.slides[-1]).blocks if b.text == "Next steps"
    )

    # -- chart data update, by name, validated
    chart_slide = next(s for s in prs.slides if any(sh.has_chart for sh in s.shapes))
    chart = chart_slide.shapes.chart_by_name("gauntlet_chart")
    chart.replace_data_safe(["Q3", "Q4"], [("Revenue", (14.2, 16.8))], number_format="0.0")

    # -- speaker notes for the talk track
    chart_slide.replace_notes_text("Walk through Q4 revenue; flag the Alpha renewal.")

    # -- declutter the chart slide: drop its decorative picture, rel-safely
    chart_slide.shapes.delete(chart_slide.shapes.picture_by_name("gauntlet_img_1"))

    # -- image swap by name, across formats (the brand asset arrives
    # -- as a JPEG), geometry preserved
    from PIL import Image as PILImage

    replacement = io.BytesIO()
    PILImage.new("RGB", (64, 64), (10, 90, 160)).save(replacement, format="JPEG")
    closing = prs.slides[-1]
    picture = closing.shapes.picture_by_name("gauntlet_img_2")
    picture.replace_image(io.BytesIO(replacement.getvalue()), allow_format_change=True)
    # -- reuse the swapped brand asset on the title slide (add_copy: image shared)
    brand_mark = prs.slides[0].shapes.add_copy(picture)
    title_slide_ids = [s.shape_id for s in prs.slides[0].shapes]
    assert title_slide_ids.count(brand_mark.shape_id) == 1  # -- fresh id on ITS slide
    # -- and push the link box behind everything on the closing slide (z-order)
    closing.shapes.move(closing.shapes[-1], 0)

    # -- normalize the autofit frame we kept prose in (freeze what the reader sees). The
    # -- box lives on the closing slide; no None-guard — a missing box must fail loudly.
    body_box = next(
        s for slide in prs.slides for s in slide.shapes if s.name == "hyperlink_box"
    )
    tf = body_box.text_frame
    for paragraph in tf.paragraphs:
        paragraph.line_spacing = 1.0
    tf.normalize_autofit(resolve=True, min_font_size=Pt(11))
    assert tf.auto_size == MSO_AUTO_SIZE.NONE

    # -- footer pass: the deck-level dialog-equivalent replaces the
    # -- hand-rolled per-slide page-number rail (the static-page-number anti-pattern,
    # -- now retired); numbers are real slidenum fields, live across any later reorder
    from pptx.inspect import effective_paragraph_format

    prs.slide_masters[0].header_footers.slide_number_visible = True
    prs.apply_footers(footer="Q4 QBR", slide_number=True)
    # -- the title slide keeps one caption with a real datetime field (the
    # -- field primitive: fields compose into arbitrary text, not just furniture)
    caption = prs.slides[0].shapes.add_textbox(Pt(20), Pt(500), Pt(200), Pt(20))
    caption.name = "prepared_on_caption"
    caption_paragraph = caption.text_frame.paragraphs[0]
    caption_paragraph.add_run().text = "Prepared "
    caption_paragraph.add_datetime_field("datetime1")
    # -- the applied footer's effective alignment resolves through the layout
    footer_ph = next(
        s
        for s in prs.slides[0].shapes
        if s.is_placeholder and s.element.ph.get("type") == "ftr"
    )
    footer_format = effective_paragraph_format(footer_ph.text_frame.paragraphs[0])
    assert footer_format.alignment.resolved is True

    # -- the early anchor is now stale (z-order + footer passes shifted indices or will);
    # -- the pinned recovery path: refuse -> refind -> retry
    from pptx.edit import refind
    from pptx.errors import StaleAnchorError, TargetNotFoundError

    try:
        replace_text_at(prs, closing_title_anchor, "Next steps", "Next steps & owners")
    except (StaleAnchorError, TargetNotFoundError):
        fresh = refind(prs, closing_title_anchor)
        replace_text_at(prs, fresh, "Next steps", "Next steps & owners")

    # -- exit gate: the deck is about to leave the pipeline. Speaker
    # -- notes are the talk track and stay; metadata, unused layouts, and any
    # -- unreferenced media go. The report's budget is the scrub's own evidence.
    scrub_report = prs.scrub(
        comments=True, metadata=True, unused_layouts=True, unreachable_media=True
    )
    assert scrub_report.notes_slides_removed == ()
    assert len(scrub_report.unused_layouts_removed) == 9
    assert scrub_report.metadata_fields_cleared != ()

    # -- narrow save: only genuinely-changed parts differ from the template
    out_path = tmp_path / "qbr.pptx"
    diff = patch_save(str(template_path), prs, str(out_path))
    assert not diff.is_empty
    return out_path, prs


def _assert_package_integrity(pptx_bytes):
    zip_map = zip_member_map(pptx_bytes)
    assert dangling_relationship_targets(zip_map) == []
    assert missing_relationship_references(zip_map) == []
    assert dangling_section_slide_ids(zip_map) == []
    assert duplicate_section_slide_ids(zip_map) == []


def test_walkthrough_end_to_end(tmp_path):
    out_path, _ = _build_qbr_deck(tmp_path)
    saved = out_path.read_bytes()
    _assert_package_integrity(saved)

    reopened = Presentation(io.BytesIO(saved))
    assert len(reopened.slides) == 5
    titles = [
        s.shapes.title.text if s.shapes.title is not None else None for s in reopened.slides
    ]
    assert titles[:3] == ["Q4 Business Review", "Wins", "Risks"]
    assert titles[-1] == "Next steps & owners"  # -- the refind recovery path landed
    chart = next(
        sh.chart for s in reopened.slides for sh in s.shapes if sh.has_chart
    )
    assert [(s.name, tuple(s.values)) for s in chart.series] == [("Revenue", (14.2, 16.8))]
    chart_slide = next(s for s in reopened.slides if any(sh.has_chart for sh in s.shapes))
    assert chart_slide.read_notes_text() == (
        "Walk through Q4 revenue; flag the Alpha renewal."
    )

    # -- the applied footer trio landed on every slide as REAL fields,
    # -- and the title slide's caption carries its datetime field
    from pptx.inspect import inspect_text as _inspect_text

    for ordinal, slide in enumerate(reopened.slides, start=1):
        blocks = _inspect_text(slide).blocks
        field_tokens = sorted(token for b in blocks for token in b.fields)
        expected = ["datetime1", "slidenum"] if ordinal == 1 else ["slidenum"]
        assert field_tokens == expected, "slide %d fields: %r" % (ordinal, field_tokens)
        assert any(b.text == "Q4 QBR" for b in blocks)

    # -- the scrub held: notes survive, personal metadata does not
    assert reopened.core_properties.author == ""
    assert reopened.core_properties.last_modified_by == ""

    # -- the brand asset swapped formats and was copied to the title slide (shared part)
    closing_pic = reopened.slides[-1].shapes.picture_by_name("gauntlet_img_2")
    assert closing_pic.image.ext in ("jpg", "jpeg")
    title_mark = next(
        s for s in reopened.slides[0].shapes if s.shape_type.name == "PICTURE"
    )
    assert title_mark.image.blob == closing_pic.image.blob

    # -- z-order held: the link box is backmost on the closing slide
    assert reopened.slides[-1].shapes[0].name == "hyperlink_box"

    # -- the normalized frame persisted explicit no-autofit with floored sizes
    normalized = next(
        s for slide in reopened.slides for s in slide.shapes if s.name == "hyperlink_box"
    )
    assert normalized.text_frame.auto_size == MSO_AUTO_SIZE.NONE
    assert reopened.slide_masters[0].header_footers.slide_number_visible is True


def test_walkthrough_self_consistency_against_deck_diff(tmp_path):
    """The self-consistency check: the job's own operation reports and diff_decks(input, output)
    are two independent evidence systems - they must tell the same story."""
    from pptx.diff import diff_decks

    out_path, _ = _build_qbr_deck(tmp_path)
    template_path = tmp_path / "template.pptx"
    report = diff_decks(str(template_path), str(out_path), detail="text")

    # -- structural story: two clones added, the autofit playground removed
    assert len(report.slides_added) == 2
    assert [r.title for r in report.slides_removed] == [None]  # -- the autofit slide
    changes = {change.slide_id: change for change in report.slide_changes}

    # -- the retitle the anchored replace performed
    retitled = changes[256]
    assert any(
        c["before"] == "Gauntlet: branded" and c["after"] == "Q4 Business Review"
        for c in retitled.text_changes
    )
    # -- the chart update replace_data_safe performed, per series/category
    chart_slide_id = 257  # -- gauntlet slide 2
    chart_deltas = changes[chart_slide_id].chart_data_changes
    assert any("categories_after" in d and d["categories_after"] == ["Q3", "Q4"]
               for d in chart_deltas)
    # -- the notes replacement
    assert changes[chart_slide_id].notes_change == {
        "before": "Gauntlet speaker notes.",
        "after": "Walk through Q4 revenue; flag the Alpha renewal.",
    }
    # -- the decorative picture delete on the chart slide
    assert "gauntlet_img_1" in changes[chart_slide_id].shapes_removed
    # -- the closing slide's image swap reads as replacement, not move/resize
    closing_id = 259  # -- gauntlet slide 4
    assert "gauntlet_img_2" in changes[closing_id].images_replaced


@pytest.mark.lo_smoke
def test_walkthrough_output_loads_in_libreoffice(tmp_path):
    out_path, _ = _build_qbr_deck(tmp_path)
    lo_load_smoke(out_path, tmp_path)


# ---------------------------------------------------------------- unshipped steps (strict)
# Each xfail names its capability item and FAILS the suite when the organ lands without the
# walkthrough growing the real step.


def test_step_survey_template_with_deck_manifest():
    """Survey step (was xfail): the job starts with a structural survey, not guesswork."""
    from pptx.inspect import inspect_deck

    prs = Presentation(str(corpus.fixture_path(GAUNTLET)))
    manifest = inspect_deck(prs)
    assert manifest.slide_count == 4
    chart_slides = [
        slide.part for slide in manifest.slides
        if any(shape.chart for shape in slide.shapes)
    ]
    assert chart_slides == ["/ppt/slides/slide2.xml"]
    # -- and the template's table is addressable by name from the survey
    table_slide_index = next(
        index for index, slide in enumerate(manifest.slides)
        if any(shape.table for shape in slide.shapes)
    )
    table = prs.slides[table_slide_index].shapes.table_by_name("gauntlet_table")
    assert table.cell(0, 0).text == "r0c0"


def test_step_check_brand_accent_via_effective_shape_format():
    """Brand-accent step (was xfail): 'is that box actually brand-colored?' is now answerable."""
    from pptx.inspect import effective_shape_format

    prs = Presentation(str(corpus.fixture_path("self_generated/clrmap_remap.pptx")))
    rect = prs.slides[0].shapes.shape_by_name("accent1_box")
    fmt = effective_shape_format(rect)
    assert fmt.fill_rgb.value == "C0504D"  # -- accent1 through the remapped clrMap
    assert fmt.fill_rgb.resolved is True


def test_step_update_libreoffice_authored_chart():
    """Chart-update step (was xfail): the externally-produced deck's chart takes new numbers."""
    prs = Presentation(str(corpus.fixture_path("libreoffice_export/lo_chart_notes.pptx")))
    chart = next(sh.chart for sh in prs.slides[0].shapes if sh.has_chart)
    chart.replace_data_safe(["A", "B"], [("S1", (1.0, 2.0))])
    reopened = Presentation(io.BytesIO(_bytes_of(prs)))
    reopened_chart = next(sh.chart for sh in reopened.slides[0].shapes if sh.has_chart)
    assert [(s.name, tuple(s.values)) for s in reopened_chart.series] == [("S1", (1.0, 2.0))]


def _bytes_of(prs):
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


