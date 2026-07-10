"""Contract tests: anchored, formatting-preserving text replacement (pptx.edit).

Includes the invariant this organ finally gives a subject to:
replace(x→y) then replace(y→x) restores visible text AND formatting.
"""

from __future__ import annotations

import io

import pytest
from lxml import etree

from pptx import Presentation
from pptx.edit import ReplaceResult, refind, replace_text, replace_text_at
from pptx.errors import (
    AmbiguousTargetError,
    StaleAnchorError,
    TargetNotFoundError,
    UnsupportedStructureError,
)
from pptx.inspect import BlockAnchor, inspect_text
from pptx.util import Emu

from . import corpus
from .contract import assert_changed_parts, save_to_bytes, snapshot_parts
from .lo import lo_load_smoke

MINIMAL = "self_generated/minimal_clean.pptx"
BRANDED = "self_generated/branded_template.pptx"
_A = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _open(relpath):
    return Presentation(str(corpus.fixture_path(relpath)))


def _reopen(pptx_bytes):
    return Presentation(io.BytesIO(pptx_bytes))


def _multi_run_paragraph(prs, texts_and_bold):
    """Add a textbox whose first paragraph has one run per (text, bold) pair."""
    box = prs.slides[0].shapes.add_textbox(0, 0, Emu(914400 * 4), Emu(914400))
    paragraph = box.text_frame.paragraphs[0]
    for text, bold in texts_and_bold:
        run = paragraph.add_run()
        run.text = text
        run.font.bold = bold
    return box


def _char_formatting_fingerprint(paragraph):
    """(char, serialized-rPr-or-None) per character — formatting equality at char level."""
    fingerprint = []
    for run in paragraph.runs:
        rPr = run._r.find("{%s}rPr" % _A)
        rPr_bytes = etree.tostring(rPr) if rPr is not None else None
        for char in run.text:
            fingerprint.append((char, rPr_bytes))
    return fingerprint


# ------------------------------------------------------------------------ deck-wide replace


def test_replace_text_round_trips_with_exact_budget():
    prs = _open(MINIMAL)
    before = save_to_bytes(prs)
    result = replace_text(prs, "Minimal clean", "Renamed")
    assert result.replacements == 1
    after = save_to_bytes(prs)
    assert_changed_parts(before, after, expect_changed=["ppt/slides/slide1.xml"])
    assert _reopen(after).slides[0].shapes.title.text == "Renamed fixture"


def test_replacement_inherits_the_match_start_runs_formatting():
    prs = _open(MINIMAL)
    box = _multi_run_paragraph(prs, [("The ", None), ("BIG", True), (" deal", None)])
    replace_text(prs, "BIG", "HUGE")

    paragraph = box.text_frame.paragraphs[0]
    assert [r.text for r in paragraph.runs] == ["The ", "HUGE", " deal"]
    assert [r.font.bold for r in paragraph.runs] == [None, True, None]


def test_untouched_runs_stay_byte_identical():
    prs = _open(MINIMAL)
    box = _multi_run_paragraph(prs, [("alpha ", None), ("beta", True), (" gamma", None)])
    paragraph = box.text_frame.paragraphs[0]
    first_before = etree.tostring(paragraph.runs[0]._r)
    last_before = etree.tostring(paragraph.runs[2]._r)
    replace_text(prs, "beta", "delta")
    assert etree.tostring(paragraph.runs[0]._r) == first_before
    assert etree.tostring(paragraph.runs[2]._r) == last_before


def test_cross_run_match_inherits_start_run_and_keeps_boundary_fragments():
    prs = _open(MINIMAL)
    # -- "Hello world" split as "Hell" + "o world"; the match starts INSIDE run 0
    box = _multi_run_paragraph(prs, [("Hell", True), ("o world", None)])
    result = replace_text(prs, "lo wo", "LO WO")
    assert result.replacements == 1
    paragraph = box.text_frame.paragraphs[0]
    # -- replacement joined run 0 (where the match starts); run 1 keeps its suffix
    assert [r.text for r in paragraph.runs] == ["HelLO WO", "rld"]
    assert [r.font.bold for r in paragraph.runs] == [True, None]


def test_match_starting_at_a_run_boundary_belongs_to_the_later_run():
    prs = _open(MINIMAL)
    box = _multi_run_paragraph(prs, [("Hel", True), ("lo world", None)])
    replace_text(prs, "lo wo", "LO WO")  # -- starts at char 3 = run 1's first char
    paragraph = box.text_frame.paragraphs[0]
    assert [r.text for r in paragraph.runs] == ["Hel", "LO WOrld"]
    assert [r.font.bold for r in paragraph.runs] == [True, None]


def test_run_consumed_whole_is_removed():
    prs = _open(MINIMAL)
    box = _multi_run_paragraph(prs, [("a", None), ("b", True), ("c", None)])
    replace_text(prs, "b", "")
    paragraph = box.text_frame.paragraphs[0]
    assert [r.text for r in paragraph.runs] == ["a", "c"]


def test_trailing_whitespace_is_content_in_both_directions():
    """The kernel's whitespace-trap doctrine applies to replacement too."""
    prs = Presentation(
        str(corpus.fixture_path("self_generated/whitespace_trailing_a.pptx"))
    )
    result = replace_text(prs, "Trailing space ", "Kept tail ")
    assert result.replacements == 1
    box = next(s for s in prs.slides[0].shapes if s.name == "whitespace_box")
    assert box.text_frame.text == "Kept tail "  # -- trailing space preserved


def test_matches_do_not_cross_line_breaks():
    prs = _open(MINIMAL)
    box = _multi_run_paragraph(prs, [("alpha", None)])
    paragraph = box.text_frame.paragraphs[0]
    paragraph.add_line_break()
    run = paragraph.add_run()
    run.text = "beta"
    assert replace_text(prs, "alphabeta", "x").replacements == 0
    assert replace_text(prs, "alpha", "ALPHA").replacements == 1


def test_replace_reaches_table_cells_and_grouped_shapes():
    prs = _open("self_generated/tables_in_group.pptx")
    before = save_to_bytes(prs)
    result = replace_text(prs, "cell r0c0", "REPLACED CELL")
    result2 = replace_text(prs, "In-group text", "IN GROUP REPLACED")
    assert result.replacements == 1
    assert result2.replacements == 1
    after = save_to_bytes(prs)
    assert_changed_parts(before, after, expect_changed=["ppt/slides/slide1.xml"])
    texts = [b.text for b in inspect_text(_reopen(after).slides[0]).blocks]
    assert "REPLACED CELL" in texts
    assert "IN GROUP REPLACED" in texts


def test_notes_replaced_only_when_asked():
    CHART_NOTES = "self_generated/chart_notes.pptx"
    prs = _open(CHART_NOTES)
    assert replace_text(prs, "clone fixture", "QBR").replacements == 0
    assert prs.slides[0].read_notes_text() == "Speaker notes for the clone fixture."

    prs = _open(CHART_NOTES)
    result = replace_text(prs, "clone fixture", "QBR", include_notes=True)
    assert result.replacements == 1
    assert prs.slides[0].read_notes_text() == "Speaker notes for the QBR."
    assert result.blocks[0].part.startswith("/ppt/notesSlides/")


def test_zero_matches_is_a_normal_result_with_empty_budget():
    prs = _open(MINIMAL)
    before = save_to_bytes(prs)
    result = replace_text(prs, "no such text anywhere", "x")
    assert result == ReplaceResult(0, ())
    assert_changed_parts(before, save_to_bytes(prs))


def test_result_payload_carries_pinned_schema():
    prs = _open(MINIMAL)
    payload = replace_text(prs, "Minimal", "Basic").to_dict()
    assert payload["schema"] == "paper-replace-result"
    assert payload["version"] == 1
    assert payload["replacements"] == 1
    assert payload["blocks"][0]["part"] == "/ppt/slides/slide1.xml"


@pytest.mark.parametrize(
    "bad_call",
    [
        lambda prs: replace_text(prs, "", "x"),
        lambda prs: replace_text(prs, 42, "x"),
        lambda prs: replace_text(prs, "x", 42),
        lambda prs: replace_text(prs, "x", "multi\nline"),
        lambda prs: replace_text(prs, "with\nbreak", "x"),
        lambda prs: replace_text(prs, "x", "bad\udc80"),
        lambda prs: replace_text(prs, "x", "bad\x01ctrl"),
        lambda prs: replace_text(prs, "bad\x00find", "x"),
        lambda prs: replace_text_at(prs, "not-an-anchor", "x", "y"),
    ],
)
def test_bad_arguments_raise_valueerror_and_touch_nothing(bad_call):
    prs = _open(MINIMAL)
    before = snapshot_parts(prs)
    with pytest.raises(ValueError):
        bad_call(prs)
    assert snapshot_parts(prs) == before


def test_overlapping_occurrences_are_matched_non_overlapping_left_to_right():
    """Pinned occurrence semantics (was only a docstring): 'aaa' has ONE match of 'aa'."""
    prs = _open(MINIMAL)
    box = _multi_run_paragraph(prs, [("aaa", None)])
    result = replace_text(prs, "aa", "X")
    assert result.replacements == 1
    assert box.text_frame.paragraphs[0].text == "Xa"


def test_multiple_matches_in_one_paragraph_are_all_replaced_and_counted():
    prs = _open(MINIMAL)
    box = _multi_run_paragraph(prs, [("banana", None)])
    result = replace_text(prs, "na", "NA")  # -- "na" occurs nowhere else in the fixture
    assert result.replacements == 2
    assert box.text_frame.paragraphs[0].text == "baNANA"


def _add_deep_group(slide, depth=17):
    group = slide.shapes.add_group_shape()
    for _ in range(depth):
        group = group.shapes.add_group_shape()
    box = group.shapes.add_textbox(0, 0, 914400, 914400)
    box.text_frame.paragraphs[0].add_run().text = "unreachable target"


def test_traversal_refusal_fires_before_any_write():
    """Regression (review, two dimensions): the depth guard on a LATER slide used to fire
    after earlier blocks were already rewritten — a refusal must leave zero edits behind."""
    prs = Presentation(str(corpus.fixture_path("self_generated/gauntlet.pptx")))
    _add_deep_group(prs.slides[3])  # -- refusal source is on the LAST slide

    def operation(p):
        replace_text(p, "Gauntlet", "REWRITTEN")  # -- would match on earlier slides

    from .contract import assert_refusal_atomic

    raised = assert_refusal_atomic(prs, operation, UnsupportedStructureError)
    assert "nested" in str(raised)


_MC = "http://schemas.openxmlformats.org/markup-compatibility/2006"
_P = "http://schemas.openxmlformats.org/presentationml/2006/main"


def _wrap_first_textbox_in_alternate_content(prs):
    """Wrap a new textbox's p:sp in mc:AlternateContent (Choice + Fallback branches)."""
    import copy as copy_module

    slide = prs.slides[0]
    box = slide.shapes.add_textbox(0, 0, 914400, 914400)
    box.text_frame.paragraphs[0].add_run().text = "mc content"
    sp = box._element
    spTree = sp.getparent()
    ac = spTree.makeelement("{%s}AlternateContent" % _MC, {})
    choice = ac.makeelement("{%s}Choice" % _MC, {"Requires": "p14"})
    fallback = ac.makeelement("{%s}Fallback" % _MC, {})
    ac.append(choice)
    ac.append(fallback)
    choice.append(copy_module.deepcopy(sp))
    sp.addprevious(ac)
    spTree.remove(sp)
    fallback.append(sp)
    return prs


def test_replace_text_refuses_alternate_content_decks_atomically():
    """Regression: mc:AlternateContent used to be silently skipped; 'replace every
    occurrence' cannot be honored over content the library cannot see into."""
    prs = _wrap_first_textbox_in_alternate_content(_open(MINIMAL))

    from .contract import assert_refusal_atomic

    raised = assert_refusal_atomic(
        prs, lambda p: replace_text(p, "Minimal", "X"), UnsupportedStructureError
    )
    assert "AlternateContent" in str(raised)


def test_alternate_content_occupies_one_block_index_consistently():
    """inspect_text reports the AC subtree as ONE blind block; anchors before/after it stay
    editable and index-aligned between inspect and edit."""
    prs = _wrap_first_textbox_in_alternate_content(_open(MINIMAL))
    inspection = inspect_text(prs.slides[0])
    ac_blocks = [b for b in inspection.blocks if b.container == "alternate-content"]
    assert len(ac_blocks) == 1
    assert ac_blocks[0].blind is True
    assert inspection.blind_region_count == 1

    # -- a normal block found by inspect is editable at the same index via edit
    title_block = next(b for b in inspection.blocks if b.text == "Minimal clean fixture")
    result = replace_text_at(prs, title_block.anchor, "Minimal clean", "Aligned")
    assert result.replacements == 1

    # -- anchoring the AC block itself refuses, typed
    with pytest.raises(UnsupportedStructureError, match="AlternateContent"):
        replace_text_at(prs, ac_blocks[0].anchor, "x", "y")


def test_cross_boundary_occurrence_in_anchored_block_refuses():
    """Regression: `find` present in the hash-text but only across a field boundary used to
    return a success-shaped ReplaceResult(0) that listed the block as touched."""
    prs = _open(MINIMAL)
    box = prs.slides[0].shapes.add_textbox(0, 0, Emu(914400 * 3), Emu(914400))
    box.name = "field_box"
    paragraph = box.text_frame.paragraphs[0]
    paragraph.add_run().text = "Page "
    paragraph.add_slide_number_field()
    run = paragraph.add_run()
    run.text = " of 10"

    anchor = next(
        b.anchor for b in inspect_text(prs.slides[0]).blocks if b.shape_name == "field_box"
    )
    from .contract import snapshot_parts

    before = snapshot_parts(prs)
    with pytest.raises(TargetNotFoundError, match="boundary"):
        replace_text_at(prs, anchor, "e  o", "X")  # -- spans the a:fld
    assert snapshot_parts(prs) == before


# --------------------------------------------------------------------------- the invariant


def test_replace_inverse_restores_text_and_formatting():
    """replace(x→y) then replace(y→x) restores text AND formatting —
    asserted on the REOPENED document, not the live object."""
    prs = _open(BRANDED)
    paragraph = prs.slides[0].placeholders[1].text_frame.paragraphs[0]
    text_before = paragraph.text
    fingerprint_before = _char_formatting_fingerprint(paragraph)

    assert replace_text(prs, "level one", "tier 1").replacements == 1
    assert replace_text(prs, "tier 1", "level one").replacements == 1

    reopened = _reopen(save_to_bytes(prs))
    reopened_paragraph = reopened.slides[0].placeholders[1].text_frame.paragraphs[0]
    assert reopened_paragraph.text == text_before
    assert _char_formatting_fingerprint(reopened_paragraph) == fingerprint_before


def test_replace_inverse_is_exact_across_identically_formatted_runs():
    prs = _open(MINIMAL)
    box = _multi_run_paragraph(prs, [("aa", True), ("bb", True), ("cc", True)])
    paragraph = box.text_frame.paragraphs[0]
    fingerprint_before = _char_formatting_fingerprint(paragraph)

    box_id = box.shape_id
    replace_text(prs, "abbc", "XYZ")
    replace_text(prs, "XYZ", "abbc")

    reopened = _reopen(save_to_bytes(prs))
    reopened_paragraph = next(
        s for s in reopened.slides[0].shapes if s.shape_id == box_id
    ).text_frame.paragraphs[0]
    assert reopened_paragraph.text == "aabbcc"
    assert _char_formatting_fingerprint(reopened_paragraph) == fingerprint_before


def test_replace_inverse_across_mixed_formatting_restores_text_and_collapses_formatting():
    """The documented limit of the invariant: a match spanning
    differently-formatted runs collapses the span to the START run's formatting — the
    consumed runs' formatting is gone by design; the alternative is guessing."""
    prs = _open(MINIMAL)
    box = _multi_run_paragraph(prs, [("aa", True), ("bb", None), ("cc", False)])
    paragraph = box.text_frame.paragraphs[0]

    replace_text(prs, "abbc", "XYZ")
    replace_text(prs, "XYZ", "abbc")

    assert paragraph.text == "aabbcc"  # -- text restores exactly
    # -- the replaced span carries the start run's (bold) formatting
    assert [(r.text, r.font.bold) for r in paragraph.runs] == [("aabbc", True), ("c", False)]


# ------------------------------------------------------------------------ anchored variant


def _body_anchor(prs):
    inspection = inspect_text(prs.slides[0])
    return next(b.anchor for b in inspection.blocks if b.text == "Body level one")


def test_replace_text_at_edits_exactly_the_anchored_block():
    prs = _open(BRANDED)
    anchor = _body_anchor(prs)
    result = replace_text_at(prs, anchor, "level one", "tier one")
    assert result.replacements == 1
    assert prs.slides[0].placeholders[1].text_frame.paragraphs[0].text == "Body tier one"
    # -- the returned anchor is fresh: usable for a follow-up edit
    follow_up = replace_text_at(prs, result.blocks[0], "tier one", "level one")
    assert follow_up.replacements == 1


def test_stale_anchor_refuses_and_is_catchable_as_target_not_found():
    prs = _open(BRANDED)
    anchor = _body_anchor(prs)
    replace_text(prs, "Body level one", "Changed body")  # -- document moves on
    with pytest.raises(StaleAnchorError):
        replace_text_at(prs, anchor, "level", "tier")
    with pytest.raises(TargetNotFoundError):  # -- subclass contract
        replace_text_at(prs, anchor, "level", "tier")


def test_stale_anchor_refusal_is_atomic():
    prs = _open(BRANDED)
    anchor = _body_anchor(prs)
    replace_text(prs, "Body level one", "Changed body")
    before = snapshot_parts(prs)
    with pytest.raises(StaleAnchorError):
        replace_text_at(prs, anchor, "level", "tier")
    assert snapshot_parts(prs) == before


def test_find_absent_from_anchored_block_refuses_atomically():
    from .contract import assert_refusal_atomic

    prs = _open(BRANDED)
    anchor = _body_anchor(prs)
    raised = assert_refusal_atomic(
        prs,
        lambda p: replace_text_at(p, anchor, "no such text", "x"),
        TargetNotFoundError,
    )
    assert "does not occur" in str(raised)


def test_unknown_part_and_out_of_range_index_refuse_atomically():
    from .contract import assert_refusal_atomic

    prs = _open(BRANDED)
    raised = assert_refusal_atomic(
        prs,
        lambda p: replace_text_at(
            p, BlockAnchor("/ppt/slides/slide99.xml", 0, "00000000"), "a", "b"
        ),
        TargetNotFoundError,
    )
    assert "no slide or notes part" in str(raised)
    raised = assert_refusal_atomic(
        prs,
        lambda p: replace_text_at(
            p, BlockAnchor("/ppt/slides/slide1.xml", 99, "00000000"), "a", "b"
        ),
        TargetNotFoundError,
    )
    assert "beyond" in str(raised)


def test_refind_recovers_a_moved_block():
    prs = _open(BRANDED)
    anchor = _body_anchor(prs)
    # -- structural change shifts indices: new empty paragraph FIRST in the title shape
    title_p = prs.slides[0].shapes.title.text_frame.paragraphs[0]._p
    title_p.addprevious(title_p.makeelement(title_p.tag, {}))

    fresh = refind(prs, anchor)
    assert fresh.block_index != anchor.block_index
    assert fresh.content_hash == anchor.content_hash
    assert replace_text_at(prs, fresh, "level one", "tier one").replacements == 1


def test_refind_refuses_when_content_is_gone_or_ambiguous():
    prs = _open(BRANDED)
    anchor = _body_anchor(prs)
    replace_text(prs, "Body level one", "gone entirely")
    with pytest.raises(TargetNotFoundError, match="gone"):
        refind(prs, anchor)

    prs = _open(BRANDED)
    anchor = _body_anchor(prs)
    box = prs.slides[0].shapes.add_textbox(0, 0, Emu(914400), Emu(914400))
    box.text_frame.paragraphs[0].add_run().text = "Body level one"  # -- now two identical
    with pytest.raises(AmbiguousTargetError):
        refind(prs, anchor)


# ------------------------------------------------------------------------------- lo_smoke


@pytest.mark.lo_smoke
def test_replaced_output_loads_in_libreoffice(tmp_path):
    prs = _open(BRANDED)
    replace_text(prs, "Branded Title", "LO smoke title")
    out = tmp_path / "replaced.pptx"
    prs.save(str(out))
    lo_load_smoke(out, tmp_path)
