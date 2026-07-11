"""Contract tests: effective-style inspection (`pptx.inspect`).

Sidecar-driven against the branded-template and clrMap fixtures, determinism-goldened, and
independently cross-checked against LibreOffice's own resolution of the same source deck.
The API is read-only: proven here by part-snapshot equality around every call.
"""

from __future__ import annotations

import hashlib
import json

import pytest
from lxml import etree

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.errors import UnsupportedStructureError
from pptx.inspect import content_hash, effective_font, inspect_text
from pptx.util import Emu

from . import corpus
from .contract import snapshot_parts

BRANDED = "self_generated/branded_template.pptx"
CLRMAP = "self_generated/clrmap_remap.pptx"
GAUNTLET = "self_generated/gauntlet.pptx"
_A = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _open(relpath):
    return Presentation(str(corpus.fixture_path(relpath)))


def _ground_truth(relpath):
    return corpus.load_sidecar(relpath)["ground_truth"]


def _supplied_levels(effective_value):
    return [step.level for step in effective_value.provenance if step.supplied]


# ------------------------------------------------- sidecar-driven resolution (branded deck)


def test_title_size_resolves_through_layout_override():
    expected = _ground_truth(BRANDED)["expected_effective"]["title_run"]
    font = _open(BRANDED).slides[0].shapes.title.text_frame.paragraphs[0].runs[0].effective_font()
    assert font.size.resolved
    assert font.size.value_pt == expected["size_centipoints"] / 100.0
    assert font.size.value == Emu(int(expected["size_centipoints"] * 127))
    assert _supplied_levels(font.size) == ["layout placeholder lstStyle lvl1"]


def test_title_name_resolves_through_master_theme_reference():
    expected = _ground_truth(BRANDED)["expected_effective"]["title_run"]
    font = _open(BRANDED).slides[0].shapes.title.text_frame.paragraphs[0].runs[0].effective_font()
    assert font.name.resolved
    assert font.name.value == expected["font_name"]
    assert _supplied_levels(font.name) == ["theme fontScheme majorFont"]
    # -- the +mj-lt reference itself is recorded on the master step of the chain
    master_steps = [s for s in font.name.provenance if "txStyles titleStyle" in s.level]
    assert master_steps
    assert "+mj-lt" in master_steps[0].detail


def test_body_level_one_resolves_from_master_body_style():
    expected = _ground_truth(BRANDED)["expected_effective"]["body_paragraph_0_run"]
    font = (
        _open(BRANDED).slides[0].placeholders[1].text_frame.paragraphs[0].runs[0].effective_font()
    )
    assert font.size.value_pt == expected["size_centipoints"] / 100.0
    assert font.name.value == expected["font_name"]  # -- "Trebuchet MS", explicit at master
    assert _supplied_levels(font.size) == ["master txStyles bodyStyle lvl1"]
    assert _supplied_levels(font.name) == ["master txStyles bodyStyle lvl1"]


def test_body_level_two_resolves_size_from_master_and_name_from_theme():
    expected = _ground_truth(BRANDED)["expected_effective"]["body_paragraph_1_run"]
    font = (
        _open(BRANDED).slides[0].placeholders[1].text_frame.paragraphs[1].runs[0].effective_font()
    )
    assert font.size.value_pt == expected["size_centipoints"] / 100.0
    assert _supplied_levels(font.size) == ["master txStyles bodyStyle lvl2"]
    assert font.name.value == expected["font_name"]  # -- Calibri via +mn-lt
    assert _supplied_levels(font.name) == ["theme fontScheme minorFont"]


def test_every_consulted_level_appears_in_provenance_in_walk_order():
    font = (
        _open(BRANDED).slides[0].placeholders[1].text_frame.paragraphs[1].runs[0].effective_font()
    )
    levels = [step.level for step in font.size.provenance]
    assert levels == [
        "run rPr",
        "paragraph defRPr",
        "shape lstStyle lvl2",
        "layout placeholder lstStyle lvl2",
        "master placeholder lstStyle lvl2",
        "master txStyles bodyStyle lvl2",
    ]


# ------------------------------------------------------------- scheme-color resolution


def test_scheme_color_resolves_through_remapped_clrmap():
    expected = _ground_truth(CLRMAP)["expected_resolution"]
    prs = _open(CLRMAP)
    box = next(s for s in prs.slides[0].shapes if s.name == "tx1_text")
    font = box.text_frame.paragraphs[0].runs[0].effective_font()
    assert font.color_rgb.resolved
    assert font.color_rgb.value == expected["text_run_rgb"]
    details = [step.detail for step in font.color_rgb.provenance]
    assert 'schemeClr val="tx1"' in details[0]
    assert 'tx1="lt1"' in details  # -- the clrMap remap step is visible
    assert _supplied_levels(font.color_rgb) == ["theme clrScheme lt1"]


def test_slide_clrmapovr_override_beats_master_clrmap():
    prs = _open(CLRMAP)
    slide = prs.slides[0]
    clrMapOvr = slide._element.get_or_add_clrMapOvr()
    override = etree.SubElement(clrMapOvr, "{%s}overrideClrMapping" % _A)
    master_map = dict(prs.slide_masters[0].element.find(
        "{http://schemas.openxmlformats.org/presentationml/2006/main}clrMap").attrib)
    master_map["tx1"] = "accent3"  # -- 9BBB59 in the default theme
    for key, value in master_map.items():
        override.set(key, value)

    box = next(s for s in slide.shapes if s.name == "tx1_text")
    font = box.text_frame.paragraphs[0].runs[0].effective_font()
    assert font.color_rgb.value == "9BBB59"
    assert any(step.level == "slide clrMapOvr" for step in font.color_rgb.provenance)


def test_direct_srgb_color_supplies_immediately():
    prs = _open(BRANDED)
    run = prs.slides[0].shapes.title.text_frame.paragraphs[0].runs[0]
    run.font.color.rgb = RGBColor(0x12, 0x34, 0xAB)
    font = run.effective_font()
    assert font.color_rgb.value == "1234AB"
    assert _supplied_levels(font.color_rgb) == ["run rPr"]


@pytest.mark.parametrize("color_tag", ["srgbClr", "schemeClr"])
def test_transformed_color_is_unresolved_instead_of_reporting_base_rgb(color_tag):
    prs = _open("self_generated/minimal_clean.pptx")
    box = prs.slides[0].shapes.add_textbox(0, 0, 914400, 914400)
    run = box.text_frame.paragraphs[0].add_run()
    run.text = "transformed"
    rPr = run._r.get_or_add_rPr()
    solidFill = etree.SubElement(rPr, "{%s}solidFill" % _A)
    color = etree.SubElement(solidFill, "{%s}%s" % (_A, color_tag))
    color.set("val", "804020" if color_tag == "srgbClr" else "accent1")
    etree.SubElement(color, "{%s}lumMod" % _A).set("val", "50000")

    effective = run.effective_font().color_rgb
    assert not effective.resolved
    assert effective.value is None
    assert any("unapplied transforms" in step.detail for step in effective.provenance)


# ----------------------------------------------------- non-placeholder + honesty cases


def test_plain_textbox_resolves_via_presentation_default_text_style():
    prs = _open("self_generated/minimal_clean.pptx")
    box = prs.slides[0].shapes.add_textbox(0, 0, 914400, 914400)
    box.text_frame.paragraphs[0].add_run().text = "plain"
    font = box.text_frame.paragraphs[0].runs[0].effective_font()
    assert (font.size.value_pt, font.name.value, font.color_rgb.value) == (
        18.0,
        "Calibri",
        "000000",
    )
    assert _supplied_levels(font.size) == ["presentation defaultTextStyle lvl1"]


def test_autoshape_fontref_supplies_theme_font_and_text_color():
    prs = _open("self_generated/minimal_clean.pptx")
    shape = prs.slides[0].shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, 914400, 914400)
    shape.text = "font ref"
    font = shape.text_frame.paragraphs[0].runs[0].effective_font()
    assert font.name.value == "Calibri"
    assert font.color_rgb.value == "FFFFFF"
    assert "shape fontRef" in [step.level for step in font.name.provenance]


def test_non_latin_run_does_not_report_the_latin_typeface_as_effective():
    prs = _open("self_generated/minimal_clean.pptx")
    box = prs.slides[0].shapes.add_textbox(0, 0, 914400, 914400)
    run = box.text_frame.paragraphs[0].add_run()
    run.text = "日本語"
    run.font.name = "Arial"
    etree.SubElement(run._r.get_or_add_rPr(), "{%s}ea" % _A).set("typeface", "Yu Mincho")

    name = run.effective_font().name
    assert not name.resolved
    assert name.value is None
    assert "non-Latin" in name.provenance[-1].detail


def test_gradient_text_fill_reports_unresolved_not_a_guess():
    prs = _open("self_generated/minimal_clean.pptx")
    box = prs.slides[0].shapes.add_textbox(0, 0, 914400, 914400)
    run = box.text_frame.paragraphs[0].add_run()
    run.text = "gradient"
    rPr = run._r.get_or_add_rPr()
    etree.SubElement(rPr, "{%s}gradFill" % _A)
    font = box.text_frame.paragraphs[0].runs[0].effective_font()
    assert not font.color_rgb.resolved
    assert font.color_rgb.value is None
    assert "gradFill" in font.color_rgb.provenance[-1].detail


def test_phclr_scheme_token_reports_unresolved():
    prs = _open("self_generated/minimal_clean.pptx")
    box = prs.slides[0].shapes.add_textbox(0, 0, 914400, 914400)
    run = box.text_frame.paragraphs[0].add_run()
    run.text = "phClr"
    rPr = run._r.get_or_add_rPr()
    solidFill = etree.SubElement(rPr, "{%s}solidFill" % _A)
    schemeClr = etree.SubElement(solidFill, "{%s}schemeClr" % _A)
    schemeClr.set("val", "phClr")
    font = run.effective_font()
    assert not font.color_rgb.resolved
    assert any("unmappable" in step.detail for step in font.color_rgb.provenance)


def test_table_cell_run_refuses_instead_of_guessing():
    prs = _open(GAUNTLET)
    table_shape = next(s for s in prs.slides[2].shapes if s.name == "gauntlet_table")
    cell_run = table_shape.table.cell(0, 0).text_frame.paragraphs[0].runs[0]
    with pytest.raises(UnsupportedStructureError):
        effective_font(cell_run)


# ---------------------------------------------------------------- read-only + determinism


def test_inspection_never_mutates_the_package():
    prs = _open(BRANDED)
    before = snapshot_parts(prs)
    inspect_text(prs.slides[0])
    for shape in prs.slides[0].shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.effective_font()
    assert snapshot_parts(prs) == before


def test_inspection_is_deterministic_across_runs_and_loads():
    first = json.dumps(inspect_text(_open(BRANDED).slides[0]).to_dict())
    second = json.dumps(inspect_text(_open(BRANDED).slides[0]).to_dict())
    assert first == second


@pytest.mark.parametrize(
    ("golden_name", "fixture_relpath", "slide_index"),
    [
        ("branded_template.inspect.json", BRANDED, 0),
        ("clrmap_remap.inspect.json", CLRMAP, 0),
        ("gauntlet_slide1.inspect.json", GAUNTLET, 0),
    ],
)
def test_inspection_matches_frozen_golden(golden_name, fixture_relpath, slide_index):
    """Byte-identical to the reviewed golden; update ONLY via update_goldens.py + PR review."""
    golden_path = corpus.FIXTURES_DIR.parent / "goldens" / golden_name
    payload = inspect_text(_open(fixture_relpath).slides[slide_index]).to_dict()
    actual = (json.dumps(payload, indent=2, ensure_ascii=False) + "\n").encode("utf-8")
    assert actual == golden_path.read_bytes()  # -- byte-exact, no newline translation


def test_libreoffice_independently_confirms_branded_effective_sizes():
    """LO resolved the same inheritance when it baked run sizes into its re-export."""
    ours = sorted(
        run.font.size.value_pt
        for block in inspect_text(_open(BRANDED).slides[0]).blocks
        for run in block.runs
    )
    lo_prs = Presentation(str(corpus.fixture_path("libreoffice_export/lo_branded_template.pptx")))
    theirs = sorted(
        r.effective_font().size.value_pt
        for s in lo_prs.slides[0].shapes
        if s.has_text_frame
        for p in s.text_frame.paragraphs
        for r in p.runs
    )
    assert ours == theirs == [22.0, 26.0, 36.0]


# --------------------------------------------------------------------------------- anchors


# -------------------------------------------------- visibility completeness


def test_inspect_text_sees_table_cells_as_counted_blind_regions():
    """Table-cell text appears (row-major) as typed blind regions, never silence."""
    inspection = inspect_text(_open("self_generated/tables_in_group.pptx").slides[0])
    by_container = {}
    for block in inspection.blocks:
        by_container.setdefault(block.container, []).append(block)

    cell_blocks = by_container.get("table-cell", [])
    assert [b.text for b in cell_blocks] == [
        "cell r0c0", "cell r0c1", "cell r1c0", "cell r1c1"
    ]
    assert [b.container_detail for b in cell_blocks] == [
        "grouped_table!r0c0", "grouped_table!r0c1",
        "grouped_table!r1c0", "grouped_table!r1c1",
    ]
    assert all(b.blind for b in cell_blocks)
    for block in cell_blocks:
        for run in block.runs:
            assert run.font.size.resolved is False  # -- honest unresolved, not a guess
    assert inspection.blind_region_count == 4

    payload = inspection.to_dict()
    assert payload["version"] == 2
    assert payload["blind_region_count"] == 4


def test_inspect_text_sees_grouped_shape_text_with_group_paths():
    """Text inside groups appears, recursively, with its group path."""
    inspection = inspect_text(_open("self_generated/nested_groups.pptx").slides[0])
    by_text = {b.text: b for b in inspection.blocks}
    assert by_text["Level zero"].container == "shape"
    assert by_text["Level zero"].container_detail is None
    assert by_text["Level one"].container == "group"
    assert by_text["Level one"].container_detail == "group_level1"
    assert by_text["Level two"].container_detail == "group_level1/group_level2"
    assert by_text["Level three"].container_detail == (
        "group_level1/group_level2/group_level3"
    )
    assert not any(b.blind for b in inspection.blocks)
    # -- grouped runs resolve through the normal (non-placeholder) chain
    assert by_text["Level three"].runs[0].font.size.resolved is True


def test_inspect_text_sees_in_group_textbox_beside_grouped_table():
    inspection = inspect_text(_open("self_generated/tables_in_group.pptx").slides[0])
    by_text = {b.text: b for b in inspection.blocks}
    assert by_text["Top-level text"].container == "shape"
    assert by_text["In-group text"].container == "group"
    assert by_text["In-group text"].container_detail == "outer_group"


def test_inspect_text_block_order_is_depth_first_document_order():
    inspection = inspect_text(_open("self_generated/tables_in_group.pptx").slides[0])
    texts = [b.text for b in inspection.blocks]
    assert texts == [
        "Top-level text", "In-group text",
        "cell r0c0", "cell r0c1", "cell r1c0", "cell r1c1",
    ]
    assert [b.anchor.block_index for b in inspection.blocks] == list(range(len(texts)))


def test_pathological_group_nesting_refuses_instead_of_recursing_forever():
    prs = _open("self_generated/minimal_clean.pptx")
    shapes = prs.slides[0].shapes
    group = shapes.add_group_shape()
    for _ in range(17):
        group = group.shapes.add_group_shape()
    box = group.shapes.add_textbox(0, 0, 914400, 914400)
    box.text_frame.paragraphs[0].add_run().text = "too deep"
    with pytest.raises(UnsupportedStructureError, match="nested"):
        inspect_text(prs.slides[0])


def test_effective_font_resolves_runs_inside_groups():
    prs = _open("self_generated/nested_groups.pptx")
    group = next(s for s in prs.slides[0].shapes if s.name == "group_level1")
    for name in ("group_level2", "group_level3"):
        group = next(s for s in group.shapes if s.name == name)
    box = next(s for s in group.shapes if s.name == "level3_box")
    info = box.text_frame.paragraphs[0].runs[0].effective_font()
    assert info.size.resolved is True


def test_content_hash_is_pinned_sha256_nfc_prefix():
    assert content_hash("Branded Title") == (
        hashlib.sha256("Branded Title".encode("utf-8")).hexdigest()[:8]
    )


def test_content_hash_applies_nfc_normalization():
    composed = "café"  # -- é as one code point
    decomposed = "café"  # -- e + combining acute
    assert composed != decomposed
    assert content_hash(composed) == content_hash(decomposed)


def test_out_of_schema_indent_level_refuses_instead_of_crashing():
    prs = _open(BRANDED)
    run = prs.slides[0].shapes.title.text_frame.paragraphs[0].runs[0]
    pPr = run._r.getparent().find("{%s}pPr" % _A)
    if pPr is None:
        pPr = etree.SubElement(run._r.getparent(), "{%s}pPr" % _A)
        run._r.getparent().insert(0, pPr)
    pPr.set("lvl", "9")  # -- outside ST_TextIndentLevelType's 0..8
    with pytest.raises(UnsupportedStructureError, match="lvl=9"):
        run.effective_font()


def test_alternate_content_is_a_typed_counted_blind_region_not_silence():
    """Regression (review): mc:AlternateContent used to be invisible to inspect_text and
    inspect_deck — a fail-silent hole in the visibility-complete contract."""
    from pptx.inspect import inspect_deck
    from tests.paper.test_edit_text import _wrap_first_textbox_in_alternate_content

    prs = _wrap_first_textbox_in_alternate_content(
        _open("self_generated/minimal_clean.pptx")
    )
    inspection = inspect_text(prs.slides[0])
    ac_blocks = [b for b in inspection.blocks if b.container == "alternate-content"]
    assert len(ac_blocks) == 1
    assert ac_blocks[0].blind is True
    assert inspection.blind_region_count == 1
    assert inspection.to_dict()["blind_region_count"] == 1

    manifest = inspect_deck(prs)
    assert manifest.slides[0].alternate_content_count == 1
    assert manifest.to_dict()["slides"][0]["alternate_content_count"] == 1


def test_boolean_effective_values_serialize_as_json_booleans():
    """Regression (review): bool is an int subclass; the payload used to emit 0/1."""
    payload = (
        _open(BRANDED).slides[0].placeholders[1].text_frame.paragraphs[0].runs[0]
        .effective_font()
        .to_dict()
    )
    assert payload["bold"]["value"] is False
    assert payload["italic"]["value"] is False


# --------------------------------------------------------- walk extensions


def test_bold_italic_underline_resolve_with_schema_defaults():
    prs = _open(BRANDED)
    run = prs.slides[0].placeholders[1].text_frame.paragraphs[0].runs[0]
    font = run.effective_font()
    assert font.bold.value is False
    assert font.bold.resolved is True
    assert font.bold.provenance[-1].level == "schema default"
    assert font.italic.value is False
    assert font.italic.resolved is True
    assert font.underline.value == "none"
    assert font.underline.resolved is True

    run.font.bold = True
    run.font.underline = True  # -- upstream writes u="sng"
    font = run.effective_font()
    assert font.bold.value is True
    assert font.bold.provenance[-1].level == "run rPr"
    assert font.underline.value == "sng"


def test_effective_paragraph_format_resolves_alignment_and_spacing():
    from pptx.enum.text import PP_ALIGN
    from pptx.inspect import effective_paragraph_format

    prs = _open(BRANDED)
    paragraph = prs.slides[0].placeholders[1].text_frame.paragraphs[0]
    fmt = effective_paragraph_format(paragraph)
    assert fmt.alignment.value == "l"  # -- master bodyStyle supplies algn="l"
    assert any(s.supplied and "bodyStyle" in s.level for s in fmt.alignment.provenance)
    assert fmt.line_spacing.value == 1.0  # -- rendering default, explicitly provenanced
    assert fmt.line_spacing.provenance[-1].level == "rendering default"

    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.line_spacing = 1.5
    fmt = effective_paragraph_format(paragraph)
    assert fmt.alignment.value == "ctr"
    assert fmt.alignment.provenance[0].supplied is True  # -- paragraph pPr wins
    assert fmt.line_spacing.value == 1.5

    payload = fmt.to_dict()
    assert payload["schema"] == "paper-effective-paragraph-format"
    assert payload["version"] == 1


def test_effective_shape_format_resolves_explicit_fill_through_clrmap():
    """A probe case: the rectangle behind the text, resolved like the text."""
    from pptx.inspect import effective_shape_format

    prs = _open(CLRMAP)
    rect = prs.slides[0].shapes.shape_by_name("accent1_box")
    fmt = effective_shape_format(rect)
    assert fmt.fill_rgb.value == "C0504D"  # -- accent1 -> clrMap -> theme accent2
    assert fmt.fill_rgb.resolved is True
    assert any("clrMap" in s.detail or "clrMap" in s.level for s in fmt.fill_rgb.provenance)

    # -- line color comes only from the style lnRef: honestly unresolved, reference
    # -- color carried in provenance
    assert fmt.line_rgb.resolved is False
    assert any("C0504D" in s.detail for s in fmt.line_rgb.provenance)

    payload = fmt.to_dict()
    assert payload["schema"] == "paper-effective-shape-format"
    assert payload["version"] == 1


def test_effective_shape_format_reports_nofill_and_absent_fill_honestly():
    from pptx.inspect import effective_shape_format

    prs = _open(CLRMAP)
    rect = prs.slides[0].shapes.shape_by_name("accent1_box")
    rect.fill.background()  # -- a:noFill
    fmt = effective_shape_format(rect)
    assert fmt.fill_rgb.value == "none"
    assert fmt.fill_rgb.resolved is True

    # -- upstream textboxes carry an explicit a:noFill too
    box = prs.slides[0].shapes.shape_by_name("tx1_text")
    assert effective_shape_format(box).fill_rgb.value == "none"

    # -- a placeholder with an empty spPr and no p:style is honestly unresolved
    branded = _open(BRANDED)
    title = branded.slides[0].shapes.title
    fmt = effective_shape_format(title)
    assert fmt.fill_rgb.resolved is False
    assert fmt.fill_rgb.value is None


def test_effective_font_payload_carries_pinned_schema_keys():
    payload = (
        _open(BRANDED).slides[0].shapes.title.text_frame.paragraphs[0].runs[0]
        .effective_font()
        .to_dict()
    )
    assert payload["schema"] == "paper-effective-font"
    assert payload["version"] == 2  # -- v2: bold/italic/underline added


def test_content_hash_treats_whitespace_as_content():
    assert content_hash("Trailing space ") != content_hash("Trailing space")


def test_blocks_carry_stable_anchors_in_sptree_order():
    inspection = inspect_text(_open(BRANDED).slides[0])
    assert [block.anchor.block_index for block in inspection.blocks] == list(
        range(len(inspection.blocks))
    )
    title_block = inspection.blocks[0]
    assert title_block.anchor.part == "/ppt/slides/slide1.xml"
    assert title_block.anchor.content_hash == content_hash("Branded Title")
    assert title_block.placeholder_type == "TITLE"
