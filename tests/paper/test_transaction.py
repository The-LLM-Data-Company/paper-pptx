"""Focused contracts for package-level mutation rollback."""

from __future__ import annotations

import io

import pytest
from lxml import etree

from pptx import Presentation
from pptx._transaction import PackageTransaction
from pptx.errors import UnsupportedStructureError
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.opc.package import Part, XmlPart
from pptx.opc.packuri import PackURI
from pptx.util import Inches

from .contract import save_to_bytes, zip_member_map


def _presentation_with_textbox(text="Before"):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    shape.text = text
    return prs, slide, shape


def test_failure_restores_live_proxies_relationships_and_saved_package():
    prs, slide, shape = _presentation_with_textbox()
    package = prs.part.package
    relationships = slide.part.rels
    before = zip_member_map(save_to_bytes(prs))
    refusal = UnsupportedStructureError("forced refusal")

    with pytest.raises(UnsupportedStructureError) as exc_info:
        with PackageTransaction(package, prs, slide, shape):
            shape.text = "Dirty"
            relationships.get_or_add_ext_rel(RT.HYPERLINK, "https://example.invalid")
            raise refusal

    assert exc_info.value is refusal
    assert shape.text == "Before"
    assert slide.part.rels is relationships
    assert all(rel.target_ref != "https://example.invalid" for rel in relationships.values())
    assert zip_member_map(save_to_bytes(prs)) == before


def test_nested_failure_restores_to_inner_entry_state():
    prs, _slide, shape = _presentation_with_textbox()
    package = prs.part.package

    with PackageTransaction(package, shape):
        shape.text = "Outer"
        with pytest.raises(ValueError, match="inner failed"):
            with PackageTransaction(package, shape):
                shape.text = "Inner"
                raise ValueError("inner failed")
        assert shape.text == "Outer"

    assert shape.text == "Outer"


def test_preconstructed_nested_transaction_snapshots_context_entry():
    prs, _slide, shape = _presentation_with_textbox()
    package = prs.part.package
    inner = PackageTransaction(package, shape)

    with PackageTransaction(package, shape):
        shape.text = "Outer"
        with pytest.raises(ValueError, match="inner failed"):
            with inner:
                shape.text = "Inner"
                raise ValueError("inner failed")
        assert shape.text == "Outer"


def test_candidate_validation_failure_rolls_back_live_state(monkeypatch):
    prs, _slide, shape = _presentation_with_textbox()
    package = prs.part.package
    before = zip_member_map(save_to_bytes(prs))

    def refuse_candidate(_transaction):
        raise UnsupportedStructureError("candidate refused")

    monkeypatch.setattr(PackageTransaction, "_validate_candidate", refuse_candidate)
    with pytest.raises(UnsupportedStructureError, match="candidate refused"):
        with PackageTransaction(package, shape):
            shape.text = "Dirty"

    assert shape.text == "Before"
    assert zip_member_map(save_to_bytes(prs)) == before


def test_successful_transaction_commits_a_reopenable_candidate():
    prs, _slide, shape = _presentation_with_textbox()

    with PackageTransaction(prs.part.package, shape):
        shape.text = "After"

    reopened = Presentation(io.BytesIO(save_to_bytes(prs)))
    assert shape.text == "After"
    assert reopened.slides[0].shapes[-1].text == "After"


def test_failure_restores_custom_xml_nodes_and_binary_payloads():
    prs = Presentation()
    package = prs.part.package
    binary = Part(
        PackURI("/custom/payload.bin"), "application/octet-stream", package, b"original"
    )
    root = etree.fromstring(b"<paper-state><child/></paper-state>")
    child = root[0]
    xml = XmlPart(PackURI("/custom/state.xml"), "application/xml", package, root)
    package.relate_to(binary, "https://paper.example/relationships/binary")
    package.relate_to(xml, "https://paper.example/relationships/xml")

    with pytest.raises(RuntimeError, match="forced failure"):
        with PackageTransaction(package):
            binary._blob = b"dirty"
            root.set("dirty", "1")
            root.remove(child)
            raise RuntimeError("forced failure")

    assert binary.blob == b"original"
    assert xml._element is root
    assert root.get("dirty") is None
    assert root[0] is child
