"""Corpus integrity tests: frozen hashes, sidecar schema, and ground truth vs. actual bytes.

Ground-truth checks read the fixture files directly (zipfile + lxml), not through the pptx
object model, so a future object-model bug cannot silently re-verify the corpus against itself.
"""

from __future__ import annotations

import zipfile

import pytest
from lxml import etree

from pptx import Presentation

from . import corpus
from .contract import save_reopen, zip_member_map
from .idlists import dangling_section_slide_ids, duplicate_section_slide_ids
from .lo import lo_load_smoke
from .relint import dangling_relationship_targets, missing_relationship_references

_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

ALL_RELPATHS = corpus.iter_fixture_relpaths()
NONCORRUPT_RELPATHS = [r for r in ALL_RELPATHS if not corpus.is_corrupt_fixture(r)]


def _member_xml(relpath, member):
    with zipfile.ZipFile(str(corpus.fixture_path(relpath))) as zipf:
        return etree.fromstring(zipf.read(member))


def _member_names(relpath):
    with zipfile.ZipFile(str(corpus.fixture_path(relpath))) as zipf:
        return zipf.namelist()


def _ground_truth(relpath):
    return corpus.load_sidecar(relpath)["ground_truth"]


# ------------------------------------------------------------------ frozen-corpus discipline


def test_manifest_lists_exactly_the_fixture_files():
    assert sorted(corpus.manifest_entries()) == ALL_RELPATHS, (
        "MANIFEST.sha256 out of sync with the fixture files on disk"
    )


def test_every_fixture_lives_in_a_known_provenance_bucket():
    unknown = [r for r in ALL_RELPATHS if r.split("/")[0] not in corpus.BUCKETS]
    assert not unknown, "fixtures outside the pinned provenance buckets: %r" % unknown


@pytest.mark.parametrize("relpath", ALL_RELPATHS)
def test_fixture_hash_matches_manifest(relpath):
    expected = corpus.manifest_entries()[relpath]
    actual = corpus.sha256_of(corpus.fixture_path(relpath))
    assert actual == expected, (
        "fixture %s hash changed; fixtures are frozen - a changed fixture requires a new "
        "manifest entry and sidecar in a reviewed PR" % relpath
    )


@pytest.mark.parametrize("relpath", ALL_RELPATHS)
def test_fixture_has_valid_sidecar(relpath):
    assert corpus.sidecar_path(relpath).is_file(), "fixture %s has no sidecar" % relpath
    problems = corpus.sidecar_schema_problems(corpus.load_sidecar(relpath), relpath)
    assert not problems, "sidecar schema violations for %s: %s" % (relpath, "; ".join(problems))


def test_no_orphan_sidecars():
    fixture_stems = {corpus.fixture_path(r) for r in ALL_RELPATHS}
    orphans = [
        p.relative_to(corpus.FIXTURES_DIR).as_posix()
        for p in corpus.FIXTURES_DIR.rglob("*.json")
        if p.with_suffix(".pptx") not in fixture_stems
    ]
    assert not orphans, "sidecars without a fixture: %r" % orphans


def test_no_fixtures_outside_buckets():
    """A .pptx directly under fixtures/ would dodge every discipline test above - forbid it."""
    strays = [p.name for p in corpus.FIXTURES_DIR.glob("*.pptx")]
    assert strays == [], "fixtures must live inside a provenance bucket: %r" % strays


def test_corpus_partition_shape():
    """Guard the corrupt/non-corrupt split itself: a regression in is_corrupt_fixture toward
    'everything corrupt' would silently empty three parametrized suites into green skips."""
    corrupt = [r for r in ALL_RELPATHS if corpus.is_corrupt_fixture(r)]
    assert corrupt == ["self_generated/corrupt_dangling_sldid.pptx"]
    assert len(NONCORRUPT_RELPATHS) == len(ALL_RELPATHS) - 1


# ------------------------------------------------------------------------ basic viability


@pytest.mark.parametrize("relpath", NONCORRUPT_RELPATHS)
def test_noncorrupt_fixture_opens_and_survives_save_reopen(relpath):
    prs = Presentation(str(corpus.fixture_path(relpath)))
    reopened = save_reopen(prs)
    assert reopened is not prs
    assert len(reopened.slides) == len(prs.slides)
    ground_truth = _ground_truth(relpath)
    if "slide_count" in ground_truth:
        assert len(reopened.slides) == ground_truth["slide_count"]


@pytest.mark.parametrize("relpath", NONCORRUPT_RELPATHS)
def test_noncorrupt_fixture_has_relationship_integrity(relpath):
    zip_map = zip_member_map(corpus.fixture_path(relpath).read_bytes())
    assert dangling_relationship_targets(zip_map) == []
    assert missing_relationship_references(zip_map) == []
    assert dangling_section_slide_ids(zip_map) == []
    assert duplicate_section_slide_ids(zip_map) == []


def test_corrupt_fixture_missing_reference_is_detected():
    relpath = "self_generated/corrupt_dangling_sldid.pptx"
    ground_truth = _ground_truth(relpath)
    zip_map = zip_member_map(corpus.fixture_path(relpath).read_bytes())
    assert missing_relationship_references(zip_map) == [
        ("ppt/presentation.xml", ground_truth["dangling_r_id"])
    ]
    assert dangling_relationship_targets(zip_map) == []


def test_corrupt_fixture_behaves_as_documented():
    relpath = "self_generated/corrupt_dangling_sldid.pptx"
    ground_truth = _ground_truth(relpath)
    prs = Presentation(str(corpus.fixture_path(relpath)))  # -- opens: parts load lazily
    assert ground_truth["opens_with_python_pptx"] is True
    with pytest.raises(KeyError):
        list(prs.slides)


# --------------------------------------------------------------- ground truth vs. real bytes


def test_minimal_clean_ground_truth():
    relpath = "self_generated/minimal_clean.pptx"
    ground_truth = _ground_truth(relpath)
    slide = _member_xml(relpath, "ppt/slides/slide1.xml")
    assert [t.text for t in slide.iter("{%s}t" % _A)] == [
        ground_truth["title_text"],
        ground_truth["subtitle_text"],
    ]
    rels = _member_xml(relpath, "ppt/slides/_rels/slide1.xml.rels")
    layout_targets = [
        rel.get("Target") for rel in rels if rel.get("Type").endswith("/slideLayout")
    ]
    assert len(layout_targets) == 1
    layout_member = "ppt/" + layout_targets[0].replace("../", "")
    layout = _member_xml(relpath, layout_member)
    assert layout.find("{%s}cSld" % _P).get("name") == ground_truth["layout_name"]


def test_branded_template_ground_truth():
    relpath = "self_generated/branded_template.pptx"
    ground_truth = _ground_truth(relpath)
    master = _member_xml(relpath, "ppt/slideMasters/slideMaster1.xml")
    txStyles = master.find("{%s}txStyles" % _P)
    title_defRPr = txStyles.find("{%s}titleStyle/{%s}lvl1pPr/{%s}defRPr" % (_P, _A, _A))
    assert int(title_defRPr.get("sz")) == ground_truth["master_title_lvl1_sz_centipoints"]
    assert (
        title_defRPr.find("{%s}latin" % _A).get("typeface")
        == ground_truth["master_title_lvl1_latin"]
    )
    body_l1 = txStyles.find("{%s}bodyStyle/{%s}lvl1pPr/{%s}defRPr" % (_P, _A, _A))
    assert int(body_l1.get("sz")) == ground_truth["master_body_lvl1_sz_centipoints"]
    assert body_l1.find("{%s}latin" % _A).get("typeface") == ground_truth["master_body_lvl1_latin"]
    body_l2 = txStyles.find("{%s}bodyStyle/{%s}lvl2pPr/{%s}defRPr" % (_P, _A, _A))
    assert int(body_l2.get("sz")) == ground_truth["master_body_lvl2_sz_centipoints"]
    assert body_l2.find("{%s}latin" % _A).get("typeface") == ground_truth["master_body_lvl2_latin"]

    layout = _member_xml(relpath, ground_truth["layout_part_with_title_override"])
    overrides = [
        defRPr.get("sz")
        for sp in layout.iter("{%s}sp" % _P)
        if (ph := sp.find(".//{%s}ph" % _P)) is not None and ph.get("type") == "title"
        for defRPr in sp.findall(".//{%s}lstStyle/{%s}lvl1pPr/{%s}defRPr" % (_A, _A, _A))
    ]
    assert overrides == [str(ground_truth["layout_title_lvl1_sz_centipoints"])]

    theme = _member_xml(relpath, "ppt/theme/theme1.xml")
    font_scheme = theme.find("{%s}themeElements/{%s}fontScheme" % (_A, _A))
    assert (
        font_scheme.find("{%s}majorFont/{%s}latin" % (_A, _A)).get("typeface")
        == ground_truth["theme_major_latin"]
    )
    assert (
        font_scheme.find("{%s}minorFont/{%s}latin" % (_A, _A)).get("typeface")
        == ground_truth["theme_minor_latin"]
    )

    slide = _member_xml(relpath, "ppt/slides/slide1.xml")
    for run in slide.iter("{%s}r" % _A):
        rPr = run.find("{%s}rPr" % _A)
        has_local = rPr is not None and (
            rPr.get("sz") is not None or rPr.find("{%s}latin" % _A) is not None
        )
        assert has_local == ground_truth["slide_runs_carry_local_size_or_font"]


def test_clrmap_remap_ground_truth():
    relpath = "self_generated/clrmap_remap.pptx"
    ground_truth = _ground_truth(relpath)
    master = _member_xml(relpath, "ppt/slideMasters/slideMaster1.xml")
    clrMap = master.find("{%s}clrMap" % _P)
    assert dict(clrMap.attrib) == ground_truth["clrmap"]

    theme = _member_xml(relpath, "ppt/theme/theme1.xml")
    scheme = theme.find("{%s}themeElements/{%s}clrScheme" % (_A, _A))
    assert (
        scheme.find("{%s}accent1" % _A)[0].get("val") == ground_truth["theme_accent1_srgb"]
    )
    assert (
        scheme.find("{%s}accent2" % _A)[0].get("val") == ground_truth["theme_accent2_srgb"]
    )
    assert (
        scheme.find("{%s}lt1" % _A)[0].get("lastClr") == ground_truth["theme_lt1_sysclr_lastclr"]
    )

    slide = _member_xml(relpath, "ppt/slides/slide1.xml")
    scheme_vals = [c.get("val") for c in slide.iter("{%s}schemeClr" % _A)]
    assert ground_truth["rectangle_fill_schemeclr_val"] in scheme_vals
    assert ground_truth["text_run_schemeclr_val"] in scheme_vals


@pytest.mark.parametrize(
    "relpath",
    ["self_generated/chart_notes.pptx", "libreoffice_export/lo_chart_notes.pptx"],
)
def test_chart_notes_ground_truth(relpath):
    ground_truth = _ground_truth(relpath)
    names = _member_names(relpath)
    assert ground_truth["chart_part"] in names
    if "embedded_workbook_part" in ground_truth:
        assert ground_truth["embedded_workbook_part"] in names
    else:
        embeddings = [n for n in names if n.startswith("ppt/embeddings/")]
        assert embeddings == ground_truth["embedded_workbook_parts"]
    notes_members = [n for n in names if n.startswith("ppt/notesSlides/") and n.endswith(".xml")]
    assert len(notes_members) == 1
    notes_texts = [t.text for t in _member_xml(relpath, notes_members[0]).iter("{%s}t" % _A)]
    if "notes_text" in ground_truth:
        assert notes_texts == [ground_truth["notes_text"]]
    else:
        assert notes_texts == ground_truth["notes_text_values"]


def test_chart_notes_series_ground_truth():
    relpath = "self_generated/chart_notes.pptx"
    ground_truth = _ground_truth(relpath)
    chart_ns = "http://schemas.openxmlformats.org/drawingml/2006/chart"
    chart = _member_xml(relpath, ground_truth["chart_part"])
    series_names = [
        v.text
        for ser in chart.iter("{%s}ser" % chart_ns)
        for v in ser.find("{%s}tx" % chart_ns).iter("{%s}v" % chart_ns)
    ]
    assert series_names == ground_truth["series_names"]


@pytest.mark.parametrize(
    "relpath",
    ["self_generated/shared_media.pptx", "libreoffice_export/lo_shared_media.pptx"],
)
def test_shared_media_ground_truth(relpath):
    ground_truth = _ground_truth(relpath)
    names = _member_names(relpath)
    assert [n for n in names if n.startswith("ppt/media/")] == ground_truth["media_parts"]
    for slide_number in (1, 2):
        rels = _member_xml(relpath, "ppt/slides/_rels/slide%d.xml.rels" % slide_number)
        targets = [rel.get("Target") for rel in rels if rel.get("Type").endswith("/image")]
        assert targets == [ground_truth["slide%d_image_target" % slide_number]]


@pytest.mark.parametrize(
    "relpath",
    [
        "self_generated/autofit_none.pptx",
        "self_generated/autofit_normal.pptx",
        "self_generated/autofit_shape.pptx",
        "libreoffice_export/lo_autofit_normal.pptx",
    ],
)
def test_autofit_ground_truth(relpath):
    ground_truth = _ground_truth(relpath)
    slide = _member_xml(relpath, "ppt/slides/slide1.xml")
    bodyPr = slide.find(".//{%s}bodyPr" % _A)
    children = list(bodyPr)
    assert len(children) == 1
    qname = etree.QName(children[0])
    assert "a:" + qname.localname == ground_truth["bodypr_autofit_element"]
    if ground_truth["bodypr_autofit_element"] == "a:normAutofit":
        # -- hard-indexed: every normAutofit sidecar must pin both attrs (null when absent)
        assert children[0].get("fontScale") == ground_truth["fontscale_attr"]
        assert children[0].get("lnSpcReduction") == ground_truth["lnspcreduction_attr"]


def test_whitespace_pair_ground_truth():
    for suffix in ("a", "b"):
        relpath = "self_generated/whitespace_trailing_%s.pptx" % suffix
        ground_truth = _ground_truth(relpath)
        slide = _member_xml(relpath, "ppt/slides/slide1.xml")
        texts = [t.text for t in slide.iter("{%s}t" % _A)]
        assert texts == [ground_truth["run_text"]]

        pair_relpath = "self_generated/" + ground_truth["pair_fixture"]
        own = zip_member_map(corpus.fixture_path(relpath).read_bytes())
        pair = zip_member_map(corpus.fixture_path(pair_relpath).read_bytes())
        assert sorted(own) == sorted(pair)
        differing = sorted(n for n in own if own[n] != pair[n])
        assert differing == ground_truth["differing_zip_members_vs_pair"]


def test_gauntlet_ground_truth():
    relpath = "self_generated/gauntlet.pptx"
    ground_truth = _ground_truth(relpath)
    names = _member_names(relpath)
    slide_members = [
        n for n in names if n.startswith("ppt/slides/slide") and n.endswith(".xml")
    ]
    assert len(slide_members) == ground_truth["slide_count"]
    assert ground_truth["chart_part"] in names
    assert ground_truth["embedded_workbook_part"] in names
    assert ground_truth["notes_slide_part"] in names
    assert [n for n in names if n.startswith("ppt/media/")] == ground_truth["media_parts"]

    master = _member_xml(relpath, "ppt/slideMasters/slideMaster1.xml")
    assert master.find("{%s}clrMap" % _P).get("accent1") == ground_truth["clrmap_accent1"]
    body_l1 = master.find(
        "{%s}txStyles/{%s}bodyStyle/{%s}lvl1pPr/{%s}defRPr" % (_P, _P, _A, _A)
    )
    assert int(body_l1.get("sz")) == ground_truth["master_body_lvl1_sz_centipoints"]

    slide3 = _member_xml(relpath, "ppt/slides/slide3.xml")
    assert [b.get("char") for b in slide3.iter("{%s}buChar" % _A)] == [
        ground_truth["real_bullet"]["bu_char"]
    ]
    assert [b.get("typeface") for b in slide3.iter("{%s}buFont" % _A)] == [
        ground_truth["real_bullet"]["bu_font"]
    ]
    src_rects = [dict(sr.attrib) for sr in slide3.iter("{%s}srcRect" % _A)]
    assert src_rects == [
        {
            "l": ground_truth["cropped_picture"]["srcrect_l"],
            "t": ground_truth["cropped_picture"]["srcrect_t"],
        }
    ]

    rels4 = _member_xml(relpath, "ppt/slides/_rels/slide4.xml.rels")
    external_targets = [
        rel.get("Target") for rel in rels4 if rel.get("TargetMode") == "External"
    ]
    assert external_targets == [ground_truth["external_hyperlink"]["target"]]

    # -- layout-level title override
    layout = _member_xml(relpath, "ppt/slideLayouts/slideLayout2.xml")
    override_sizes = [
        defRPr.get("sz")
        for defRPr in layout.iter("{%s}defRPr" % _A)
        if defRPr.getparent().tag == "{%s}lvl1pPr" % _A
    ]
    assert str(ground_truth["layout_title_lvl1_sz_centipoints"]) in override_sizes

    # -- per-shape autofit map (every named box, including the explicit noAutofit ones)
    autofit_seen = {}
    for slide_number in (3, 4):
        slide = _member_xml(relpath, "ppt/slides/slide%d.xml" % slide_number)
        for sp in slide.iter("{%s}sp" % _P):
            name = sp.find(".//{%s}cNvPr" % _P).get("name")
            bodyPr = sp.find(".//{%s}bodyPr" % _A)
            if name in ground_truth["autofit_by_shape"] and bodyPr is not None:
                children = list(bodyPr)
                assert len(children) == 1, name
                autofit_seen[name] = "a:" + etree.QName(children[0]).localname
    assert autofit_seen == ground_truth["autofit_by_shape"]

    # -- notes text, table shape, media sharing, empty body placeholder
    notes = _member_xml(relpath, ground_truth["notes_slide_part"])
    assert ground_truth["notes_text"] in [t.text for t in notes.iter("{%s}t" % _A)]

    slide3 = _member_xml(relpath, "ppt/slides/slide3.xml")
    tables = list(slide3.iter("{%s}tbl" % _A))
    assert len(tables) == 1
    assert len(tables[0].findall("{%s}tr" % _A)) == ground_truth["table"]["rows"]
    grid_cols = tables[0].find("{%s}tblGrid" % _A).findall("{%s}gridCol" % _A)
    assert len(grid_cols) == ground_truth["table"]["cols"]

    for slide_number in ground_truth["media_shared_by_slides"]:
        rels = _member_xml(relpath, "ppt/slides/_rels/slide%d.xml.rels" % slide_number)
        image_targets = [
            rel.get("Target") for rel in rels if rel.get("Type").endswith("/image")
        ]
        assert image_targets == ["../media/image1.png"], "slide %d" % slide_number

    empty_slide = _member_xml(
        relpath, "ppt/slides/slide%d.xml" % ground_truth["empty_body_placeholder_slide"]
    )
    for sp in empty_slide.iter("{%s}sp" % _P):
        ph = sp.find(".//{%s}ph" % _P)
        if ph is not None and ph.get("idx") == "1":
            assert [t.text for t in sp.iter("{%s}t" % _A)] in ([], [None], [""])


def test_large_smoke_ground_truth():
    relpath = "self_generated/large_smoke.pptx"
    ground_truth = _ground_truth(relpath)
    names = _member_names(relpath)
    slide_members = [
        n for n in names if n.startswith("ppt/slides/slide") and n.endswith(".xml")
    ]
    assert len(slide_members) == ground_truth["slide_count"]
    assert [n for n in names if n.startswith("ppt/media/")] == ground_truth["media_parts"]
    with_picture = 0
    for member in slide_members:
        rels_member = "ppt/slides/_rels/" + member.rsplit("/", 1)[1] + ".rels"
        rels = _member_xml(relpath, rels_member)
        if any(rel.get("Type").endswith("/image") for rel in rels):
            with_picture += 1
    assert with_picture == ground_truth["slides_with_picture_count"]


def test_shared_media_image_pixel_size():
    import io

    from PIL import Image

    relpath = "self_generated/shared_media.pptx"
    ground_truth = _ground_truth(relpath)
    with zipfile.ZipFile(str(corpus.fixture_path(relpath))) as zipf:
        image = Image.open(io.BytesIO(zipf.read(ground_truth["media_parts"][0])))
    assert list(image.size) == ground_truth["image_pixel_size"]


def test_lo_chart_notes_style_parts_ground_truth():
    relpath = "libreoffice_export/lo_chart_notes.pptx"
    ground_truth = _ground_truth(relpath)
    names = _member_names(relpath)
    for style_part in ground_truth["chart_style_parts"]:
        assert style_part in names


def test_corrupt_fixture_sldid_entries_ground_truth():
    relpath = "self_generated/corrupt_dangling_sldid.pptx"
    ground_truth = _ground_truth(relpath)
    presentation = _member_xml(relpath, "ppt/presentation.xml")
    entries = [
        {"id": sldId.get("id"), "r_id": sldId.get("{%s}id" % _R)}
        for sldId in presentation.iter("{%s}sldId" % _P)
    ]
    assert entries == ground_truth["sldid_entries"]


def test_lo_branded_template_ground_truth():
    relpath = "libreoffice_export/lo_branded_template.pptx"
    ground_truth = _ground_truth(relpath)
    master = _member_xml(relpath, "ppt/slideMasters/slideMaster1.xml")
    assert (master.find("{%s}txStyles" % _P) is not None) == ground_truth["master_has_txstyles"]
    slide = _member_xml(relpath, "ppt/slides/slide1.xml")
    run_sizes = [
        int(run.find("{%s}rPr" % _A).get("sz")) for run in slide.iter("{%s}r" % _A)
    ]
    baked = ground_truth["baked_run_sizes_centipoints"]
    assert run_sizes == [baked["title"], baked["body_paragraph_0"], baked["body_paragraph_1"]]


def test_sections_ground_truth():
    relpath = "self_generated/sections.pptx"
    ground_truth = _ground_truth(relpath)
    p14 = "http://schemas.microsoft.com/office/powerpoint/2010/main"
    presentation = _member_xml(relpath, "ppt/presentation.xml")

    entries = [
        {"id": sldId.get("id"), "r_id": sldId.get("{%s}id" % _R)}
        for sldId in presentation.find("{%s}sldIdLst" % _P)
    ]
    assert entries == ground_truth["sldid_entries"]

    sections = [
        {
            "name": section.get("name"),
            "slide_ids": [s.get("id") for s in section.iter("{%s}sldId" % p14)],
        }
        for section in presentation.iter("{%s}section" % p14)
    ]
    assert sections == ground_truth["sections"]
    ext_uris = [e.get("uri") for e in presentation.iter("{%s}ext" % _P)]
    assert ground_truth["section_ext_uri"] in ext_uris

    custom_shows = [
        {
            "name": show.get("name"),
            "slide_r_ids": [s.get("{%s}id" % _R) for s in show.iter("{%s}sld" % _P)],
        }
        for show in presentation.iter("{%s}custShow" % _P)
    ]
    assert custom_shows == ground_truth["custom_shows"]


def test_tables_in_group_ground_truth():
    relpath = "self_generated/tables_in_group.pptx"
    ground_truth = _ground_truth(relpath)
    slide = _member_xml(relpath, "ppt/slides/slide1.xml")
    group = slide.find(".//{%s}grpSp" % _P)
    assert [etree.QName(c).localname for c in group] == (
        ground_truth["group_children_localnames"]
    )
    cell_texts = [
        t.text
        for tbl in group.iter("{%s}tbl" % _A)
        for t in tbl.iter("{%s}t" % _A)
    ]
    assert cell_texts == ground_truth["grouped_table"]["cell_texts_row_major"]
    in_group_texts = [
        t.text for sp in group.iter("{%s}sp" % _P) for t in sp.iter("{%s}t" % _A)
    ]
    assert in_group_texts == [ground_truth["in_group_textbox"]["text"]]


def test_nested_groups_ground_truth():
    relpath = "self_generated/nested_groups.pptx"
    ground_truth = _ground_truth(relpath)
    slide = _member_xml(relpath, "ppt/slides/slide1.xml")

    def max_depth(element, depth=0):
        children = element.findall("{%s}grpSp" % _P)
        return depth if not children else max(max_depth(c, depth + 1) for c in children)

    assert max_depth(slide.find(".//{%s}spTree" % _P)) == ground_truth["max_group_depth"]
    all_texts = sorted(t.text for t in slide.iter("{%s}t" % _A))
    expected = sorted(v["text"] for v in ground_truth["texts_by_level"].values())
    assert all_texts == expected


def test_autofit_inherited_ground_truth():
    relpath = "self_generated/autofit_inherited.pptx"
    ground_truth = _ground_truth(relpath)
    slide = _member_xml(relpath, "ppt/slides/slide1.xml")
    autofits = [
        (etree.QName(child).localname, dict(child.attrib))
        for bodyPr in slide.iter("{%s}bodyPr" % _A)
        for child in bodyPr
    ]
    assert autofits == [("normAutofit", {"fontScale": ground_truth["fontscale_attr"]})]
    body_runs_with_size = [
        r for r in slide.iter("{%s}r" % _A)
        if r.find("{%s}rPr" % _A) is not None and r.find("{%s}rPr" % _A).get("sz")
    ]
    assert bool(body_runs_with_size) == ground_truth["body_runs_carry_local_size"]
    master = _member_xml(relpath, "ppt/slideMasters/slideMaster1.xml")
    body_l1 = master.find("{%s}txStyles/{%s}bodyStyle/{%s}lvl1pPr/{%s}defRPr" % (_P, _P, _A, _A))
    assert int(body_l1.get("sz")) == ground_truth["master_body_lvl1_sz_centipoints"]


def test_hf_flags_ground_truth():
    relpath = "self_generated/hf_flags.pptx"
    ground_truth = _ground_truth(relpath)
    master = _member_xml(relpath, "ppt/slideMasters/slideMaster1.xml")
    master_hf = master.find("{%s}hf" % _P)
    assert dict(master_hf.attrib) == ground_truth["master_hf"]
    layout = _member_xml(relpath, "ppt/slideLayouts/slideLayout1.xml")
    layout_hf = layout.find("{%s}hf" % _P)
    assert dict(layout_hf.attrib) == ground_truth["layout1_hf"]

    # -- the HeaderFooters proxy reads authored-elsewhere flags per the sidecar
    prs = Presentation(str(corpus.fixture_path(relpath)))
    expected = ground_truth["expected_reads"]
    master_hf_proxy = prs.slide_masters[0].header_footers
    assert master_hf_proxy.slide_number_visible == expected["master"]["slide_number_visible"]
    assert master_hf_proxy.footer_visible == expected["master"]["footer_visible"]
    assert master_hf_proxy.date_visible == expected["master"]["date_visible"]
    layout_hf_proxy = prs.slides[0].slide_layout.header_footers
    assert layout_hf_proxy.slide_number_visible == expected["layout1"]["slide_number_visible"]
    assert layout_hf_proxy.footer_visible is None


# ------------------------------------------------------------------- fixtures


def _table_rows(relpath):
    slide = _member_xml(relpath, "ppt/slides/slide1.xml")
    tbl = slide.find(".//{%s}tbl" % _A)
    assert tbl is not None
    grid_cols = tbl.find("{%s}tblGrid" % _A).findall("{%s}gridCol" % _A)
    rows = tbl.findall("{%s}tr" % _A)
    return grid_cols, [tr.findall("{%s}tc" % _A) for tr in rows]


@pytest.mark.parametrize(
    "relpath", ["self_generated/merged_tables.pptx", "libreoffice_export/lo_merged_tables.pptx"]
)
def test_merged_tables_ground_truth(relpath):
    ground_truth = _ground_truth(relpath)
    grid_cols, rows = _table_rows(relpath)
    assert len(grid_cols) == ground_truth["grid_col_count"]
    assert len(rows) == ground_truth["row_count"]
    # -- the grid invariant: every row holds exactly one a:tc per a:gridCol
    assert all(len(row) == len(grid_cols) for row in rows)
    assert rows[0][0].get("gridSpan") == "4"
    assert all(rows[0][col].get("hMerge") == "1" for col in (1, 2, 3))
    assert rows[2][0].get("rowSpan") == "2"
    assert rows[3][0].get("vMerge") == "1"
    header_text = "".join(t.text or "" for t in rows[0][0].iter("{%s}t" % _A))
    assert header_text == "Merged header"


def _hf_placeholders(slide):
    """Map ph type -> (idx, fld types, all text) for dt/ftr/sldNum placeholders on `slide`."""
    found = {}
    for sp in slide.iter("{%s}sp" % _P):
        ph = sp.find(".//{%s}nvSpPr/{%s}nvPr/{%s}ph" % (_P, _P, _P))
        if ph is None or ph.get("type") not in ("dt", "ftr", "sldNum"):
            continue
        flds = [f.get("type") for f in sp.findall(".//{%s}fld" % _A)]
        text = "".join(t.text or "" for t in sp.iter("{%s}t" % _A))
        found[ph.get("type")] = (int(ph.get("idx")), flds, text)
    return found


@pytest.mark.parametrize(
    "relpath",
    ["self_generated/footers_applied.pptx", "libreoffice_export/lo_footers_applied.pptx"],
)
def test_footers_applied_ground_truth(relpath):
    ground_truth = _ground_truth(relpath)
    for ordinal in range(1, 6):
        slide = _member_xml(relpath, "ppt/slides/slide%d.xml" % ordinal)
        found = _hf_placeholders(slide)
        assert ("ftr" in found) == (ordinal in ground_truth["slides_with_ftr_placeholder"])
        assert "dt" in found
        assert "sldNum" in found
        sldnum_idx, sldnum_flds, sldnum_text = found["sldNum"]
        assert sldnum_flds == ["slidenum"]
        if ground_truth["sldNum_cached_text_equals_ordinal"]:
            assert sldnum_text == str(ordinal)
        if "ftr" in found:
            assert found["ftr"][2] == ground_truth["footer_text"]
            assert found["ftr"][1] == []  # -- footer is a literal, never a field
    if "hidden_slide_ordinal" in ground_truth:
        # -- `show` lives on the slide part's ROOT `p:sld` element (CT_Slide), never on
        # -- the presentation's p:sldId entries (a mechanism finding)
        hidden = [
            ordinal
            for ordinal in range(1, 6)
            if _member_xml(relpath, "ppt/slides/slide%d.xml" % ordinal).get("show") == "0"
        ]
        assert hidden == [ground_truth["hidden_slide_ordinal"]]


def test_footers_applied_hf_absent_everywhere():
    ground_truth = _ground_truth("self_generated/footers_applied.pptx")
    assert ground_truth["p_hf_absent_everywhere"] is True
    for member in _member_names("self_generated/footers_applied.pptx"):
        if member.startswith(("ppt/slides/slide", "ppt/slideLayouts/", "ppt/slideMasters/")):
            if not member.endswith(".xml"):
                continue
            root = _member_xml("self_generated/footers_applied.pptx", member)
            assert root.find("{%s}hf" % _P) is None, member


def test_scrub_gauntlet_ground_truth():
    relpath = "self_generated/scrub_gauntlet.pptx"
    ground_truth = _ground_truth(relpath)
    members = _member_names(relpath)
    assert ground_truth["comments_part"] in members
    assert ground_truth["comment_authors_part"] in members
    assert ground_truth["custom_props_part"] in members
    assert ground_truth["embedded_font_part"] in members
    assert ground_truth["unused_layout_media_part"] in members

    comments = _member_xml(relpath, ground_truth["comments_part"])
    texts = [t.text for t in comments.iter("{%s}text" % _P)]
    assert texts == ground_truth["comment_texts"]
    authors = _member_xml(relpath, ground_truth["comment_authors_part"])
    assert [a.get("name") for a in authors] == [ground_truth["comment_author_name"]]

    presentation = _member_xml(relpath, "ppt/presentation.xml")
    fontLst = presentation.find("{%s}embeddedFontLst" % _P)
    assert fontLst is not None
    typeface = fontLst.find(".//{%s}font" % _P).get("typeface")
    assert typeface == ground_truth["embedded_font_typeface"]

    # -- the unused-layout picture is the ONLY reference to its media part
    layout = _member_xml(relpath, ground_truth["unused_layout_part"])
    assert layout.find(".//{%s}pic" % _P) is not None
    # -- no slide references the unused layout (that is what makes it unused)
    for member in members:
        if member.startswith("ppt/slides/_rels/"):
            rels = _member_xml(relpath, member)
            targets = [r.get("Target") for r in rels]
            assert not any("slideLayout4.xml" in t for t in targets)

    notes_members = [m for m in members if m.startswith("ppt/notesSlides/notesSlide")]
    assert len(notes_members) == len(ground_truth["slides_with_notes"])


@pytest.mark.parametrize(
    ("relpath", "expect"),
    [
        ("self_generated/template_alpha.pptx", ("Paper Alpha", "Georgia", "Verdana", "AA3311")),
        (
            "libreoffice_export/lo_template_alpha.pptx",
            ("Paper Alpha", "Georgia", "Verdana", "AA3311"),
        ),
        (
            "self_generated/template_beta.pptx",
            ("Paper Beta", "Courier New", "Times New Roman", "1166BB"),
        ),
    ],
)
def test_template_identity_ground_truth(relpath, expect):
    name, major, minor, accent1 = expect
    theme = _member_xml(relpath, "ppt/theme/theme1.xml")
    assert theme.get("name") == name
    font_scheme = theme.find("{%s}themeElements/{%s}fontScheme" % (_A, _A))
    assert font_scheme.find("{%s}majorFont/{%s}latin" % (_A, _A)).get("typeface") == major
    assert font_scheme.find("{%s}minorFont/{%s}latin" % (_A, _A)).get("typeface") == minor
    accent = theme.find(
        "{%s}themeElements/{%s}clrScheme/{%s}accent1/{%s}srgbClr" % (_A, _A, _A, _A)
    )
    assert accent.get("val") == accent1


def test_template_pair_layout_name_contract():
    """Every alpha layout name exists in beta except beta's renamed 'Beta Special' —
    the import-collision fixture requirement (a real-PowerPoint fixture a human must author)."""

    def layout_names(relpath):
        names = set()
        for member in _member_names(relpath):
            if member.startswith("ppt/slideLayouts/slideLayout") and member.endswith(".xml"):
                root = _member_xml(relpath, member)
                names.add(root.find("{%s}cSld" % _P).get("name"))
        return names

    alpha = layout_names("self_generated/template_alpha.pptx")
    beta = layout_names("self_generated/template_beta.pptx")
    assert "Beta Special" in beta
    assert "Beta Special" not in alpha
    assert "Title Only" in alpha
    assert "Title Only" not in beta
    assert alpha - {"Title Only"} == beta - {"Beta Special"}  # -- everything else collides


def _slide_ids(relpath):
    presentation = _member_xml(relpath, "ppt/presentation.xml")
    return [int(s.get("id")) for s in presentation.find("{%s}sldIdLst" % _P)]


def test_lineage_slide_id_ground_truth():
    assert _slide_ids("self_generated/lineage_v1.pptx") == [256, 257, 258, 259, 260]
    assert _slide_ids("self_generated/lineage_v2.pptx") == [257, 256, 258, 259, 261]
    assert _slide_ids("self_generated/lineage_reorder.pptx") == [260, 256, 257, 258, 259]
    # -- the sidecars carry the same claim (they are the diff organ's ground truth)
    for name in ("lineage_v1", "lineage_v2", "lineage_reorder"):
        ground_truth = _ground_truth("self_generated/%s.pptx" % name)
        assert ground_truth["slide_ids"] == _slide_ids("self_generated/%s.pptx" % name)


def test_lineage_reorder_is_content_identical_to_v1():
    """The reorder-only variant differs from v1 ONLY in presentation.xml slide order:
    every slide/notes/chart member byte-matches some member of v1."""
    v1_map = zip_member_map(corpus.fixture_path("self_generated/lineage_v1.pptx").read_bytes())
    reorder_map = zip_member_map(
        corpus.fixture_path("self_generated/lineage_reorder.pptx").read_bytes()
    )
    assert set(v1_map) == set(reorder_map)
    for member, payload in reorder_map.items():
        if member.startswith(("ppt/slides/slide", "ppt/notesSlides/", "ppt/charts/")):
            assert payload in v1_map.values(), "%s not byte-identical to any v1 member" % member


def test_lineage_v2_edits_ground_truth():
    relpath = "self_generated/lineage_v2.pptx"
    ground_truth = _ground_truth(relpath)
    ops = {edit["op"] for edit in ground_truth["edits"]}
    assert ops == {
        "text", "chart_data", "notes", "geometry", "image_replace",
        "add_slide", "delete_slide", "move_slide",
    }
    # -- spot-verify against bytes: retitle landed, deleted title gone, new title present.
    # -- Slide partnames are NOT contiguous after delete+add (slide5.xml gone, slide6.xml
    # -- new) - iterate the actual members.
    def all_text(root):
        return "".join(t.text or "" for t in root.iter("{%s}t" % _A))

    slide_members = sorted(
        m
        for m in _member_names(relpath)
        if m.startswith("ppt/slides/slide") and m.endswith(".xml")
    )
    assert len(slide_members) == 5
    texts = [all_text(_member_xml(relpath, member)) for member in slide_members]
    assert any("Lineage slide one, retitled" in t for t in texts)
    assert not any("Lineage slide five" in t for t in texts)
    assert any("Lineage slide six, new" in t for t in texts)


# ------------------------------------------------------------------- independent loader smoke


@pytest.mark.lo_smoke
@pytest.mark.parametrize("relpath", NONCORRUPT_RELPATHS)
def test_fixture_loads_in_libreoffice(relpath, tmp_path):
    lo_load_smoke(corpus.fixture_path(relpath), tmp_path)
